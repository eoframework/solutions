#!/usr/bin/env python3
import os, sys, json, argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor
import yaml

BRAND_RED = RGBColor(204, 0, 0)  # #c00

def get_layout(prs, name_contains):
    for layout in prs.slide_layouts:
        if name_contains.lower() in layout.name.lower():
            return layout
    raise ValueError(f"Layout containing '{name_contains}' not found. Available layouts: {[l.name for l in prs.slide_layouts]}")

def set_title_and_subtitle(slide, title_text=None, subtitle_text=None):
    if title_text and slide.shapes.title:
        slide.shapes.title.text = title_text
    # try most common subtitle placeholder
    subs = [ph for ph in slide.placeholders if getattr(ph, "name", "").lower().startswith("subtitle") or ph.placeholder_format.idx == 1]
    if subtitle_text and subs:
        subs[0].text = subtitle_text

def first_body_placeholder(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.type == 2:  # BODY
            return ph
    return None

def first_picture_placeholder(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            return ph
    return None

def fill_bullets(body_placeholder, bullets):
    tf = body_placeholder.text_frame
    tf.clear()
    first = True
    for b in bullets:
        if first:
            tf.text = b
            first = False
        else:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0

def place_picture_exact(slide, anchor_placeholder, image_path):
    left, top = anchor_placeholder.left, anchor_placeholder.top
    width, height = anchor_placeholder.width, anchor_placeholder.height
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)

def add_slide_from_spec(prs, spec):
    layout_name = spec.get("layout", "Content")
    layout = get_layout(prs, layout_name)

    slide = prs.slides.add_slide(layout)
    set_title_and_subtitle(slide, spec.get("title"), spec.get("subtitle"))

    bullets = spec.get("bullets")
    if bullets:
        body = first_body_placeholder(slide)
        if body:
            fill_bullets(body, bullets)

    if "left" in spec:
        bodies = [ph for ph in slide.placeholders if ph.placeholder_format.type == 2]
        bodies_sorted = sorted(bodies, key=lambda p: (p.left, p.top))
        if bodies_sorted:
            tf = bodies_sorted[0].text_frame
            tf.clear()
            first = True
            for para in spec["left"]:
                if first:
                    tf.text = para
                    first = False
                else:
                    p = tf.add_paragraph()
                    p.text = para
                    p.level = 0

    if "right_image" in spec and spec["right_image"]:
        pic_ph = first_picture_placeholder(slide)
        if pic_ph:
            pic_ph.insert_picture(spec["right_image"])
        else:
            bodies = [ph for ph in slide.placeholders if ph.placeholder_format.type == 2]
            if bodies:
                bodies_sorted = sorted(bodies, key=lambda p: (p.left, p.top))
                anchor = bodies_sorted[-1]
                place_picture_exact(slide, anchor, spec["right_image"])

    if "image" in spec and spec["image"]:
        pic_ph = first_picture_placeholder(slide)
        if pic_ph:
            pic_ph.insert_picture(spec["image"])
        else:
            body = first_body_placeholder(slide)
            if body:
                place_picture_exact(slide, body, spec["image"])

    if "caption" in spec and spec["caption"]:
        bodies = [ph for ph in slide.placeholders if ph.placeholder_format.type == 2]
        if len(bodies) >= 2:
            tf = bodies[1].text_frame
            tf.clear()
            tf.text = spec["caption"]

def load_content(path):
    with open(path, "r") as f:
        if path.endswith((".yml",".yaml")):
            return yaml.safe_load(f)
        return json.load(f)

def main():
    ap = argparse.ArgumentParser(description="Generate EO deck from JSON/YAML using a .potx template.")
    ap.add_argument("--template", required=True, help="Path to .potx template")
    ap.add_argument("--content", required=True, help="Path to JSON/YAML content file")
    ap.add_argument("--out", default="out/output.pptx", help="Output .pptx path")
    args = ap.parse_args()

    content = load_content(args.content)
    prs = Presentation(args.template)

    if content.get("title") or content.get("subtitle"):
        slide = prs.slides.add_slide(get_layout(prs, "Title"))
        set_title_and_subtitle(slide, content.get("title"), content.get("subtitle"))

    for spec in content.get("slides", []):
        add_slide_from_spec(prs, spec)

    os.makedirs(os.path.dirname(args.out), exist_ok=True)
    prs.save(args.out)
    print(f"Saved: {args.out}")

if __name__ == "__main__":
    main()
