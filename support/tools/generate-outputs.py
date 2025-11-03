#!/usr/bin/env python3
"""
Generate professional output files from raw source files.

Converts:
- Markdown (.md) → Word (.docx) or PowerPoint (.pptx)
- CSV (.csv) → Excel (.xlsx) with formatting

Usage:
    python3 generate-outputs.py --path solutions/provider/category/solution
    python3 generate-outputs.py --path solution-template/
"""

import argparse
import os
import sys
from pathlib import Path
import pandas as pd
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.styles.numbers import FORMAT_CURRENCY_USD_SIMPLE
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.utils import get_column_letter
import numpy as np
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER
from markdown_it import MarkdownIt
from datetime import datetime
import re


class OutputGenerator:
    """Generate professional output files from raw source files."""

    def __init__(self, base_path):
        self.base_path = Path(base_path)
        self.stats = {
            'csv_to_xlsx': 0,
            'md_to_docx': 0,
            'md_to_pptx': 0,
            'errors': []
        }

    def generate_all(self):
        """Generate all output files from raw directories."""
        print(f"🔍 Scanning: {self.base_path}")

        # Find all raw directories
        raw_dirs = list(self.base_path.rglob('raw'))

        if not raw_dirs:
            print("❌ No 'raw' directories found.")
            return False

        print(f"📁 Found {len(raw_dirs)} raw directories\n")

        for raw_dir in raw_dirs:
            self.process_raw_directory(raw_dir)

        self.print_summary()
        return len(self.stats['errors']) == 0

    def process_raw_directory(self, raw_dir):
        """Process all files in a raw directory."""
        output_dir = raw_dir.parent

        print(f"📂 Processing: {raw_dir.relative_to(self.base_path)}")

        # Process CSV files
        for csv_file in raw_dir.glob('*.csv'):
            self.convert_csv_to_xlsx(csv_file, output_dir)

        # Process Markdown files
        for md_file in raw_dir.glob('*.md'):
            if 'presentation' in md_file.stem.lower():
                self.convert_md_to_pptx(md_file, output_dir)
            else:
                self.convert_md_to_docx(md_file, output_dir)

        print()

    def _is_currency_column(self, column_name, sample_values):
        """Detect if a column contains currency values."""
        # Check column name for currency keywords
        currency_keywords = ['cost', 'price', 'amount', 'total', 'budget', 'fee', 'rate',
                            'salary', 'revenue', 'expense', 'payment', 'charge', 'value',
                            'investment', 'profit', 'loss', 'income', 'worth', 'usd', '$']

        column_lower = column_name.lower()
        if any(keyword in column_lower for keyword in currency_keywords):
            return True

        # Check sample values for $ signs
        for val in sample_values[:10]:  # Check first 10 values
            if val and isinstance(val, str) and '$' in val:
                return True

        return False

    def _convert_currency_to_number(self, value):
        """Convert currency string to numeric value."""
        if pd.isna(value) or value == '':
            return None

        if isinstance(value, (int, float)):
            return float(value)

        # Remove currency symbols and formatting
        if isinstance(value, str):
            # Remove $, commas, spaces, parentheses (for negative)
            clean = value.replace('$', '').replace(',', '').replace(' ', '').strip()

            # Handle parentheses for negative numbers
            if clean.startswith('(') and clean.endswith(')'):
                clean = '-' + clean[1:-1]

            try:
                return float(clean)
            except (ValueError, AttributeError):
                return None

        return None

    def _generate_sheet_name(self, filename):
        """Generate appropriate sheet name from filename."""
        # Common abbreviations and mappings
        name_mappings = {
            'level-of-effort-estimate': 'LOE',
            'requirements-questionnaire': 'Requirements',
            'roi-calculator': 'ROI',
            'communication-plan': 'Comm Plan',
            'configuration': 'Configuration',
            'project-plan': 'Project Plan',
            'requirements': 'Requirements',
            'roles': 'Roles',
            'test-plan': 'Test Plan',
            'training-plan': 'Training Plan',
        }

        # Check for exact match
        if filename.lower() in name_mappings:
            return name_mappings[filename.lower()]

        # Generate from filename (title case, remove dashes/underscores)
        sheet_name = filename.replace('-', ' ').replace('_', ' ').title()

        # Limit to 31 characters (Excel sheet name limit)
        if len(sheet_name) > 31:
            # Try to abbreviate
            words = sheet_name.split()
            if len(words) > 1:
                # Use first letter of each word except last
                abbrev = ''.join([w[0] for w in words[:-1]]) + ' ' + words[-1]
                sheet_name = abbrev[:31]
            else:
                sheet_name = sheet_name[:31]

        return sheet_name

    def convert_csv_to_xlsx(self, csv_file, output_dir):
        """Convert CSV to styled Excel file using branded template."""
        try:
            output_file = output_dir / f"{csv_file.stem}.xlsx"

            # Load branded template
            template_path = Path(__file__).parent.parent / 'doc-templates' / 'excel' / 'EOFramework-Excel-Template-01.xlsx'

            if not template_path.exists():
                raise FileNotFoundError(f"Excel template not found: {template_path}")

            wb = load_workbook(str(template_path))

            # Read CSV
            df = pd.read_csv(csv_file)

            # Detect currency columns
            currency_columns = []
            for col_idx, col_name in enumerate(df.columns):
                if self._is_currency_column(col_name, df[col_name].dropna().tolist()):
                    currency_columns.append(col_idx)
                    # Convert currency strings to numbers
                    df[col_name] = df[col_name].apply(self._convert_currency_to_number)

            # Update Cover sheet with metadata
            if 'Cover' in wb.sheetnames:
                cover = wb['Cover']

                # Update document title
                cover['B4'] = f"{csv_file.stem.replace('-', ' ').replace('_', ' ').title()}"

                # Update document type
                if 'presales' in str(csv_file):
                    doc_type = 'Presales Document'
                elif 'delivery' in str(csv_file):
                    doc_type = 'Delivery Document'
                else:
                    doc_type = 'Reference Document'
                cover['B5'] = doc_type

                # Update metadata (using new compact layout B7-C10)
                cover['C7'] = datetime.now().strftime('%B %d, %Y')

                # Extract solution name from path
                solution_name = 'EO Framework Solution'
                try:
                    parts = csv_file.parts
                    if 'solution-template' in parts:
                        solution_name = 'Sample Solution'
                    else:
                        # Try to extract from path
                        for i, part in enumerate(parts):
                            if part in ['presales', 'delivery'] and i > 0:
                                solution_name = parts[i-1].replace('-', ' ').title()
                                break
                except:
                    pass

                cover['C8'] = solution_name
                cover['C9'] = f"Data table for {csv_file.stem.replace('-', ' ')}"

            # Get Data sheet and rename it based on file name
            ws = wb['Data']

            # Generate appropriate sheet name from file name
            sheet_name = self._generate_sheet_name(csv_file.stem)
            ws.title = sheet_name

            # Clear sample data (rows 2 onwards)
            ws.delete_rows(2, ws.max_row)

            # Define styling
            header_font = Font(name='Calibri', size=12, bold=True, color='FFFFFF')  # Size 12
            header_fill = PatternFill(start_color='808080', end_color='808080', fill_type='solid')  # Gray header
            header_alignment = Alignment(horizontal='center', vertical='center')

            data_font = Font(name='Calibri', size=12, color='000000')  # Size 12
            data_alignment = Alignment(horizontal='left', vertical='center', wrap_text=True, indent=1)  # Center vertical + indent for padding

            # Alternating row colors for better readability
            even_fill = PatternFill(start_color='F2F2F2', end_color='F2F2F2', fill_type='solid')  # Light gray
            odd_fill = PatternFill(start_color='FFFFFF', end_color='FFFFFF', fill_type='solid')   # White

            thin_border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )

            # Write data
            for r_idx, row in enumerate(dataframe_to_rows(df, index=False, header=True), 1):
                for c_idx, value in enumerate(row, 1):
                    # Handle formula text - write as text to prevent Excel errors
                    # If value starts with '=' and is in data row, it's likely formula text that should be displayed, not executed
                    if r_idx > 1 and value and isinstance(value, str) and value.startswith('='):
                        # Write as text by prefixing with apostrophe (Excel will display without the apostrophe)
                        cell = ws.cell(row=r_idx, column=c_idx)
                        cell.value = value  # Set as plain string
                        # Force it to be text format
                        cell.number_format = '@'  # Text format
                    else:
                        cell = ws.cell(row=r_idx, column=c_idx, value=value)

                    cell.border = thin_border

                    # Style header row
                    if r_idx == 1:
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = header_alignment
                    else:
                        # Style data rows
                        cell.font = data_font
                        cell.alignment = data_alignment

                        # Apply currency formatting to currency columns (but not if it's formula text)
                        if (c_idx - 1) in currency_columns and not (value and isinstance(value, str) and value.startswith('=')):
                            cell.number_format = '$#,##0.00'  # Currency format with 2 decimals

                        # Alternating row colors
                        if r_idx % 2 == 0:
                            cell.fill = even_fill
                        else:
                            cell.fill = odd_fill

            # Auto-adjust column widths based on content
            for col_idx, column in enumerate(ws.columns, 1):
                max_length = 0
                column_letter = get_column_letter(col_idx)

                for cell in column:
                    try:
                        if cell.value:
                            # Calculate display length based on value and format
                            if cell.number_format and '$' in cell.number_format:
                                # Currency values need extra space for $ and formatting
                                # Example: $180,000.00 needs more space than raw value 180000
                                if isinstance(cell.value, (int, float)):
                                    # Format as currency to get actual display length
                                    formatted = f"${cell.value:,.2f}"
                                    cell_length = len(formatted)
                                else:
                                    cell_length = len(str(cell.value))
                            else:
                                cell_length = len(str(cell.value))

                            max_length = max(max_length, cell_length)
                    except:
                        pass

                # Add padding (3 characters) and set reasonable bounds
                adjusted_width = min(max(max_length + 3, 10), 50)  # Min 10, Max 50
                ws.column_dimensions[column_letter].width = adjusted_width

            # Set minimum row heights with auto-adjust capability (maintains padding while allowing content expansion)
            ws.row_dimensions[1].height = 32  # Header row - minimum height with padding
            for row_idx in range(2, len(df) + 2):
                ws.row_dimensions[row_idx].height = 26  # Data rows - minimum height, will auto-expand with content

            # Update auto-filter range to match actual data
            num_cols = len(df.columns)
            ws.auto_filter.ref = f'A1:{get_column_letter(num_cols)}{len(df) + 1}'

            # Freeze header row
            ws.freeze_panes = ws['A2']

            wb.save(output_file)
            print(f"  ✅ {csv_file.name} → {output_file.name}")
            self.stats['csv_to_xlsx'] += 1

        except Exception as e:
            error_msg = f"Error converting {csv_file.name}: {str(e)}"
            print(f"  ❌ {error_msg}")
            self.stats['errors'].append(error_msg)

    def convert_md_to_docx(self, md_file, output_dir):
        """Convert Markdown to styled Word document."""
        try:
            output_file = output_dir / f"{md_file.stem}.docx"

            # Read markdown
            with open(md_file, 'r', encoding='utf-8') as f:
                md_content = f.read()

            # Parse markdown
            md = MarkdownIt()
            tokens = md.parse(md_content)

            # Create Word document
            doc = Document()

            # Process tokens
            self._process_md_tokens(doc, tokens, md_content)

            doc.save(output_file)
            print(f"  ✅ {md_file.name} → {output_file.name}")
            self.stats['md_to_docx'] += 1

        except Exception as e:
            error_msg = f"Error converting {md_file.name}: {str(e)}"
            print(f"  ❌ {error_msg}")
            self.stats['errors'].append(error_msg)

    def _process_md_tokens(self, doc, tokens, md_content):
        """Process markdown tokens and add to Word document."""
        i = 0
        while i < len(tokens):
            token = tokens[i]

            if token.type == 'heading_open':
                level = int(token.tag[1])  # h1 -> 1, h2 -> 2, etc.
                i += 1
                if i < len(tokens) and tokens[i].type == 'inline':
                    text = tokens[i].content
                    doc.add_heading(text, level=level)
                i += 1  # Skip heading_close

            elif token.type == 'paragraph_open':
                i += 1
                if i < len(tokens) and tokens[i].type == 'inline':
                    text = tokens[i].content
                    # Simple text formatting
                    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)  # Remove bold markers
                    text = re.sub(r'\*(.+?)\*', r'\1', text)      # Remove italic markers
                    if text.strip():
                        doc.add_paragraph(text)
                i += 1  # Skip paragraph_close

            elif token.type == 'bullet_list_open':
                i += 1
                while i < len(tokens) and tokens[i].type != 'bullet_list_close':
                    if tokens[i].type == 'list_item_open':
                        i += 1
                        if i < len(tokens) and tokens[i].type == 'paragraph_open':
                            i += 1
                            if i < len(tokens) and tokens[i].type == 'inline':
                                text = tokens[i].content
                                text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
                                text = re.sub(r'\*(.+?)\*', r'\1', text)
                                doc.add_paragraph(text, style='List Bullet')
                    i += 1
            else:
                i += 1

    def _get_layout(self, prs, name_contains):
        """Get layout by partial name match."""
        for layout in prs.slide_layouts:
            if name_contains.lower() in layout.name.lower():
                return layout
        raise ValueError(f"Layout containing '{name_contains}' not found")

    def _insert_logos(self, slide):
        """Insert logos into picture placeholders on a slide."""
        # Path to logo files (common assets directory)
        logo_path = Path(__file__).parent.parent / 'doc-templates' / 'assets' / 'logos'

        # Map of placeholder indices to logo files
        # These match the template structure
        logos = {
            10: 'client_logo.png',
            13: 'consulting_company_logo.png',
            14: 'consulting_company_logo.png',  # Some layouts use idx 14 for logo
        }

        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            if idx in logos and ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                logo_file = logo_path / logos[idx]
                if logo_file.exists():
                    try:
                        ph.insert_picture(str(logo_file))
                    except Exception as e:
                        # Silently continue if logo insertion fails
                        pass

    def convert_md_to_pptx(self, md_file, output_dir):
        """Convert Markdown to PowerPoint presentation using branded template."""
        try:
            output_file = output_dir / f"{md_file.stem}.pptx"

            # Load branded template
            template_path = Path(__file__).parent.parent / 'doc-templates' / 'powerpoint' / 'EOFramework-Template-01.pptx'

            if not template_path.exists():
                raise FileNotFoundError(f"Template not found: {template_path}")

            prs = Presentation(str(template_path))

            # Clear demo slides (layouts-only approach)
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]

            # Read markdown
            with open(md_file, 'r', encoding='utf-8') as f:
                md_content = f.read()

            # Parse markdown into slides
            slides_data = self._parse_md_for_presentation(md_content)

            # Generate slides from layouts
            for slide_data in slides_data:
                self._add_slide_from_layout(prs, slide_data)

            # Set the presentation to open in Normal (Slide) view, not Slide Master view
            # This ensures it opens in presentation view instead of master view
            try:
                # Access the presentation part and set the view properties
                prs._element.attrib.pop('showMasterSp', None)  # Remove any master view settings
            except:
                pass  # If this fails, the file will still work

            prs.save(output_file)
            print(f"  ✅ {md_file.name} → {output_file.name}")
            self.stats['md_to_pptx'] += 1

        except Exception as e:
            error_msg = f"Error converting {md_file.name}: {str(e)}"
            print(f"  ❌ {error_msg}")
            self.stats['errors'].append(error_msg)

    def _parse_md_for_presentation(self, md_content):
        """Parse markdown content into presentation slides.

        Expects format:
        # Main Title (H1) → Title slide
        ### Slide Title (H3) → Content slides
        --- → Slide separator
        """
        slides = []
        lines = md_content.split('\n')
        current_slide = None
        first_h1_found = False

        i = 0
        while i < len(lines):
            line = lines[i].strip()

            # H1 creates title slide (only first one)
            if line.startswith('# ') and not first_h1_found:
                if current_slide:
                    slides.append(current_slide)

                title_text = line[2:].strip()

                # Look ahead for subtitle (next non-empty line that's not a header)
                subtitle_text = ''
                j = i + 1
                while j < len(lines):
                    next_line = lines[j].strip()
                    if next_line and not next_line.startswith('#') and not next_line.startswith('---'):
                        subtitle_text = next_line
                        break
                    j += 1

                current_slide = {
                    'type': 'title',
                    'title': title_text,
                    'subtitle': subtitle_text,
                    'presenter': 'EO Framework Team',
                    'date': datetime.now().strftime('%B %d, %Y')
                }
                first_h1_found = True

            # H3 creates content slide
            elif line.startswith('### '):
                if current_slide:
                    slides.append(current_slide)

                title_text = line[4:].strip()
                current_slide = {
                    'type': 'content',
                    'title': title_text,
                    'bullets': [],
                    'content': []
                }

            # Bullet points
            elif (line.startswith('- ') or line.startswith('* ')) and current_slide:
                bullet_text = line[2:].strip()
                # Remove markdown formatting
                bullet_text = re.sub(r'\*\*(.+?)\*\*', r'\1', bullet_text)
                bullet_text = re.sub(r'\*(.+?)\*', r'\1', bullet_text)
                bullet_text = re.sub(r'🔴|✅', '', bullet_text).strip()  # Remove emoji markers

                if current_slide['type'] == 'content':
                    current_slide['bullets'].append(bullet_text)

            # Numbered list items
            elif re.match(r'^\d+\.', line) and current_slide:
                item_text = re.sub(r'^\d+\.\s*', '', line).strip()
                item_text = re.sub(r'\*\*(.+?)\*\*', r'\1', item_text)
                if current_slide['type'] == 'content':
                    current_slide['bullets'].append(item_text)

            # Regular content lines
            elif line and not line.startswith('#') and not line.startswith('---') and current_slide:
                if current_slide['type'] == 'content':
                    current_slide['content'].append(line)

            i += 1

        if current_slide:
            slides.append(current_slide)

        return slides

    def _add_slide_from_layout(self, prs, slide_data):
        """Add slide using appropriate template layout."""
        try:
            if slide_data['type'] == 'title':
                # Use "EO Title Slide" layout
                layout = self._get_layout(prs, "Title")
                slide = prs.slides.add_slide(layout)

                # Fill placeholders: [12]=Title, [14]=Subtitle, [15]=Presenter|Date
                try:
                    slide.placeholders[12].text = slide_data['title']
                except:
                    pass

                try:
                    slide.placeholders[14].text = slide_data.get('subtitle', '')
                except:
                    pass

                try:
                    presenter_text = f"{slide_data.get('presenter', 'EO Framework Team')} | {slide_data.get('date', datetime.now().strftime('%B %d, %Y'))}"
                    slide.placeholders[15].text = presenter_text
                except:
                    pass

                # Insert logos
                self._insert_logos(slide)

            else:  # content slide
                # Use "EO Single Column" layout for content slides
                layout = self._get_layout(prs, "Single")
                slide = prs.slides.add_slide(layout)

                # Fill title: [10]=Title
                try:
                    slide.placeholders[10].text = slide_data['title']
                except:
                    pass

                # Fill body content: [11]=Body
                if slide_data.get('bullets'):
                    try:
                        body = slide.placeholders[11]
                        tf = body.text_frame
                        tf.clear()

                        # First bullet
                        tf.text = slide_data['bullets'][0]

                        # Remaining bullets
                        for bullet in slide_data['bullets'][1:]:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                    except:
                        pass

                # Insert logos
                self._insert_logos(slide)

        except Exception as e:
            # Log error but continue with other slides
            print(f"    ⚠️  Warning: Could not create slide '{slide_data.get('title', 'Unknown')}': {str(e)}")

    def print_summary(self):
        """Print generation summary."""
        print("\n" + "="*60)
        print("📊 Generation Summary")
        print("="*60)
        print(f"✅ CSV → Excel:      {self.stats['csv_to_xlsx']} files")
        print(f"✅ MD → Word:        {self.stats['md_to_docx']} files")
        print(f"✅ MD → PowerPoint:  {self.stats['md_to_pptx']} files")

        total = self.stats['csv_to_xlsx'] + self.stats['md_to_docx'] + self.stats['md_to_pptx']
        print(f"\n📁 Total generated:  {total} files")

        if self.stats['errors']:
            print(f"\n❌ Errors: {len(self.stats['errors'])}")
            for error in self.stats['errors']:
                print(f"   • {error}")
        else:
            print("\n✨ All files generated successfully!")
        print("="*60)


def main():
    parser = argparse.ArgumentParser(
        description='Generate professional output files from raw source files.'
    )
    parser.add_argument(
        '--path',
        type=str,
        required=True,
        help='Path to solution directory or solution-template'
    )

    args = parser.parse_args()

    # Validate path exists
    path = Path(args.path)
    if not path.exists():
        print(f"❌ Error: Path does not exist: {path}")
        sys.exit(1)

    # Generate outputs
    generator = OutputGenerator(path)
    success = generator.generate_all()

    sys.exit(0 if success else 1)


if __name__ == '__main__':
    main()
