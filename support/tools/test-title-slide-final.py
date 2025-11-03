#!/usr/bin/env python3
"""
Final test: Full title slide with logos, title, subtitle, and presenter info.
"""

from pptx import Presentation
from pathlib import Path
import re
from datetime import datetime

def extract_metadata_from_markdown(md_file):
    """Extract title, subtitle, and other metadata from markdown file."""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    metadata = {
        'title': 'Untitled Presentation',
        'subtitle': '',
        'presenter': 'Presenter Name',
        'date': datetime.now().strftime('%B %d, %Y')
    }

    # Find first H1 header (title)
    match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
    if match:
        metadata['title'] = match.group(1).strip()

    # Find first H2 header (subtitle) - optional
    match = re.search(r'^##\s+(.+)$', content, re.MULTILINE)
    if match:
        metadata['subtitle'] = match.group(1).strip()

    return metadata

def replace_image_in_placeholder(slide_part, relationship_id, new_image_path, shape):
    """Replace image in a relationship."""
    try:
        # Get the current relationship
        image_rel = slide_part.rels[relationship_id]

        # Get the current image part
        old_image_part = image_rel.target_part

        # Read new image
        with open(new_image_path, 'rb') as f:
            image_blob = f.read()

        # Replace the blob in the existing image part
        old_image_part._blob = image_blob

        # Remove cropping (srcRect) to prevent image from being cropped
        pic = shape._element
        blipFill = pic.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}blipFill')
        if blipFill is not None:
            srcRect = blipFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srcRect')
            if srcRect is not None:
                blipFill.remove(srcRect)
                print(f"      Removed cropping")

        return True
    except Exception as e:
        print(f"      ❌ Error: {e}")
        return False

def create_title_slide_final():
    """Create complete title slide with all placeholders filled."""

    # Paths
    template_path = Path('support/doc-templates/powerpoint/EOFramework-Template-01.pptx')
    logo_path = Path('support/doc-templates/powerpoint')
    md_file = Path('solution-template/sample-provider/sample-category/sample-solution/presales/raw/executive-presentation.md')
    output_file = Path('test-title-slide-final.pptx')

    # Load template
    print(f"📂 Loading template: {template_path}")
    prs = Presentation(str(template_path))

    # Use the demo slide
    if len(prs.slides) == 0:
        print("❌ No demo slide found in template!")
        return

    slide = prs.slides[0]
    print(f"✅ Using demo slide: {slide.slide_layout.name}")

    # Extract metadata from markdown
    metadata = extract_metadata_from_markdown(md_file)
    print(f"📝 Metadata extracted:")
    print(f"   Title: {metadata['title']}")
    print(f"   Subtitle: {metadata['subtitle']}")
    print(f"   Presenter: {metadata['presenter']}")
    print(f"   Date: {metadata['date']}")
    print()

    # Map placeholder indices to content
    # Based on template analysis:
    # [10] = Picture (Logo 1)
    # [12] = Text (Title)
    # [13] = Picture (Logo 2)
    # [14] = Text (Subtitle)
    # [15] = Text (Presenter | Date)

    logo_mapping = {
        10: ('acme_logo_rectangle.png', 'Acme Inc'),
        13: ('zplus_consulting_logo.png', 'ZplusConsulting'),
    }

    text_mapping = {
        12: ('title', 'Title'),
        14: ('subtitle', 'Subtitle'),
        15: ('presenter_date', 'Presenter | Date'),
    }

    print("🎨 Filling placeholders:")

    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx

        # Handle logo placeholders
        if idx in logo_mapping:
            logo_file, logo_name = logo_mapping[idx]
            logo_full_path = logo_path / logo_file

            print(f"  [{idx}] {logo_name}:")

            if not logo_full_path.exists():
                print(f"      ❌ File not found: {logo_full_path}")
                continue

            # Find the relationship ID
            pic = shape._element
            blipFill = pic.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}blipFill')

            if blipFill is None:
                print(f"      ❌ No blipFill found")
                continue

            blip = blipFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')

            if blip is None:
                print(f"      ❌ No blip found")
                continue

            embed_rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')

            if not embed_rId:
                print(f"      ❌ No embed relationship found")
                continue

            print(f"      Relationship ID: {embed_rId}")

            # Replace the image
            success = replace_image_in_placeholder(slide.part, embed_rId, logo_full_path, shape)

            if success:
                print(f"      ✅ Replaced with {logo_file}")

        # Handle text placeholders
        elif idx in text_mapping:
            field_name, display_name = text_mapping[idx]

            if hasattr(shape, 'text_frame'):
                if field_name == 'title':
                    shape.text_frame.text = metadata['title']
                    print(f"  [{idx}] ✅ {display_name}: '{metadata['title']}'")

                elif field_name == 'subtitle':
                    shape.text_frame.text = metadata['subtitle']
                    print(f"  [{idx}] ✅ {display_name}: '{metadata['subtitle']}'")

                elif field_name == 'presenter_date':
                    presenter_text = f"{metadata['presenter']} | {metadata['date']}"
                    shape.text_frame.text = presenter_text
                    print(f"  [{idx}] ✅ {display_name}: '{presenter_text}'")

    # Save
    print(f"\n💾 Saving: {output_file}")
    prs.save(str(output_file))

    print(f"\n✨ Done! Open {output_file} to verify.")
    print(f"   Windows path: {output_file.absolute()}")
    print(f"\nThe slide should have:")
    print(f"  ✓ 2 logos (Acme Inc, ZplusConsulting)")
    print(f"  ✓ Title: {metadata['title']}")
    print(f"  ✓ Subtitle: {metadata['subtitle']}")
    print(f"  ✓ Presenter: {metadata['presenter']} | {metadata['date']}")

if __name__ == '__main__':
    create_title_slide_final()
