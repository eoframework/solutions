#!/usr/bin/env python3
"""
Test v5: Replace image relationship targets (like the Word example).
"""

from pptx import Presentation
from pathlib import Path
from pptx.parts.image import Image
import re

def extract_title_from_markdown(md_file):
    """Extract the first H1 title from markdown file."""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
    if match:
        return match.group(1).strip()
    return "Untitled Presentation"

def replace_placeholder_image(placeholder_shape, new_image_path):
    """Replace the image in a picture placeholder by updating relationships."""
    try:
        # Get the placeholder's image part relationship
        if not hasattr(placeholder_shape, '_element'):
            return False

        # Find the blipFill element (contains image reference)
        pic = placeholder_shape._element
        blipFill = pic.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}blipFill')

        if blipFill is None:
            print(f"      No blipFill found")
            return False

        # Find the blip (image reference)
        blip = blipFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')

        if blip is None:
            print(f"      No blip found")
            return False

        # Get the relationship ID
        embed_rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')

        if not embed_rId:
            print(f"      No embed relationship found")
            return False

        print(f"      Found image relationship: {embed_rId}")

        # Get the slide part
        slide_part = placeholder_shape.part

        # Get the current image relationship
        old_image_rel = slide_part.rels[embed_rId]
        print(f"      Old image: {old_image_rel.target_ref}")

        # Create new image part
        with open(new_image_path, 'rb') as f:
            image_data = f.read()

        image_part = Image.from_blob(image_data, content_type='image/png')

        # Replace the relationship target
        slide_part.rels[embed_rId]._target = image_part

        print(f"      ✅ Replaced with: {new_image_path.name}")
        return True

    except Exception as e:
        print(f"      ❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return False

def create_title_slide_v5():
    """Create title slide by replacing placeholder image relationships."""

    # Paths
    template_path = Path('support/doc-templates/powerpoint/EOFramework-Template-01.pptx')
    logo_path = Path('support/doc-templates/powerpoint')
    md_file = Path('solution-template/sample-provider/sample-category/sample-solution/presales/raw/executive-presentation.md')
    output_file = Path('test-title-slide-v5.pptx')

    # Load template
    print(f"📂 Loading template: {template_path}")
    prs = Presentation(str(template_path))

    # Create slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    print(f"✅ Created slide from layout: {title_layout.name}")

    # Extract title
    title_text = extract_title_from_markdown(md_file)
    print(f"📝 Title from markdown: {title_text}")
    print()

    # Define logo mappings
    logo_mapping = {
        1: ('acme_logo_rectangle.png', 'Acme Inc'),
        3: ('zplus_consulting_logo.png', 'ZplusConsulting'),
        4: ('eo-framework-logo-converted.png', 'EO Framework')
    }

    print("🎨 Replacing placeholder images via relationship update:")

    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx

        # Handle logo placeholders
        if idx in logo_mapping:
            logo_file, logo_name = logo_mapping[idx]
            logo_full_path = logo_path / logo_file

            print(f"  [{idx}] {logo_name}:")

            if logo_full_path.exists():
                success = replace_placeholder_image(shape, logo_full_path)
                if not success:
                    print(f"      Failed to replace")
            else:
                print(f"      ❌ File not found: {logo_full_path}")

        # Handle title placeholder
        elif idx == 12:
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = title_text
                print(f"  [{idx}] ✅ Title: Set to '{title_text}'")

    # Save
    print(f"\n💾 Saving: {output_file}")
    prs.save(str(output_file))

    print(f"\n✨ Done! Open {output_file} to verify.")
    print(f"   Windows path: {output_file.absolute()}")

if __name__ == '__main__':
    create_title_slide_v5()
