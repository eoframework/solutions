#!/usr/bin/env python3
"""
Comprehensive test of all 7 layouts using layouts-only approach.
"""

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pathlib import Path

def get_layout(prs, name_contains):
    """Get layout by partial name match."""
    for layout in prs.slide_layouts:
        if name_contains.lower() in layout.name.lower():
            return layout
    raise ValueError(f"Layout containing '{name_contains}' not found")

def insert_logos(slide, logo_path):
    """Insert logos into picture placeholders."""
    logos = {
        10: 'client_logo.png',
        13: 'consulting_company_logo.png',
        14: 'consulting_company_logo.png',  # Some layouts use idx 14 for logo
    }

    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx in logos and ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            logo_file = logo_path / logos[idx]
            if logo_file.exists():
                try:
                    ph.insert_picture(str(logo_file))
                    print(f"      ✅ Logo inserted at [{idx}]")
                except Exception as e:
                    print(f"      ❌ Logo failed at [{idx}]: {e}")

def test_all_layouts():
    """Test all 7 layouts."""

    template_path = Path('support/doc-templates/powerpoint/EOFramework-Template-01.pptx')
    logo_path = Path('support/doc-templates/powerpoint')
    output_file = Path('test-all-layouts.pptx')

    print(f"📂 Loading template: {template_path}")
    prs = Presentation(str(template_path))

    # Clear demo slides
    print(f"🗑️  Removing {len(prs.slides)} demo slides...")
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    print(f"✅ Starting fresh with {len(prs.slide_layouts)} layouts")
    print()

    # =================================================================
    # LAYOUT 0: EO Title Slide
    # =================================================================
    print("=" * 70)
    print("LAYOUT 0: EO Title Slide")
    print("=" * 70)

    layout0 = get_layout(prs, "Title")
    slide = prs.slides.add_slide(layout0)

    # [10] = Picture (Client Logo)
    # [12] = Text (Title)
    # [13] = Picture (Consulting Logo)
    # [14] = Text (Subtitle)
    # [15] = Text (Presenter | Date)

    slide.placeholders[12].text = "EO Framework Presentation"
    slide.placeholders[14].text = "Testing All Layouts"
    slide.placeholders[15].text = "Test User | November 02, 2025"
    print("  ✅ Text placeholders filled")

    insert_logos(slide, logo_path)
    print()

    # =================================================================
    # LAYOUT 1: EO Single Column
    # =================================================================
    print("=" * 70)
    print("LAYOUT 1: EO Single Column")
    print("=" * 70)

    layout1 = get_layout(prs, "Single")
    slide = prs.slides.add_slide(layout1)

    # [10] = Text (Title)
    # [11] = Text (Body/Content)
    # [13] = Picture (Logo)

    slide.placeholders[10].text = "Single Column Content"

    body = slide.placeholders[11]
    tf = body.text_frame
    tf.clear()
    tf.text = "First bullet point"
    for text in ["Second bullet point", "Third bullet point", "Fourth bullet point"]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0

    print("  ✅ Title and bullets filled")
    insert_logos(slide, logo_path)
    print()

    # =================================================================
    # LAYOUT 2: EO Two Column
    # =================================================================
    print("=" * 70)
    print("LAYOUT 2: EO Two Column")
    print("=" * 70)

    layout2 = get_layout(prs, "Two Column")
    slide = prs.slides.add_slide(layout2)

    # [10] = Text (Title)
    # [11] = Text (Left column header)
    # [12] = Text (Right column header)
    # [13] = Text (Left column content)
    # [14] = Picture (Logo)
    # [15] = Text (Right column content)

    slide.placeholders[10].text = "Two Column Comparison"
    slide.placeholders[11].text = "Traditional Approach"
    slide.placeholders[12].text = "Our Solution"

    left_content = slide.placeholders[13]
    left_content.text = "Manual processes\nHigh costs\nSlow deployment"

    right_content = slide.placeholders[15]
    right_content.text = "Automated workflow\nCost-effective\nRapid deployment"

    print("  ✅ Title and two columns filled")
    insert_logos(slide, logo_path)
    print()

    # =================================================================
    # LAYOUT 3: EO Key Points
    # =================================================================
    print("=" * 70)
    print("LAYOUT 3: EO Key Points")
    print("=" * 70)

    layout3 = get_layout(prs, "Key Points")
    slide = prs.slides.add_slide(layout3)

    # [10] = Text (Title)
    # [13] = Picture (Logo)
    # [14] = Text (Key points content)

    slide.placeholders[10].text = "Key Takeaways"

    content = slide.placeholders[14]
    tf = content.text_frame
    tf.clear()
    tf.text = "Critical point #1: Foundation for success"
    for text in ["Critical point #2: Proven methodology", "Critical point #3: Measurable results"]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0

    print("  ✅ Title and key points filled")
    insert_logos(slide, logo_path)
    print()

    # =================================================================
    # LAYOUT 4: EO Visual Content
    # =================================================================
    print("=" * 70)
    print("LAYOUT 4: EO Visual Content")
    print("=" * 70)

    layout4 = get_layout(prs, "Visual")
    slide = prs.slides.add_slide(layout4)

    # [10] = Text (Title)
    # [11] = Text (Subtitle/Section)
    # [13] = Text (Description)
    # [14] = Picture (Logo)
    # [15] = Picture (Main visual content)

    slide.placeholders[10].text = "Architecture Overview"
    slide.placeholders[11].text = "System Components"
    slide.placeholders[13].text = "The diagram shows the high-level architecture with key components and data flow."

    print("  ✅ Title and description filled")

    # Insert main visual (use logo as placeholder image)
    visual_logo = logo_path / 'eo-framework-logo-real.png'
    if visual_logo.exists():
        try:
            slide.placeholders[15].insert_picture(str(visual_logo))
            print("  ✅ Visual content image inserted")
        except Exception as e:
            print(f"  ⚠️  Visual image failed: {e}")

    insert_logos(slide, logo_path)
    print()

    # =================================================================
    # LAYOUT 5: EO Data Visualization
    # =================================================================
    print("=" * 70)
    print("LAYOUT 5: EO Data Visualization")
    print("=" * 70)

    layout5 = get_layout(prs, "Data")
    slide = prs.slides.add_slide(layout5)

    # [10] = Text (Title)
    # [11] = Text (Subtitle)
    # [13] = Text (Description/Insights)
    # [14] = Picture (Logo)
    # [15] = Chart (placeholder - skip for now)

    slide.placeholders[10].text = "Performance Metrics"
    slide.placeholders[11].text = "Q4 Results"
    slide.placeholders[13].text = "Key insights:\n• 40% improvement in efficiency\n• 60% reduction in costs\n• 95% customer satisfaction"

    print("  ✅ Title, subtitle and insights filled")
    print("  ⚠️  Chart placeholder [15] not populated (requires chart data)")
    insert_logos(slide, logo_path)
    print()

    # =================================================================
    # LAYOUT 6: EO Thank You
    # =================================================================
    print("=" * 70)
    print("LAYOUT 6: EO Thank You")
    print("=" * 70)

    layout6 = get_layout(prs, "Thank")
    slide = prs.slides.add_slide(layout6)

    # [10] = Picture (Top logo)
    # [12] = Text (Thank you message)
    # [13] = Picture (Bottom logo)
    # [14] = Text (Contact info)

    slide.placeholders[12].text = "Thank You"
    slide.placeholders[14].text = "Questions?\nwww.eoframework.com"

    print("  ✅ Thank you message and contact filled")
    insert_logos(slide, logo_path)
    print()

    # Save
    print("=" * 70)
    print(f"💾 Saving: {output_file}")
    prs.save(str(output_file))

    print()
    print("✨ Done! Open test-all-layouts.pptx to verify.")
    print(f"   Windows path: {output_file.absolute()}")
    print()
    print("Generated slides:")
    print("  1. Title Slide - with logos and metadata")
    print("  2. Single Column - with bullets")
    print("  3. Two Column - with comparison")
    print("  4. Key Points - with emphasized points")
    print("  5. Visual Content - with image")
    print("  6. Data Visualization - with insights")
    print("  7. Thank You - with closing")

if __name__ == '__main__':
    test_all_layouts()
