#!/usr/bin/env python3
"""
Test script to create a title slide with logos and title from markdown.
"""

from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
import re

def extract_title_from_markdown(md_file):
    """Extract the first H1 title from markdown file."""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Find first H1 header
    match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
    if match:
        return match.group(1).strip()

    return "Untitled Presentation"

def create_title_slide_test():
    """Create a test presentation with title slide."""

    # Paths
    template_path = Path('support/doc-templates/powerpoint/EOFramework-Template-01.pptx')
    logo_path = Path('support/doc-templates/powerpoint')
    md_file = Path('solution-template/sample-provider/sample-category/sample-solution/presales/raw/executive-presentation.md')
    output_file = Path('test-title-slide.pptx')

    # Load template
    print(f"📂 Loading template: {template_path}")
    prs = Presentation(str(template_path))

    # Remove all existing demo slides
    print(f"🗑️  Removing {len(prs.slides)} demo slides from template...")
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    print(f"✅ Demo slides removed")

    # Get title layout
    title_layout = prs.slide_layouts[0]
    print(f"✅ Using layout: {title_layout.name}")

    # Create slide
    slide = prs.slides.add_slide(title_layout)
    print(f"✅ Created slide")

    # Extract title from markdown
    title_text = extract_title_from_markdown(md_file)
    print(f"📝 Title from markdown: {title_text}")

    # Fill placeholders
    print("\n🎨 Filling placeholders:")

    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        name = shape.name
        ph_type = str(shape.placeholder_format.type).split('.')[-1]

        print(f"  [{idx}] {name} ({ph_type})")

        # Customer Logo (idx 1)
        if idx == 1:
            logo_file = logo_path / 'acme_logo_rectangle.png'
            if logo_file.exists():
                try:
                    # Try picture fill method
                    pic = shape.insert_picture(str(logo_file))
                    print(f"      ✅ Inserted (method 1): acme_logo_rectangle.png")
                except Exception as e:
                    print(f"      ⚠️  Method 1 failed: {e}")
                    try:
                        # Alternative: direct picture element replacement
                        from pptx.util import Inches
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        shape.element.getparent().remove(shape.element)
                        pic = slide.shapes.add_picture(str(logo_file), left, top, width, height)
                        print(f"      ✅ Inserted (method 2): acme_logo_rectangle.png")
                    except Exception as e2:
                        print(f"      ❌ All methods failed: {e2}")
            else:
                print(f"      ❌ Logo not found: {logo_file}")

        # Consulting Company Logo (idx 3)
        elif idx == 3:
            logo_file = logo_path / 'zplus_consulting_logo.png'
            if logo_file.exists():
                try:
                    pic = shape.insert_picture(str(logo_file))
                    print(f"      ✅ Inserted (method 1): zplus_consulting_logo.png")
                except Exception as e:
                    print(f"      ⚠️  Method 1 failed: {e}")
                    try:
                        from pptx.util import Inches
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        shape.element.getparent().remove(shape.element)
                        pic = slide.shapes.add_picture(str(logo_file), left, top, width, height)
                        print(f"      ✅ Inserted (method 2): zplus_consulting_logo.png")
                    except Exception as e2:
                        print(f"      ❌ All methods failed: {e2}")
            else:
                print(f"      ❌ Logo not found: {logo_file}")

        # EO Framework Logo (idx 4)
        elif idx == 4:
            logo_file = logo_path / 'eo-framework-logo-converted.png'
            if logo_file.exists():
                try:
                    pic = shape.insert_picture(str(logo_file))
                    print(f"      ✅ Inserted (method 1): eo-framework-logo.png")
                except Exception as e:
                    print(f"      ⚠️  Method 1 failed: {e}")
                    try:
                        from pptx.util import Inches
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        shape.element.getparent().remove(shape.element)
                        pic = slide.shapes.add_picture(str(logo_file), left, top, width, height)
                        print(f"      ✅ Inserted (method 2): eo-framework-logo.png")
                    except Exception as e2:
                        print(f"      ❌ All methods failed: {e2}")
            else:
                print(f"      ❌ Logo not found: {logo_file}")

        # Text Placeholder (idx 12) - Title
        elif idx == 12:
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = title_text
                print(f"      ✅ Set text: {title_text}")

    # Save
    print(f"\n💾 Saving: {output_file}")
    prs.save(str(output_file))

    print(f"\n✨ Done! Open {output_file} to verify.")
    print(f"   Windows path: {output_file.absolute()}")

if __name__ == '__main__':
    create_title_slide_test()
