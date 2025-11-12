#!/usr/bin/env python3
"""
Build professional output documents from raw source files.

Builds:
- Markdown (.md) → Word (.docx) or PowerPoint (.pptx)
- CSV (.csv) → Excel (.xlsx) with formatting

Usage:
    # Generate all files (default)
    python3 solution-doc-builder.py --path solutions/provider/category/solution

    # Generate specific file types
    python3 solution-doc-builder.py --path <path> --type excel
    python3 solution-doc-builder.py --path <path> --type word pptx

    # Generate specific file
    python3 solution-doc-builder.py --path <path> --file statement-of-work.md

    # Generate from specific directory
    python3 solution-doc-builder.py --path <path> --dir presales

    # Dry run to preview
    python3 solution-doc-builder.py --path <path> --dry-run
"""

import argparse
import os
import sys
from pathlib import Path
import pandas as pd
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.styles.numbers import FORMAT_CURRENCY_USD_SIMPLE
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.utils import get_column_letter
import numpy as np
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER
from markdown_it import MarkdownIt
from datetime import datetime
import re


class OutputGenerator:
    """Build professional output files from raw source files."""

    def __init__(self, base_path, type_filter=None, file_filter=None,
                 dir_filter=None, force=False, dry_run=False, quiet=False, verbose=False):
        self.base_path = Path(base_path)
        self.type_filter = set(type_filter) if type_filter else None
        self.file_filter = set(file_filter) if file_filter else None
        self.dir_filter = dir_filter
        self.force = force
        self.dry_run = dry_run
        self.quiet = quiet
        self.verbose = verbose
        self.stats = {
            'csv_to_xlsx': 0,
            'md_to_docx': 0,
            'md_to_pptx': 0,
            'skipped': 0,
            'errors': []
        }

    def should_process_file(self, source_file, output_type, raw_dir):
        """
        Determine if a file should be processed based on filters.

        Returns True if:
        - No filters specified (process everything)
        - File matches all specified filters
        """
        # No filters = process everything (backward compatible)
        if not any([self.type_filter, self.file_filter, self.dir_filter]):
            return True

        # Check directory filter
        if self.dir_filter:
            dir_name = raw_dir.parent.name  # e.g., 'presales' or 'delivery'
            if dir_name != self.dir_filter:
                return False

        # Check file filter (match source filename)
        if self.file_filter:
            if source_file.name not in self.file_filter:
                return False

        # Check type filter
        if self.type_filter:
            if output_type not in self.type_filter:
                return False

        return True

    def generate_all(self):
        """Generate all output files from raw directories."""
        if not self.quiet:
            print(f"🔍 Scanning: {self.base_path}")
            if self.dry_run:
                print("🔍 DRY RUN MODE - No files will be generated\n")

        # Find all raw directories
        raw_dirs = list(self.base_path.rglob('raw'))

        if not raw_dirs:
            if not self.quiet:
                print("❌ No 'raw' directories found.")
            return False

        if not self.quiet:
            print(f"📁 Found {len(raw_dirs)} raw directories\n")

        for raw_dir in raw_dirs:
            self.process_raw_directory(raw_dir)

        self.print_summary()
        return len(self.stats['errors']) == 0

    def process_raw_directory(self, raw_dir):
        """Process all files in a raw directory."""
        output_dir = raw_dir.parent

        if not self.quiet:
            print(f"📂 Processing: {raw_dir.relative_to(self.base_path)}")

        # Process CSV files
        for csv_file in raw_dir.glob('*.csv'):
            if self.should_process_file(csv_file, 'excel', raw_dir):
                self.convert_csv_to_xlsx(csv_file, output_dir)
            elif self.verbose:
                print(f"  ⏭️  Skipped (filtered): {csv_file.name}")
                self.stats['skipped'] += 1

        # Process Markdown files
        for md_file in raw_dir.glob('*.md'):
            # Check if file should be converted to PowerPoint (presentation or briefing)
            if 'presentation' in md_file.stem.lower() or 'briefing' in md_file.stem.lower():
                if self.should_process_file(md_file, 'pptx', raw_dir):
                    self.convert_md_to_pptx(md_file, output_dir)
                elif self.verbose:
                    print(f"  ⏭️  Skipped (filtered): {md_file.name}")
                    self.stats['skipped'] += 1
            else:
                if self.should_process_file(md_file, 'word', raw_dir):
                    self.convert_md_to_docx(md_file, output_dir)
                elif self.verbose:
                    print(f"  ⏭️  Skipped (filtered): {md_file.name}")
                    self.stats['skipped'] += 1

        if not self.quiet:
            print()

    def _is_currency_column(self, column_name, sample_values):
        """Detect if a column contains currency values."""
        # Check column name for currency keywords
        currency_keywords = ['cost', 'price', 'amount', 'total', 'budget', 'fee', 'rate',
                            'salary', 'revenue', 'expense', 'payment', 'charge', 'value',
                            'investment', 'profit', 'loss', 'income', 'worth', 'usd', '$']

        column_lower = column_name.lower()
        if any(keyword in column_lower for keyword in currency_keywords):
            return True

        # Check sample values for $ signs
        for val in sample_values[:10]:  # Check first 10 values
            if val and isinstance(val, str) and '$' in val:
                return True

        return False

    def _convert_currency_to_number(self, value):
        """Convert currency string to numeric value."""
        if pd.isna(value) or value == '':
            return None

        if isinstance(value, (int, float)):
            return float(value)

        # Remove currency symbols and formatting
        if isinstance(value, str):
            # Remove $, commas, spaces, parentheses (for negative)
            clean = value.replace('$', '').replace(',', '').replace(' ', '').strip()

            # Handle parentheses for negative numbers
            if clean.startswith('(') and clean.endswith(')'):
                clean = '-' + clean[1:-1]

            try:
                return float(clean)
            except (ValueError, AttributeError):
                return None

        return None

    def _generate_sheet_name(self, filename):
        """Generate appropriate sheet name from filename."""
        # Common abbreviations and mappings
        name_mappings = {
            'level-of-effort-estimate': 'LOE',
            'requirements-questionnaire': 'Requirements',
            'roi-calculator': 'ROI',
            'communication-plan': 'Comm Plan',
            'configuration': 'Configuration',
            'project-plan': 'Project Plan',
            'requirements': 'Requirements',
            'roles': 'Roles',
            'test-plan': 'Test Plan',
            'training-plan': 'Training Plan',
        }

        # Check for exact match
        if filename.lower() in name_mappings:
            return name_mappings[filename.lower()]

        # Generate from filename (title case, remove dashes/underscores)
        sheet_name = filename.replace('-', ' ').replace('_', ' ').title()

        # Limit to 31 characters (Excel sheet name limit)
        if len(sheet_name) > 31:
            # Try to abbreviate
            words = sheet_name.split()
            if len(words) > 1:
                # Use first letter of each word except last
                abbrev = ''.join([w[0] for w in words[:-1]]) + ' ' + words[-1]
                sheet_name = abbrev[:31]
            else:
                sheet_name = sheet_name[:31]

        return sheet_name

    def convert_loe_to_xlsx(self, csv_file, output_dir):
        """Convert multi-section LOE CSV to multi-sheet Excel with branded template."""
        import csv as csvmod
        from io import StringIO

        try:
            output_file = output_dir / f"{csv_file.stem}.xlsx"

            # Load branded template
            template_path = Path(__file__).parent.parent / 'doc-templates' / 'excel' / 'EOFramework-Excel-Template-01.xlsx'
            if not template_path.exists():
                raise FileNotFoundError(f"Excel template not found: {template_path}")

            wb = load_workbook(str(template_path))

            # Update Cover sheet
            if 'Cover' in wb.sheetnames:
                cover = wb['Cover']
                cover['B4'] = "Level Of Effort Estimate"
                cover['B5'] = 'Presales Document' if 'presales' in str(csv_file) else 'Project Planning Document'
                cover['C7'] = datetime.now().strftime('%B %d, %Y')

                # Extract solution name from path
                solution_name = 'Sample Solution'
                try:
                    parts = csv_file.parts
                    if 'solution-template' in parts:
                        solution_name = 'Sample Solution'
                    else:
                        # Try to extract from path
                        for i, part in enumerate(parts):
                            if part in ['presales', 'delivery'] and i > 0:
                                solution_name = parts[i-1].replace('-', ' ').title()
                                break
                except:
                    pass

                cover['C8'] = solution_name
                cover['C9'] = "[Customer Name]"

            # Parse multi-section CSV
            sections = {'assumptions': [], 'multipliers': [], 'estimate': []}
            current_section = None

            with open(csv_file, 'r', encoding='utf-8') as f:
                for line in f:
                    line = line.rstrip('\n')
                    if line.startswith('# SCOPE ASSUMPTIONS'):
                        current_section = 'assumptions'
                        continue
                    elif line.startswith('# ESTIMATE SETTINGS') or line.startswith('# MULTIPLIER SETTINGS'):
                        current_section = 'multipliers'
                        continue
                    elif line.startswith('# EFFORT ESTIMATE'):
                        current_section = 'estimate'
                        continue
                    elif line.startswith('#') or not line.strip():
                        continue
                    if current_section:
                        sections[current_section].append(line)

            # Standard styling from template
            header_font = Font(name='Calibri', size=12, bold=True, color='FFFFFF')
            header_fill = PatternFill(start_color='808080', end_color='808080', fill_type='solid')
            data_font = Font(name='Calibri', size=12, color='000000')
            thin_border = Border(left=Side(style='thin'), right=Side(style='thin'),
                               top=Side(style='thin'), bottom=Side(style='thin'))

            # Sheet 1: Scope Assumptions
            ws_assumptions = wb['Data']
            ws_assumptions.title = 'Scope Assumptions'
            ws_assumptions.delete_rows(2, ws_assumptions.max_row)

            reader = csvmod.reader(StringIO('\n'.join(sections['assumptions'])))
            rows = list(reader)
            for r_idx, row in enumerate(rows, 1):
                for c_idx, value in enumerate(row, 1):
                    cell = ws_assumptions.cell(row=r_idx, column=c_idx, value=value)
                    if r_idx == 1:
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = Alignment(horizontal='center', vertical='center')
                    else:
                        cell.font = data_font
                        cell.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
                    cell.border = thin_border

            ws_assumptions.column_dimensions['A'].width = 25
            ws_assumptions.column_dimensions['B'].width = 35
            ws_assumptions.column_dimensions['C'].width = 30
            ws_assumptions.column_dimensions['D'].width = 30
            ws_assumptions.column_dimensions['E'].width = 30
            ws_assumptions.column_dimensions['F'].width = 40

            # Sheet 2: Estimate Settings with instructions
            ws_mult = wb.create_sheet('Estimate Settings', 2)

            # Add instructions
            instructions = [
                ["INSTRUCTIONS: How to Use Estimate Settings"],
                [""],
                ["This sheet contains phase-based multiplier variables that control effort scaling for each project phase."],
                ["The EO Sales Engineer updates these multipliers based on customer scope compared to baseline assumptions."],
                [""],
                ["How It Works:"],
                ["1. Review customer requirements from Discovery Questionnaire"],
                ["2. Compare customer scope to Baseline Assumptions (see 'Scope Assumptions' sheet)"],
                ["3. Update YELLOW HIGHLIGHTED multiplier values below (Column C) for each phase"],
                ["4. All tasks in that phase on the 'LOE' sheet automatically update their Actual Hours"],
                ["5. Total cost recalculates automatically"],
                [""],
                ["Multiplier Guidelines:"],
                ["• 1.0 = Matches baseline exactly (Small scope)"],
                ["• 1.5-2.0 = Medium scope (moderately above baseline)"],
                ["• 2.0-3.0 = Large scope (significantly above baseline)"],
                ["• 3.0+ = Very large/complex (substantially above baseline, document rationale)"],
                [""],
                ["Example: Customer needs comprehensive testing due to compliance → Set TESTING_MX to 2.0"],
                ["Example: Complex architecture requiring more technical oversight → Set LEADERSHIP_MX to 1.5"],
                [""],
                ["ESTIMATE SETTINGS - Update phase multipliers in Column C below:"],
            ]

            for r_idx, row in enumerate(instructions, 1):
                for c_idx, value in enumerate(row, 1):
                    cell = ws_mult.cell(row=r_idx, column=c_idx, value=value)
                    if r_idx == 1:
                        cell.font = Font(name='Calibri', size=14, bold=True, color='FFFFFF')
                        cell.fill = PatternFill(start_color='808080', end_color='808080', fill_type='solid')
                    elif "How It Works:" in str(value) or "Multiplier Guidelines:" in str(value):
                        cell.font = Font(name='Calibri', size=12, bold=True)
                    elif "ESTIMATE SETTINGS" in str(value) or "MULTIPLIER SETTINGS" in str(value):
                        cell.font = Font(name='Calibri', size=12, bold=True, color='FFFFFF')
                        cell.fill = PatternFill(start_color='808080', end_color='808080', fill_type='solid')
                    else:
                        cell.font = data_font

            ws_mult.merge_cells('A1:D1')
            for row_num in range(3, 21):
                if ws_mult.cell(row=row_num, column=1).value:
                    ws_mult.merge_cells(f'A{row_num}:D{row_num}')
            ws_mult.merge_cells('A22:D22')

            # Add multiplier data
            reader = csvmod.reader(StringIO('\n'.join(sections['multipliers'])))
            rows = list(reader)
            start_row = 23
            for r_idx, row in enumerate(rows, start_row):
                for c_idx, value in enumerate(row, 1):
                    cell = ws_mult.cell(row=r_idx, column=c_idx, value=value)
                    if r_idx == start_row:
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = Alignment(horizontal='center', vertical='center')
                    else:
                        cell.font = data_font
                        if c_idx == 3:  # Multiplier column - highlight and convert to number
                            cell.fill = PatternFill(start_color='FFF2CC', end_color='FFF2CC', fill_type='solid')
                            cell.font = Font(name='Calibri', size=12, bold=True)
                            try:
                                cell.value = float(value)
                            except:
                                pass
                    cell.border = thin_border

            ws_mult.column_dimensions['A'].width = 25
            ws_mult.column_dimensions['B'].width = 20
            ws_mult.column_dimensions['C'].width = 15
            ws_mult.column_dimensions['D'].width = 80

            # Sheet 3: LOE
            ws_est = wb.create_sheet('LOE', 3)
            reader = csvmod.reader(StringIO('\n'.join(sections['estimate'])))
            rows = list(reader)

            for r_idx, row in enumerate(rows, 1):
                for c_idx, value in enumerate(row, 1):
                    cell = ws_est.cell(row=r_idx, column=c_idx)

                    # Write formulas as formulas and adjust row numbers (CSV uses row 33+, Excel uses row 2+)
                    if isinstance(value, str) and value.startswith('='):
                        # Adjust row references: subtract 31 to convert CSV row numbers to Excel row numbers
                        # CSV row 33 becomes Excel row 2, CSV row 34 becomes Excel row 3, etc.
                        # BUT: Don't adjust cross-sheet references (containing !)
                        adjusted_value = value
                        import re

                        # Only adjust row numbers that are NOT part of cross-sheet references
                        # Split by '!' to separate cross-sheet refs
                        if '!' in value:
                            # Has cross-sheet reference - only adjust the local part
                            parts = value.split('!')
                            # Adjust only references in the last part (after the last !)
                            for i in range(len(parts)):
                                if i < len(parts) - 1:
                                    # Don't adjust cross-sheet references
                                    continue
                                else:
                                    # Adjust local references
                                    parts[i] = re.sub(r'(?<!\w)\$([A-Z])(\$)?(\d+)',
                                                    lambda m: f"${m.group(1)}{m.group(2) or ''}{int(m.group(3)) - 31 if int(m.group(3)) >= 33 else m.group(3)}",
                                                    parts[i])
                            adjusted_value = '!'.join(parts)
                        else:
                            # No cross-sheet reference - adjust all row numbers
                            adjusted_value = re.sub(r'\$([A-Z])(\$)?(\d+)',
                                                  lambda m: f"${m.group(1)}{m.group(2) or ''}{int(m.group(3)) - 31 if int(m.group(3)) >= 33 else m.group(3)}",
                                                  adjusted_value)
                        cell.value = adjusted_value
                    else:
                        # Column E (Base Hours) - convert to number
                        if c_idx == 5 and r_idx > 1:
                            try:
                                cell.value = float(value) if value and value.strip() else value
                                cell.number_format = '0'
                            except:
                                cell.value = value
                        # Column H (Rate) - convert to number with currency format
                        elif c_idx == 8 and r_idx > 1:
                            try:
                                # Remove $ and convert to number
                                rate_value = value.replace('$', '').replace(',', '').strip() if isinstance(value, str) else value
                                cell.value = float(rate_value) if rate_value else value
                                cell.number_format = '$#,##0'
                            except:
                                cell.value = value
                        else:
                            cell.value = value

                    if r_idx == 1:
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = Alignment(horizontal='center', vertical='center')
                    else:
                        cell.font = data_font
                        cell.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
                        if c_idx == 6:  # Effort Multiplier column
                            cell.fill = PatternFill(start_color='E8F4F8', end_color='E8F4F8', fill_type='solid')
                    cell.border = thin_border

            ws_est.column_dimensions['A'].width = 12
            ws_est.column_dimensions['B'].width = 25
            ws_est.column_dimensions['C'].width = 50
            ws_est.column_dimensions['D'].width = 20
            ws_est.column_dimensions['E'].width = 12
            ws_est.column_dimensions['F'].width = 18
            ws_est.column_dimensions['G'].width = 15
            ws_est.column_dimensions['H'].width = 12
            ws_est.column_dimensions['I'].width = 15
            ws_est.column_dimensions['J'].width = 25
            ws_est.column_dimensions['K'].width = 60

            ws_est.freeze_panes = 'A2'

            wb.save(str(output_file))
            print(f"  ✅ {csv_file.name} → {output_file.name}")
            self.stats['csv_to_xlsx'] += 1
            return True

        except Exception as e:
            error_msg = f"Error converting {csv_file.name}: {str(e)}"
            print(f"  ❌ {error_msg}")
            self.stats['errors'].append(error_msg)
            return False

    def convert_discovery_to_xlsx(self, csv_file, output_dir):
        """Convert Discovery Questionnaire CSV to multi-sheet Excel with category sheets."""
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, Protection
        from openpyxl.utils.dataframe import dataframe_to_rows
        import pandas as pd

        try:
            output_file = output_dir / f"{csv_file.stem}.xlsx"

            # Load branded template
            template_path = Path(__file__).parent.parent / 'doc-templates' / 'excel' / 'EOFramework-Excel-Template-01.xlsx'

            if not template_path.exists():
                raise FileNotFoundError(f"Excel template not found: {template_path}")

            wb = load_workbook(str(template_path))

            # Update Cover sheet
            if 'Cover' in wb.sheetnames:
                cover = wb['Cover']
                cover['B4'] = "Discovery Questionnaire"
                cover['B5'] = 'Presales Document' if 'presales' in str(csv_file) else 'Project Planning Document'
                cover['C7'] = datetime.now().strftime('%B %d, %Y')

                # Extract solution name
                solution_name = 'Sample Solution'
                try:
                    parts = csv_file.parts
                    if 'solution-template' in parts:
                        solution_name = 'Sample Solution'
                    else:
                        for i, part in enumerate(parts):
                            if part in ['presales', 'delivery'] and i > 0:
                                solution_name = parts[i-1].replace('-', ' ').title()
                                break
                except:
                    pass

                cover['C8'] = solution_name
                cover['C9'] = "[Customer Name]"

            # Read CSV and group by category
            df = pd.read_csv(csv_file)
            categories = df['Category'].unique()

            # Extract template styles from Data sheet
            data_sheet = wb['Data']
            header_cell = data_sheet['A1']

            # Extract font with defaults
            if header_cell.font and header_cell.font.name:
                header_font = Font(name=header_cell.font.name, size=header_cell.font.size or 12)
            else:
                header_font = Font(name='Times New Roman', size=12)

            # Extract body font
            if data_sheet.max_row > 1:
                body_cell = data_sheet['A2']
                if body_cell.font and body_cell.font.name:
                    body_font = Font(name=body_cell.font.name, size=body_cell.font.size or 12)
                else:
                    body_font = Font(name='Times New Roman', size=12)
            else:
                body_font = Font(name='Times New Roman', size=12)

            # Define borders
            thin_border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )

            # Remove the Data sheet (we'll create category-specific sheets)
            if 'Data' in wb.sheetnames:
                wb.remove(wb['Data'])

            # Create a sheet for each category (starting at position 2, after Cover and All Categories)
            category_sheet_names = {}
            for idx, category in enumerate(categories):
                # Create safe sheet name (max 31 chars, no special characters)
                sheet_name = category[:28] + "..." if len(category) > 31 else category
                sheet_name = sheet_name.replace('/', '-').replace('\\', '-').replace('[', '').replace(']', '')

                category_sheet_names[category] = sheet_name
                # Position starts at 2 (0=Cover, 1=All Categories, 2+=category sheets)
                ws = wb.create_sheet(sheet_name, idx + 2)

                # Filter data for this category
                category_df = df[df['Category'] == category].copy()

                # Write headers
                headers = df.columns.tolist()
                for c_idx, header in enumerate(headers, 1):
                    cell = ws.cell(row=1, column=c_idx, value=header)
                    cell.font = Font(name=header_font.name, size=header_font.size, bold=True, color='FFFFFF')
                    cell.fill = PatternFill(start_color='808080', end_color='808080', fill_type='solid')
                    cell.alignment = Alignment(horizontal='left', vertical='center', wrap_text=True)
                    cell.border = thin_border

                # Write data
                for r_idx, row in enumerate(category_df.values, 2):
                    for c_idx, value in enumerate(row, 1):
                        # Convert empty strings to None so they display as blank
                        cell_value = None if (isinstance(value, str) and value.strip() == '') or pd.isna(value) else value
                        cell = ws.cell(row=r_idx, column=c_idx, value=cell_value)
                        cell.font = Font(name=body_font.name, size=body_font.size)
                        cell.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
                        cell.border = thin_border

                        # Highlight Customer Response column (column 6)
                        if c_idx == 6:
                            cell.fill = PatternFill(start_color='FFF2CC', end_color='FFF2CC', fill_type='solid')

                # Set column widths
                ws.column_dimensions['A'].width = 12  # Question ID
                ws.column_dimensions['B'].width = 20  # Category
                ws.column_dimensions['C'].width = 25  # Sub-Category
                ws.column_dimensions['D'].width = 50  # Question
                ws.column_dimensions['E'].width = 40  # Guidance
                ws.column_dimensions['F'].width = 40  # Customer Response
                ws.column_dimensions['G'].width = 30  # Notes

                # Freeze panes
                ws.freeze_panes = 'A2'

            # Create consolidated "All Categories (Read Only)" sheet at position 1 (right after Cover)
            consolidated_sheet = wb.create_sheet("All Categories (Read Only)", 1)

            # Calculate total questions
            total_questions = len(df)

            # Row 1: Stats Labels
            stats_labels = ['Total Questions', 'Answered', 'Remaining', '% Complete']
            for c_idx, label in enumerate(stats_labels, 1):
                cell = consolidated_sheet.cell(row=1, column=c_idx, value=label)
                cell.font = Font(name=header_font.name, size=11, bold=True)
                cell.fill = PatternFill(start_color='D1ECF1', end_color='D1ECF1', fill_type='solid')
                cell.alignment = Alignment(horizontal='center', vertical='center')
                cell.border = thin_border

            # Row 2: Stats Values
            stats_values = [
                total_questions,
                f'=SUMPRODUCT(--(LEN(F5:F{4 + total_questions})>0))',  # Count non-empty Customer Response cells
                f'=A2-B2',  # Remaining = Total - Answered
                f'=B2/A2'   # % Complete = Answered / Total
            ]

            for c_idx, value in enumerate(stats_values, 1):
                cell = consolidated_sheet.cell(row=2, column=c_idx, value=value)
                cell.font = Font(name=header_font.name, size=12, bold=True)
                cell.fill = PatternFill(start_color='E8F5E9', end_color='E8F5E9', fill_type='solid')
                cell.alignment = Alignment(horizontal='center', vertical='center')
                cell.border = thin_border

                # Format percentage for column 4
                if c_idx == 4:
                    cell.number_format = '0%'

            # Row 3: Blank separator (no styling needed)

            # Row 4: Write headers
            for c_idx, header in enumerate(df.columns.tolist(), 1):
                cell = consolidated_sheet.cell(row=4, column=c_idx, value=header)
                cell.font = Font(name=header_font.name, size=header_font.size, bold=True, color='FFFFFF')
                cell.fill = PatternFill(start_color='4472C4', end_color='4472C4', fill_type='solid')
                cell.alignment = Alignment(horizontal='left', vertical='center', wrap_text=True)
                cell.border = thin_border

            # Write formulas to pull data from category sheets (starting at row 5)
            current_row = 5
            for category in categories:
                sheet_name = category_sheet_names[category]
                category_df = df[df['Category'] == category]
                num_rows = len(category_df)

                for row_offset in range(num_rows):
                    source_row = row_offset + 2  # +2 because source has header at row 1
                    for c_idx in range(1, len(df.columns) + 1):
                        col_letter = chr(64 + c_idx)  # A=65, B=66, etc.
                        # Use IF to check for empty strings and show blank instead of 0
                        formula = f"=IF('{sheet_name}'!{col_letter}{source_row}=\"\",\"\",'{sheet_name}'!{col_letter}{source_row})"
                        cell = consolidated_sheet.cell(row=current_row, column=c_idx, value=formula)
                        cell.font = Font(name=body_font.name, size=body_font.size)
                        cell.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
                        cell.border = thin_border

                        # Highlight Customer Response column
                        if c_idx == 6:
                            cell.fill = PatternFill(start_color='E8F4F8', end_color='E8F4F8', fill_type='solid')

                    current_row += 1

            # Set column widths for consolidated sheet
            consolidated_sheet.column_dimensions['A'].width = 12
            consolidated_sheet.column_dimensions['B'].width = 20
            consolidated_sheet.column_dimensions['C'].width = 25
            consolidated_sheet.column_dimensions['D'].width = 50
            consolidated_sheet.column_dimensions['E'].width = 40
            consolidated_sheet.column_dimensions['F'].width = 40
            consolidated_sheet.column_dimensions['G'].width = 30

            # Freeze panes at row 5 (keeps stats and headers visible when scrolling)
            consolidated_sheet.freeze_panes = 'A5'

            # Protect consolidated sheet (read-only, no password)
            consolidated_sheet.protection.sheet = True
            consolidated_sheet.protection.enable()

            # Save workbook
            wb.save(str(output_file))
            print(f"  ✅ {csv_file.name} → {output_file.name}")
            self.stats['csv_to_xlsx'] += 1
            return True

        except Exception as e:
            import traceback
            error_msg = f"Error converting {csv_file.name}: {str(e)}"
            print(f"  ❌ {error_msg}")
            print("Full traceback:")
            traceback.print_exc()
            self.stats['errors'].append(error_msg)
            return False

    def convert_csv_to_xlsx(self, csv_file, output_dir):
        """Convert CSV to styled Excel file using branded template."""
        try:
            output_file = output_dir / f"{csv_file.stem}.xlsx"

            # Dry-run mode - just show what would be generated
            if self.dry_run:
                print(f"  🔍 Would generate: {csv_file.name} → {output_file.name}")
                self.stats['csv_to_xlsx'] += 1
                return

            # Special handling for level-of-effort-estimate (multi-section CSV)
            if 'level-of-effort' in csv_file.stem.lower():
                return self.convert_loe_to_xlsx(csv_file, output_dir)

            # Special handling for discovery-questionnaire (multi-sheet by category)
            if 'discovery-questionnaire' in csv_file.stem.lower():
                return self.convert_discovery_to_xlsx(csv_file, output_dir)

            # Special handling for cost-breakdown (5-tab workbook)
            if 'cost-breakdown' in csv_file.stem.lower():
                return self.convert_cost_breakdown_to_xlsx(csv_file, output_dir)

            # Load branded template
            template_path = Path(__file__).parent.parent / 'doc-templates' / 'excel' / 'EOFramework-Excel-Template-01.xlsx'

            if not template_path.exists():
                raise FileNotFoundError(f"Excel template not found: {template_path}")

            wb = load_workbook(str(template_path))

            # Read CSV
            df = pd.read_csv(csv_file)

            # Detect currency columns
            currency_columns = []
            for col_idx, col_name in enumerate(df.columns):
                if self._is_currency_column(col_name, df[col_name].dropna().tolist()):
                    currency_columns.append(col_idx)
                    # Convert currency strings to numbers
                    df[col_name] = df[col_name].apply(self._convert_currency_to_number)

            # Update Cover sheet with metadata
            if 'Cover' in wb.sheetnames:
                cover = wb['Cover']

                # Update document title
                cover['B4'] = f"{csv_file.stem.replace('-', ' ').replace('_', ' ').title()}"

                # Update document type
                if 'presales' in str(csv_file):
                    doc_type = 'Presales Document'
                elif 'delivery' in str(csv_file):
                    doc_type = 'Delivery Document'
                else:
                    doc_type = 'Reference Document'
                cover['B5'] = doc_type

                # Update metadata (using new compact layout B7-C10)
                cover['C7'] = datetime.now().strftime('%B %d, %Y')

                # Extract solution name from path
                solution_name = 'EO Framework Solution'
                try:
                    parts = csv_file.parts
                    if 'solution-template' in parts:
                        solution_name = 'Sample Solution'
                    else:
                        # Try to extract from path
                        for i, part in enumerate(parts):
                            if part in ['presales', 'delivery'] and i > 0:
                                solution_name = parts[i-1].replace('-', ' ').title()
                                break
                except:
                    pass

                cover['C8'] = solution_name
                cover['C9'] = f"Data table for {csv_file.stem.replace('-', ' ')}"

            # Get Data sheet and rename it based on file name
            ws = wb['Data']

            # Generate appropriate sheet name from file name
            sheet_name = self._generate_sheet_name(csv_file.stem)
            ws.title = sheet_name

            # Clear sample data (rows 2 onwards)
            ws.delete_rows(2, ws.max_row)

            # Define styling
            header_font = Font(name='Calibri', size=12, bold=True, color='FFFFFF')  # Size 12
            header_fill = PatternFill(start_color='808080', end_color='808080', fill_type='solid')  # Gray header
            header_alignment = Alignment(horizontal='center', vertical='center')

            data_font = Font(name='Calibri', size=12, color='000000')  # Size 12
            data_alignment = Alignment(horizontal='left', vertical='center', wrap_text=True, indent=1)  # Center vertical + indent for padding

            # Alternating row colors for better readability
            even_fill = PatternFill(start_color='F2F2F2', end_color='F2F2F2', fill_type='solid')  # Light gray
            odd_fill = PatternFill(start_color='FFFFFF', end_color='FFFFFF', fill_type='solid')   # White

            thin_border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )

            # Write data
            for r_idx, row in enumerate(dataframe_to_rows(df, index=False, header=True), 1):
                for c_idx, value in enumerate(row, 1):
                    # Handle formula text - write as text to prevent Excel errors
                    # If value starts with '=' and is in data row, it's likely formula text that should be displayed, not executed
                    if r_idx > 1 and value and isinstance(value, str) and value.startswith('='):
                        # Write as text by prefixing with apostrophe (Excel will display without the apostrophe)
                        cell = ws.cell(row=r_idx, column=c_idx)
                        cell.value = value  # Set as plain string
                        # Force it to be text format
                        cell.number_format = '@'  # Text format
                    else:
                        cell = ws.cell(row=r_idx, column=c_idx, value=value)

                    cell.border = thin_border

                    # Style header row
                    if r_idx == 1:
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = header_alignment
                    else:
                        # Style data rows
                        cell.font = data_font
                        cell.alignment = data_alignment

                        # Apply currency formatting to currency columns (but not if it's formula text)
                        if (c_idx - 1) in currency_columns and not (value and isinstance(value, str) and value.startswith('=')):
                            cell.number_format = '$#,##0.00'  # Currency format with 2 decimals

                        # Alternating row colors
                        if r_idx % 2 == 0:
                            cell.fill = even_fill
                        else:
                            cell.fill = odd_fill

            # Auto-adjust column widths based on content
            for col_idx, column in enumerate(ws.columns, 1):
                max_length = 0
                column_letter = get_column_letter(col_idx)

                for cell in column:
                    try:
                        if cell.value:
                            # Calculate display length based on value and format
                            if cell.number_format and '$' in cell.number_format:
                                # Currency values need extra space for $ and formatting
                                # Example: $180,000.00 needs more space than raw value 180000
                                if isinstance(cell.value, (int, float)):
                                    # Format as currency to get actual display length
                                    formatted = f"${cell.value:,.2f}"
                                    cell_length = len(formatted)
                                else:
                                    cell_length = len(str(cell.value))
                            else:
                                cell_length = len(str(cell.value))

                            max_length = max(max_length, cell_length)
                    except:
                        pass

                # Add padding (3 characters) and set reasonable bounds
                adjusted_width = min(max(max_length + 3, 10), 50)  # Min 10, Max 50
                ws.column_dimensions[column_letter].width = adjusted_width

            # Set minimum row heights with auto-adjust capability (maintains padding while allowing content expansion)
            ws.row_dimensions[1].height = 32  # Header row - minimum height with padding
            for row_idx in range(2, len(df) + 2):
                ws.row_dimensions[row_idx].height = 26  # Data rows - minimum height, will auto-expand with content

            # Update auto-filter range to match actual data
            num_cols = len(df.columns)
            ws.auto_filter.ref = f'A1:{get_column_letter(num_cols)}{len(df) + 1}'

            # Freeze header row
            ws.freeze_panes = ws['A2']

            wb.save(output_file)
            print(f"  ✅ {csv_file.name} → {output_file.name}")
            self.stats['csv_to_xlsx'] += 1

        except Exception as e:
            error_msg = f"Error converting {csv_file.name}: {str(e)}"
            print(f"  ❌ {error_msg}")
            self.stats['errors'].append(error_msg)

    def _update_word_cover(self, doc, md_file, md_content):
        """Update Word template cover page with document metadata from YAML frontmatter."""
        # Parse YAML metadata from markdown
        metadata, content = self._parse_yaml_frontmatter(md_content)

        # Get current date if not in metadata
        current_date = metadata.get('document_date', datetime.now().strftime('%B %d, %Y'))

        # Extract document title for footer from YAML or filename
        # e.g., "statement-of-work" → "Statement of Work"
        file_stem = md_file.stem
        default_doc_title = file_stem.replace('-', ' ').title()
        document_title = metadata.get('document_title', default_doc_title)

        # Extract project name from first H1 if not in metadata
        project_name = metadata.get('project_name', 'Project Name')
        if not project_name or project_name == 'Project Name':
            lines = content.split('\n')
            for line in lines:
                if line.startswith('# '):
                    project_name = line[2:].strip()
                    break

        # Update cover page table (Table 0) with metadata - simple approach
        if len(doc.tables) > 0:
            cover_table = doc.tables[0]

            # Define the metadata values in order
            table_data = [
                metadata.get('client_name', '[Client Company Name]'),
                current_date,
                metadata.get('vendor_name', '[Consulting Company Name]'),
                f"{metadata.get('vendor_contact_name', '[Full Name]')}\n{metadata.get('vendor_contact_email', 'email@example.com')}",
                metadata.get('opportunity_no', 'OPP-XXXXX'),
                metadata.get('document_version', 'v1.0')
            ]

            # Update only the second column (values) in each row
            # Leave the first column (labels) and all formatting untouched
            for row_idx in range(min(len(cover_table.rows), len(table_data))):
                if len(cover_table.rows[row_idx].cells) >= 2:
                    # Update only the value cell (column 2)
                    cover_table.rows[row_idx].cells[1].text = table_data[row_idx]

        # Also update [Project Name] in paragraphs (for title/subtitle)
        from docx.shared import Pt

        for para in doc.paragraphs:
            if '[Project Name]' in para.text:
                # Replace text while preserving formatting
                for run in para.runs:
                    if '[Project Name]' in run.text:
                        run.text = run.text.replace('[Project Name]', project_name)
                        # Set font size to 20pt and ensure bold for cover page project name
                        run.font.size = Pt(20)
                        run.font.bold = True

            # Update [Vendor/Consultant Name] placeholder
            elif '[Vendor/Consultant Name]' in para.text:
                para.text = para.text.replace('[Vendor/Consultant Name]', 'EO Framework Consulting')

            # Update vendor contact info placeholder
            elif '[Address] • [Phone] • [Email] • [Website]' in para.text:
                para.text = para.text.replace(
                    '[Address] • [Phone] • [Email] • [Website]',
                    '123 Business Street, Suite 100 • (555) 123-4567 • info@eoframework.com • www.eoframework.com'
                )

        # Update footer - replace placeholders with actual values
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        from docx.shared import Pt

        for section in doc.sections:
            footer = section.footer

            for para in footer.paragraphs:
                # Get full paragraph text to check for placeholders
                full_text = para.text

                # Replace PROJECT NAME placeholder (case-insensitive, handles split runs)
                if 'PROJECT NAME' in full_text.upper():
                    # Replace across all runs
                    for run in para.runs:
                        if run.text:
                            # Replace case-insensitively
                            import re
                            run.text = re.sub(r'PROJECT NAME', project_name, run.text, flags=re.IGNORECASE)

                # Replace "PAGE NO" placeholder with page number fields
                if 'PAGE NO' in full_text.upper():
                    # Find the run containing PAGE NO
                    for i, run in enumerate(para.runs):
                        if run.text and 'PAGE NO' in run.text.upper():
                            # Clear the run text
                            run.text = ''

                            # Add "Page "
                            run.text = 'Page '
                            run.font.name = 'Times New Roman'
                            run.font.size = Pt(10)

                            # Add PAGE field (current page number)
                            fldChar1 = OxmlElement('w:fldChar')
                            fldChar1.set(qn('w:fldCharType'), 'begin')
                            run._element.append(fldChar1)

                            instrText = OxmlElement('w:instrText')
                            instrText.set(qn('xml:space'), 'preserve')
                            instrText.text = 'PAGE'
                            run._element.append(instrText)

                            fldChar2 = OxmlElement('w:fldChar')
                            fldChar2.set(qn('w:fldCharType'), 'end')
                            run._element.append(fldChar2)

                            # Add " of " in a new run
                            run_of = para.add_run(' of ')
                            run_of.font.name = 'Times New Roman'
                            run_of.font.size = Pt(10)

                            # Add NUMPAGES field (total pages)
                            run_total = para.add_run()
                            run_total.font.name = 'Times New Roman'
                            run_total.font.size = Pt(10)

                            fldChar3 = OxmlElement('w:fldChar')
                            fldChar3.set(qn('w:fldCharType'), 'begin')
                            run_total._element.append(fldChar3)

                            instrText2 = OxmlElement('w:instrText')
                            instrText2.set(qn('xml:space'), 'preserve')
                            instrText2.text = 'NUMPAGES'
                            run_total._element.append(instrText2)

                            fldChar4 = OxmlElement('w:fldChar')
                            fldChar4.set(qn('w:fldCharType'), 'end')
                            run_total._element.append(fldChar4)

                            break  # Only replace first occurrence

            # Also update footer table if it exists (2-column table: project | page#)
            if footer.tables:
                footer_table = footer.tables[0]

                if len(footer_table.rows) > 0 and len(footer_table.rows[0].cells) >= 2:
                    from docx.enum.text import WD_ALIGN_PARAGRAPH
                    from docx.shared import Inches
                    from docx.enum.table import WD_TABLE_ALIGNMENT
                    from docx.oxml.shared import OxmlElement

                    # Set table to full width - critical for proper layout
                    footer_table.alignment = WD_TABLE_ALIGNMENT.LEFT
                    footer_table.autofit = False

                    # Set table width to 100% of page width
                    tbl = footer_table._element
                    tblPr = tbl.tblPr

                    # Remove existing tblW element if present
                    existing_tblW = tblPr.find(qn('w:tblW'))
                    if existing_tblW is not None:
                        tblPr.remove(existing_tblW)

                    # Add new tblW element for 100% width
                    tblW = OxmlElement('w:tblW')
                    tblW.set(qn('w:w'), '5000')  # 5000 = 100% in Word
                    tblW.set(qn('w:type'), 'pct')
                    tblPr.append(tblW)

                    row = footer_table.rows[0]
                    cell_1 = row.cells[0]
                    cell_2 = row.cells[1]

                    # Set explicit widths for full page coverage (6.5" total for letter size with 1" margins)
                    cell_1.width = Inches(5.0)  # Left cell (project name) - wider to prevent wrapping
                    cell_2.width = Inches(1.5)  # Right cell (page numbers)

                    # Cell 1 (left): Project name
                    if cell_1.paragraphs:
                        para_1 = cell_1.paragraphs[0]
                        para_1.clear()
                        para_1.alignment = WD_ALIGN_PARAGRAPH.LEFT
                        run = para_1.add_run(project_name)
                        run.font.name = 'Times New Roman'
                        run.font.size = Pt(10)
                        run.font.bold = True

                    # Cell 2 (right): Page X of Y
                    # Add paragraph if cell doesn't have one
                    if not cell_2.paragraphs:
                        para_2 = cell_2.add_paragraph()
                    else:
                        para_2 = cell_2.paragraphs[0]

                    para_2.clear()
                    para_2.alignment = WD_ALIGN_PARAGRAPH.RIGHT

                    # Add "Page "
                    run_page = para_2.add_run('Page ')
                    run_page.font.name = 'Times New Roman'
                    run_page.font.size = Pt(10)

                    # Add PAGE field (current page number)
                    run_page_num = para_2.add_run()
                    run_page_num.font.name = 'Times New Roman'
                    run_page_num.font.size = Pt(10)

                    fldChar1 = OxmlElement('w:fldChar')
                    fldChar1.set(qn('w:fldCharType'), 'begin')
                    run_page_num._element.append(fldChar1)

                    instrText = OxmlElement('w:instrText')
                    instrText.set(qn('xml:space'), 'preserve')
                    instrText.text = 'PAGE'
                    run_page_num._element.append(instrText)

                    fldChar2 = OxmlElement('w:fldChar')
                    fldChar2.set(qn('w:fldCharType'), 'separate')
                    run_page_num._element.append(fldChar2)

                    # Add text element for field result
                    t = OxmlElement('w:t')
                    t.text = '1'
                    run_page_num._element.append(t)

                    fldChar3 = OxmlElement('w:fldChar')
                    fldChar3.set(qn('w:fldCharType'), 'end')
                    run_page_num._element.append(fldChar3)

                    # Add " of "
                    run_of = para_2.add_run(' of ')
                    run_of.font.name = 'Times New Roman'
                    run_of.font.size = Pt(10)

                    # Add NUMPAGES field (total pages)
                    run_total = para_2.add_run()
                    run_total.font.name = 'Times New Roman'
                    run_total.font.size = Pt(10)

                    fldChar4 = OxmlElement('w:fldChar')
                    fldChar4.set(qn('w:fldCharType'), 'begin')
                    run_total._element.append(fldChar4)

                    instrText2 = OxmlElement('w:instrText')
                    instrText2.set(qn('xml:space'), 'preserve')
                    instrText2.text = 'NUMPAGES'
                    run_total._element.append(instrText2)

                    fldChar5 = OxmlElement('w:fldChar')
                    fldChar5.set(qn('w:fldCharType'), 'separate')
                    run_total._element.append(fldChar5)

                    # Add text element for field result
                    t2 = OxmlElement('w:t')
                    t2.text = '1'
                    run_total._element.append(t2)

                    fldChar6 = OxmlElement('w:fldChar')
                    fldChar6.set(qn('w:fldCharType'), 'end')
                    run_total._element.append(fldChar6)

    def _add_eo_framework_branding(self, doc):
        """Add EO Framework logo and website at the end of the document."""
        from pathlib import Path
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        try:
            # Add spacing before branding (3 empty paragraphs)
            for _ in range(3):
                doc.add_paragraph()

            # Add horizontal line separator
            separator_para = doc.add_paragraph()
            separator_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = separator_para.add_run('_______________________________________________')
            run.font.size = Pt(10)
            run.font.color.rgb = None  # Use default color

            # Add small spacing
            doc.add_paragraph()

            # Add centered paragraph for logo
            logo_para = doc.add_paragraph()
            logo_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Add EO Framework logo
            logo_path = Path(__file__).parent.parent / 'doc-templates' / 'assets' / 'logos' / 'eo-framework-logo-real.png'
            if logo_path.exists():
                run = logo_para.add_run()
                run.add_picture(str(logo_path), width=Inches(1.5))

            # Add website URL below logo
            url_para = doc.add_paragraph()
            url_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = url_para.add_run('www.eoframework.org')
            run.font.size = Pt(10)
            run.font.name = 'Times New Roman'

        except Exception as e:
            print(f"  ⚠️  Warning: Could not add EO Framework branding: {e}")

    def _add_formatted_text(self, paragraph, text):
        """Add text to paragraph with markdown bold/italic formatting applied.

        Parses **bold** and *italic* markers and applies Word formatting.
        """
        # Pattern to match **bold** or *italic*
        # Match **text** first (bold), then *text* (italic)
        import re

        pos = 0
        while pos < len(text):
            # Find next bold or italic marker
            bold_match = re.search(r'\*\*(.+?)\*\*', text[pos:])
            italic_match = re.search(r'\*(.+?)\*', text[pos:])

            # Determine which comes first
            next_match = None
            is_bold = False

            if bold_match and italic_match:
                # Both found, prioritize bold (longer match) when they start at same position
                if bold_match.start() <= italic_match.start():
                    next_match = bold_match
                    is_bold = True
                else:
                    next_match = italic_match
                    is_bold = False
            elif bold_match:
                next_match = bold_match
                is_bold = True
            elif italic_match:
                next_match = italic_match
                is_bold = False

            if next_match:
                # Add text before the match as normal text
                if next_match.start() > 0:
                    before_text = text[pos:pos + next_match.start()]
                    run = paragraph.add_run(before_text)
                    run.font.name = self.template_fonts['body_font']
                    run.font.size = Pt(self.template_fonts['body_size'])

                # Add the formatted text
                formatted_text = next_match.group(1)
                run = paragraph.add_run(formatted_text)
                run.font.name = self.template_fonts['body_font']
                run.font.size = Pt(self.template_fonts['body_size'])
                if is_bold:
                    run.font.bold = True
                else:
                    run.font.italic = True

                # Move position past this match
                pos = pos + next_match.end()
            else:
                # No more matches, add remaining text
                remaining = text[pos:]
                if remaining:
                    run = paragraph.add_run(remaining)
                    run.font.name = self.template_fonts['body_font']
                    run.font.size = Pt(self.template_fonts['body_size'])
                break

    def _generate_word_toc(self, doc, tokens):
        """Keep the template's TOC as-is. User will update it in Word using built-in TOC functionality."""
        # Do nothing - preserve the template's TOC field so it can be updated in Word
        pass

    def _extract_template_fonts(self, doc):
        """Extract font properties from Word template."""
        font_info = {
            'body_font': None,
            'body_size': None,
            'heading_font': None,
            'heading_size': None,
        }

        try:
            # Access the document's XML to get default font settings
            from docx.oxml.ns import qn

            # Get the document XML element
            doc_element = doc._element

            # Try to find styles.xml through the document part
            styles_part = None
            for rel in doc.part.rels.values():
                if 'styles' in rel.target_ref:
                    styles_part = rel.target_part
                    break

            if styles_part:
                styles_element = styles_part.element

                # Get docDefaults
                doc_defaults = styles_element.find(qn('w:docDefaults'))
                if doc_defaults is not None:
                    rpr_default = doc_defaults.find(qn('w:rPrDefault'))
                    if rpr_default is not None:
                        rpr = rpr_default.find(qn('w:rPr'))
                        if rpr is not None:
                            # Get font name
                            rfonts = rpr.find(qn('w:rFonts'))
                            if rfonts is not None:
                                font_info['body_font'] = rfonts.get(qn('w:ascii'))

                            # Get font size (in half-points)
                            sz = rpr.find(qn('w:sz'))
                            if sz is not None:
                                size_half_pts = sz.get(qn('w:val'))
                                if size_half_pts:
                                    font_info['body_size'] = int(size_half_pts) // 2

                # Try to get Normal style as fallback
                normal_style = None
                for style in styles_element.findall(qn('w:style')):
                    style_id = style.get(qn('w:styleId'))
                    if style_id == 'Normal':
                        normal_style = style
                        break

                if normal_style is not None:
                    rpr = normal_style.find(qn('w:rPr'))
                    if rpr is not None:
                        # Get font from Normal style if not found in defaults
                        if not font_info['body_font']:
                            rfonts = rpr.find(qn('w:rFonts'))
                            if rfonts is not None:
                                font_info['body_font'] = rfonts.get(qn('w:ascii'))

                        # Get size from Normal style if not found in defaults
                        if not font_info['body_size']:
                            sz = rpr.find(qn('w:sz'))
                            if sz is not None:
                                size_half_pts = sz.get(qn('w:val'))
                                if size_half_pts:
                                    font_info['body_size'] = int(size_half_pts) // 2

                # Get Heading 1 style for heading font
                heading1_style = None
                for style in styles_element.findall(qn('w:style')):
                    style_id = style.get(qn('w:styleId'))
                    if style_id == 'Heading1':
                        heading1_style = style
                        break

                if heading1_style is not None:
                    rpr = heading1_style.find(qn('w:rPr'))
                    if rpr is not None:
                        rfonts = rpr.find(qn('w:rFonts'))
                        if rfonts is not None:
                            font_info['heading_font'] = rfonts.get(qn('w:ascii'))

                        sz = rpr.find(qn('w:sz'))
                        if sz is not None:
                            size_half_pts = sz.get(qn('w:val'))
                            if size_half_pts:
                                font_info['heading_size'] = int(size_half_pts) // 2

            # Set defaults if nothing found
            if not font_info['body_font']:
                font_info['body_font'] = 'Times New Roman'  # Word default
            if not font_info['body_size']:
                font_info['body_size'] = 11  # Word default
            if not font_info['heading_font']:
                font_info['heading_font'] = font_info['body_font']
            if not font_info['heading_size']:
                font_info['heading_size'] = 16

        except Exception as e:
            print(f"Warning: Could not extract template fonts: {e}")
            # Use Word defaults
            font_info = {
                'body_font': 'Times New Roman',
                'body_size': 11,
                'heading_font': 'Times New Roman',
                'heading_size': 16,
            }

        return font_info

    def _extract_template_table_format(self, doc):
        """Extract table formatting properties from Word template."""
        table_format = {
            'header_bg_color': 'E8F0F8',  # Light blue from template
            'header_padding_before': 6,
            'header_padding_after': 6,
            'border_color': 'CCCCCC',  # Gray borders
            'border_size': '1',
        }

        try:
            # Find a content table in the template (not the metadata table)
            # Look for tables after the first one
            if len(doc.tables) > 1:
                template_table = doc.tables[1]  # Second table (first content table)

                # Extract header row formatting from first row
                if len(template_table.rows) > 0:
                    header_row = template_table.rows[0]
                    if len(header_row.cells) > 0:
                        header_cell = header_row.cells[0]

                        # Extract background color
                        from docx.oxml.ns import qn
                        tcPr = header_cell._element.get_or_add_tcPr()
                        shd = tcPr.find(qn('w:shd'))
                        if shd is not None:
                            bg_color = shd.get(qn('w:fill'))
                            if bg_color:
                                table_format['header_bg_color'] = bg_color

                        # Extract padding from paragraph
                        if header_cell.paragraphs:
                            para = header_cell.paragraphs[0]
                            space_before = para.paragraph_format.space_before
                            space_after = para.paragraph_format.space_after
                            if space_before:
                                table_format['header_padding_before'] = int(space_before.pt)
                            if space_after:
                                table_format['header_padding_after'] = int(space_after.pt)

                        # Extract border color and size
                        tcBorders = tcPr.find(qn('w:tcBorders'))
                        if tcBorders is not None:
                            top_border = tcBorders.find(qn('w:top'))
                            if top_border is not None:
                                border_color = top_border.get(qn('w:color'))
                                border_size = top_border.get(qn('w:sz'))
                                if border_color:
                                    table_format['border_color'] = border_color
                                if border_size:
                                    table_format['border_size'] = border_size

        except Exception as e:
            print(f"Warning: Could not extract table formatting: {e}")
            # Use defaults already set

        return table_format

    def _extract_template_signature_format(self, doc):
        """Extract signature table formatting from template."""
        sig_format = {
            'rows': 5,
            'cols': 2,
            'header_bg': 'E8F0F8',
            'header_padding': (6, 6),
            'sig_padding': (12, 12),
            'field_padding': (18, 6),
            'border_color': 'CCCCCC',
            'border_size': '1',
        }

        try:
            # Find signature table (last table in template)
            if doc.tables:
                sig_table = doc.tables[-1]

                # Verify it looks like a signature table
                if len(sig_table.rows) == 5 and len(sig_table.columns) == 2:
                    # Extract header formatting (row 0)
                    from docx.oxml.ns import qn
                    header_cell = sig_table.rows[0].cells[0]
                    tcPr = header_cell._element.get_or_add_tcPr()

                    shd = tcPr.find(qn('w:shd'))
                    if shd is not None:
                        bg_color = shd.get(qn('w:fill'))
                        if bg_color:
                            sig_format['header_bg'] = bg_color

                    if header_cell.paragraphs:
                        para = header_cell.paragraphs[0]
                        space_before = para.paragraph_format.space_before
                        space_after = para.paragraph_format.space_after
                        if space_before and space_after:
                            sig_format['header_padding'] = (int(space_before.pt), int(space_after.pt))

                    # Extract signature row formatting (row 1)
                    sig_cell = sig_table.rows[1].cells[0]
                    if sig_cell.paragraphs:
                        para = sig_cell.paragraphs[0]
                        space_before = para.paragraph_format.space_before
                        space_after = para.paragraph_format.space_after
                        if space_before and space_after:
                            sig_format['sig_padding'] = (int(space_before.pt), int(space_after.pt))

                    # Extract field row formatting (row 2 - Name)
                    field_cell = sig_table.rows[2].cells[0]
                    if field_cell.paragraphs:
                        para = field_cell.paragraphs[0]
                        space_before = para.paragraph_format.space_before
                        space_after = para.paragraph_format.space_after
                        if space_before and space_after:
                            sig_format['field_padding'] = (int(space_before.pt), int(space_after.pt))

                    # Extract border properties
                    tcBorders = tcPr.find(qn('w:tcBorders'))
                    if tcBorders is not None:
                        top_border = tcBorders.find(qn('w:top'))
                        if top_border is not None:
                            border_color = top_border.get(qn('w:color'))
                            border_size = top_border.get(qn('w:sz'))
                            if border_color:
                                sig_format['border_color'] = border_color
                            if border_size:
                                sig_format['border_size'] = border_size

        except Exception as e:
            print(f"Warning: Could not extract signature table format: {e}")
            # Use defaults already set

        return sig_format

    def _adjust_toc_list_spacing(self, doc):
        """Reduce spacing around Table of Contents, List of Tables, and List of Figures headings.

        Focuses on reducing the gap between heading and content below it.
        """
        from docx.shared import Pt

        headings_to_adjust = {
            'Table of Contents': {'space_after': 3, 'space_before': None},
            'List of Tables': {'space_after': 3, 'space_before': 6},
            'List of Figures': {'space_after': 3, 'space_before': 6}
        }

        for para in doc.paragraphs:
            if para.style.name == 'Heading 1':
                heading_text = para.text.strip()
                if heading_text in headings_to_adjust:
                    settings = headings_to_adjust[heading_text]

                    # Set space after (gap to content below)
                    para.paragraph_format.space_after = Pt(settings['space_after'])

                    # Set space before if specified (gap from previous section)
                    if settings['space_before'] is not None:
                        para.paragraph_format.space_before = Pt(settings['space_before'])

                    if self.verbose:
                        print(f"  Adjusted spacing for: {heading_text}")

    def convert_md_to_docx(self, md_file, output_dir):
        """Convert Markdown to styled Word document using branded template."""
        try:
            output_file = output_dir / f"{md_file.stem}.docx"

            # Dry-run mode - just show what would be generated
            if self.dry_run:
                print(f"  🔍 Would generate: {md_file.name} → {output_file.name}")
                self.stats['md_to_docx'] += 1
                return

            # Load branded template
            template_path = Path(__file__).parent.parent / 'doc-templates' / 'word' / 'EOFramework-Word-Template-01.docx'

            if not template_path.exists():
                raise FileNotFoundError(f"Word template not found: {template_path}")

            doc = Document(str(template_path))

            # Extract font properties from template
            template_fonts = self._extract_template_fonts(doc)
            print(f"\n{'='*70}")
            print(f"TEMPLATE FONT PROPERTIES EXTRACTED:")
            print(f"{'='*70}")
            print(f"  Body Font:    {template_fonts['body_font']}")
            print(f"  Body Size:    {template_fonts['body_size']}pt")
            print(f"  Heading Font: {template_fonts['heading_font']}")
            print(f"  Heading Size: {template_fonts['heading_size']}pt")
            print(f"{'='*70}\n")

            # Store for use throughout conversion
            self.template_fonts = template_fonts

            # Extract table formatting from template
            self.template_table_format = self._extract_template_table_format(doc)

            print(f"\n{'='*70}")
            print(f"TEMPLATE TABLE PROPERTIES EXTRACTED:")
            print(f"{'='*70}")
            print(f"  Header BG Color:  #{self.template_table_format['header_bg_color']}")
            print(f"  Header Padding:   {self.template_table_format['header_padding_before']}pt / {self.template_table_format['header_padding_after']}pt")
            print(f"  Border Color:     #{self.template_table_format['border_color']}")
            print(f"  Border Size:      {self.template_table_format['border_size']}")
            print(f"{'='*70}\n")

            # Extract signature table format from template
            self.template_signature_format = self._extract_template_signature_format(doc)

            # Read markdown
            with open(md_file, 'r', encoding='utf-8') as f:
                md_content = f.read()

            # Parse YAML frontmatter and get clean content (without YAML)
            metadata, clean_content = self._parse_yaml_frontmatter(md_content)

            # Parse markdown (using clean content without YAML)
            md = MarkdownIt()
            tokens = md.parse(clean_content)

            # Update cover page with original content (including YAML for metadata extraction)
            self._update_word_cover(doc, md_file, md_content)

            # Clear sample content after TOC (page 3 onwards), but keep TOC placeholder
            # Also update List of Figures section with actual figures from markdown
            # Find and remove both paragraphs AND tables after the second page break
            from docx.oxml.ns import qn
            page_break_count = 0
            removal_start_index = None
            list_of_figures_index = None

            for i, para in enumerate(doc.paragraphs):
                # Track List of Figures heading location
                if 'List of Figures' in para.text and para.style.name == 'Heading 1':
                    list_of_figures_index = i

                # Check if paragraph contains page break
                # Look for w:br elements with type="page"
                for run in para.runs:
                    br_list = run._element.findall(qn('w:br'))
                    for br in br_list:
                        br_type = br.get(qn('w:type'))
                        if br_type == 'page':
                            page_break_count += 1
                            break
                    if page_break_count >= 2:
                        break

                if page_break_count >= 2:
                    # Found the paragraph with second page break
                    removal_start_index = i
                    break

            # Remove all content after second page break (including sample figures)
            # BUT preserve section elements
            # Do this BEFORE clearing List of Figures content to avoid index shifting issues
            if removal_start_index is not None:
                # Get the parent element (body)
                body = doc._element.body

                # Get the paragraph element with the second page break
                para_elem = doc.paragraphs[removal_start_index]._element

                # Remove all siblings after this paragraph, EXCEPT sectPr (section properties)
                found_start = False
                elements_to_remove = []
                for elem in body:
                    if found_start:
                        # Don't remove section elements (sectPr)
                        if not elem.tag.endswith('}sectPr'):
                            elements_to_remove.append(elem)
                    elif elem == para_elem:
                        found_start = True

                # Remove elements
                for elem in elements_to_remove:
                    body.remove(elem)

            # Clear sample figure entries from List of Figures section (keep heading only)
            # This happens AFTER removing content post-page-break, so we need to find the heading again
            if list_of_figures_index is not None:
                # Re-find the List of Figures heading (index may have changed after removal)
                list_of_figures_elem = None
                for para in doc.paragraphs:
                    if 'List of Figures' in para.text and para.style.name == 'Heading 1':
                        list_of_figures_elem = para._element
                        break

                if list_of_figures_elem is not None:
                    # Remove all paragraphs after List of Figures heading until we hit another heading or section break
                    body = doc._element.body
                    found_heading = False
                    elements_to_remove = []

                    for elem in body:
                        if elem == list_of_figures_elem:
                            found_heading = True
                            continue

                        if found_heading:
                            # Stop if we hit another heading (Heading 1) or a page break
                            if elem.tag.endswith('}p'):
                                # Check if it's a heading
                                pPr = elem.find(qn('w:pPr'))
                                if pPr is not None:
                                    pStyle = pPr.find(qn('w:pStyle'))
                                    if pStyle is not None:
                                        style_val = pStyle.get(qn('w:val'))
                                        if style_val and style_val.startswith('Heading'):
                                            break

                                # Check for page break
                                has_page_break = False
                                for r in elem.findall(qn('w:r')):
                                    for br in r.findall(qn('w:br')):
                                        if br.get(qn('w:type')) == 'page':
                                            has_page_break = True
                                            break
                                    if has_page_break:
                                        break

                                if has_page_break:
                                    break

                                # This is a regular paragraph under List of Figures, remove it
                                elements_to_remove.append(elem)
                            elif elem.tag.endswith('}tbl'):
                                # Remove tables under List of Figures
                                elements_to_remove.append(elem)

                    for elem in elements_to_remove:
                        body.remove(elem)

            # Generate table of contents (after clearing sample content)
            self._generate_word_toc(doc, tokens)

            # Add content from markdown (using clean content without YAML, pass md_file for image path resolution)
            # Returns lists of tables and figures found
            table_list, figure_list = self._process_md_tokens(doc, tokens, clean_content, md_file)

            # Populate List of Tables section with actual tables (before List of Figures)
            if table_list:
                self._populate_list_of_tables(doc, table_list)

            # Populate List of Figures section with actual figures
            if figure_list:
                self._populate_list_of_figures(doc, figure_list)

            # Adjust spacing for TOC and List sections (do this AFTER all content is generated)
            self._adjust_toc_list_spacing(doc)

            # Add EO Framework branding at the end of the document
            self._add_eo_framework_branding(doc)

            doc.save(output_file)
            print(f"  ✅ {md_file.name} → {output_file.name}")
            self.stats['md_to_docx'] += 1

        except Exception as e:
            error_msg = f"Error converting {md_file.name}: {str(e)}"
            print(f"  ❌ {error_msg}")
            self.stats['errors'].append(error_msg)

    def _parse_md_table(self, md_content, start_pos):
        """Parse a markdown table from content starting at given position.

        Returns: (table_data, end_pos) where table_data is list of lists, end_pos is position after table
        """
        lines = md_content.split('\n')
        table_rows = []
        i = start_pos

        while i < len(lines):
            line = lines[i].strip()

            # Check if this is a table row (starts and ends with |)
            if line.startswith('|') and line.endswith('|'):
                # Skip separator rows (contains only |, -, and spaces)
                if set(line.replace('|', '').replace('-', '').replace(' ', '').replace(':', '')) == set():
                    i += 1
                    continue

                # Parse cells
                cells = [cell.strip() for cell in line.split('|')[1:-1]]
                table_rows.append(cells)
                i += 1
            else:
                # End of table
                break

        return table_rows, i

    def _calculate_text_width(self, text):
        """Calculate weighted character width for proportional fonts.

        Returns a relative width score accounting for character variations.
        """
        if not text:
            return 0

        # Character width weights (relative to average character)
        # Based on typical proportional font (Times New Roman) characteristics
        wide_chars = 'WMQOCDGHAB@%&'  # ~1.4x average
        medium_wide = 'NRUKXYVZ$#0123456789'  # ~1.1x average
        narrow_chars = 'iljft!|'  # ~0.5x average
        space_chars = ' '  # ~0.6x average

        width_score = 0
        for char in str(text):
            if char in wide_chars:
                width_score += 1.4
            elif char in medium_wide:
                width_score += 1.1
            elif char in narrow_chars:
                width_score += 0.5
            elif char in space_chars:
                width_score += 0.6
            else:
                width_score += 1.0  # Average character

        return width_score

    def _calculate_optimal_column_widths(self, table_data, available_width):
        """Calculate optimal column widths based on content.

        Uses hybrid approach: max of header and content (first 5 rows).
        Applies constraints: min 8%, max 35% of total width.
        """
        if not table_data or len(table_data) == 0:
            return None

        num_cols = len(table_data[0])

        # Calculate width scores for each column
        col_scores = [0] * num_cols

        # Sample rows: header + first 5 data rows (or all if fewer)
        sample_rows = table_data[:min(6, len(table_data))]

        for row in sample_rows:
            for col_idx, cell_text in enumerate(row):
                # Remove markdown formatting for width calculation
                clean_text = cell_text.replace('**', '').replace('*', '')
                width_score = self._calculate_text_width(clean_text)
                col_scores[col_idx] = max(col_scores[col_idx], width_score)

        # Convert scores to percentages with constraints
        total_score = sum(col_scores)
        if total_score == 0:
            return None

        # Calculate initial percentages
        col_percentages = [(score / total_score) for score in col_scores]

        # Apply constraints: min 8%, max 35%
        min_pct = 0.08
        max_pct = 0.35

        adjusted = True
        iterations = 0
        while adjusted and iterations < 10:  # Prevent infinite loops
            adjusted = False
            iterations += 1

            # Find columns that violate constraints
            constrained_pct = 0
            constrained_count = 0

            for i in range(num_cols):
                if col_percentages[i] < min_pct:
                    col_percentages[i] = min_pct
                    constrained_pct += min_pct
                    constrained_count += 1
                    adjusted = True
                elif col_percentages[i] > max_pct:
                    col_percentages[i] = max_pct
                    constrained_pct += max_pct
                    constrained_count += 1
                    adjusted = True

            # Redistribute excess/deficit among unconstrained columns
            if adjusted and constrained_count < num_cols:
                remaining_pct = 1.0 - constrained_pct
                unconstrained_total = sum(col_percentages[i] for i in range(num_cols)
                                         if col_percentages[i] > min_pct and col_percentages[i] < max_pct)

                if unconstrained_total > 0:
                    for i in range(num_cols):
                        if col_percentages[i] > min_pct and col_percentages[i] < max_pct:
                            col_percentages[i] = (col_percentages[i] / unconstrained_total) * remaining_pct

        # Normalize to ensure sum = 1.0
        total_pct = sum(col_percentages)
        col_percentages = [p / total_pct for p in col_percentages]

        # Convert percentages to actual widths
        col_widths = [available_width * pct for pct in col_percentages]

        return col_widths

    def _generate_table_caption(self, table_number, section_title, table_data):
        """Generate table caption based on section context and table content.

        Args:
            table_number: Sequential table number
            section_title: Current section heading (lowercase)
            table_data: Table data to analyze

        Returns:
            str: Generated caption (e.g., "Table 1: Project Deliverables")
        """
        # Special case handling
        if not table_data or len(table_data) == 0:
            return f"Table {table_number}"

        # Check table headers for specific patterns
        first_row = ' '.join(table_data[0]).lower() if table_data else ''

        # RACI table
        if 'vendor' in first_row and 'client' in first_row and 'task' in first_row:
            return f"Table {table_number}: Roles and Responsibilities (RACI Matrix)"

        # Deliverables table
        if 'deliverable' in first_row and 'acceptance' in first_row:
            return f"Table {table_number}: Project Deliverables and Acceptance Criteria"

        # Milestones table
        if 'milestone' in first_row and ('target date' in first_row or 'description' in first_row):
            return f"Table {table_number}: Project Milestones and Timeline"

        # Tooling table
        if 'primary tools' in first_row or 'alternative' in first_row:
            return f"Table {table_number}: Implementation Tools and Technologies"

        # Generic caption based on section title
        # Clean up section title for use in caption
        clean_section = section_title.replace('&', 'and').strip()

        # Remove common prefixes
        for prefix in ['overview', 'description', 'details']:
            if clean_section.startswith(prefix):
                clean_section = clean_section[len(prefix):].strip()

        # Capitalize words
        caption_text = ' '.join(word.capitalize() for word in clean_section.split())

        # If section title is too generic or empty, try to infer from headers
        if not caption_text or len(caption_text) < 3:
            # Use first header as basis
            if table_data and len(table_data[0]) > 0:
                caption_text = table_data[0][0].replace('**', '').strip()

        return f"Table {table_number}: {caption_text}"

    def _add_word_table(self, doc, table_data, caption=None):
        """Add a formatted table to Word document with optional caption.

        Args:
            doc: Word document object
            table_data: List of lists containing table data
            caption: Optional caption text (e.g., "Table 1: Deliverables")

        Returns:
            table: The created table object
        """
        if not table_data or len(table_data) < 2:
            return None

        try:
            # Create table
            num_rows = len(table_data)
            num_cols = len(table_data[0])
            table = doc.add_table(rows=num_rows, cols=num_cols)

            # Try to apply style, but don't fail if it doesn't exist
            try:
                table.style = 'Table Grid'
            except:
                pass  # Continue without style if not available

            # Set table to use full page width
            table.autofit = False
            table.allow_autofit = False

            # Calculate available width (page width - margins)
            # Standard page: 8.5" width, 1" margins on each side = 6.5" available
            from docx.shared import Inches
            available_width = Inches(6.5)

            # Set table width to full available width
            table.width = available_width

            # Add borders to table
            from docx.oxml import OxmlElement
            from docx.oxml.ns import qn

            def set_cell_border(cell, **kwargs):
                """Set cell borders using template formatting"""
                tc = cell._element
                tcPr = tc.get_or_add_tcPr()

                # Create border elements using template border color and size
                tcBorders = OxmlElement('w:tcBorders')
                border_color = self.template_table_format['border_color']
                border_size = self.template_table_format['border_size']

                for edge in ('top', 'left', 'bottom', 'right'):
                    edge_elem = OxmlElement(f'w:{edge}')
                    edge_elem.set(qn('w:val'), 'single')
                    edge_elem.set(qn('w:sz'), border_size)
                    edge_elem.set(qn('w:space'), '0')
                    edge_elem.set(qn('w:color'), border_color)
                    tcBorders.append(edge_elem)

                tcPr.append(tcBorders)

                # Add cell margins for padding (top/bottom/left/right)
                tcMar = OxmlElement('w:tcMar')
                for margin_side in ['top', 'bottom', 'left', 'right']:
                    margin_elem = OxmlElement(f'w:{margin_side}')
                    margin_elem.set(qn('w:w'), '80')  # 80 twips = approx 0.06 inches (moderate padding)
                    margin_elem.set(qn('w:type'), 'dxa')
                    tcMar.append(margin_elem)
                tcPr.append(tcMar)

            # Populate and format table
            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_data in enumerate(row_data):
                    cell = table.rows[row_idx].cells[col_idx]

                    # Clear existing content and add formatted text
                    if cell.paragraphs:
                        para = cell.paragraphs[0]
                        para.clear()
                        self._add_formatted_text(para, cell_data)
                    else:
                        cell.text = cell_data

                    # Add borders to cell
                    set_cell_border(cell)

                    # Format header row using template formatting
                    if row_idx == 0:
                        # Set vertical alignment to center
                        from docx.enum.table import WD_ALIGN_VERTICAL
                        cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER

                        # Add padding to header cell
                        if cell.paragraphs:
                            para = cell.paragraphs[0]
                            para.paragraph_format.space_before = Pt(self.template_table_format['header_padding_before'])
                            para.paragraph_format.space_after = Pt(self.template_table_format['header_padding_after'])

                            if para.runs:
                                para.runs[0].font.bold = True
                                # Don't set white text color for light blue headers

                        # Set cell background color using template color (light blue #E8F0F8)
                        from docx.oxml import parse_xml
                        header_bg = self.template_table_format['header_bg_color']
                        shading_elm = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(
                            'xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"',
                            header_bg
                        ))
                        cell._element.get_or_add_tcPr().append(shading_elm)
                    else:
                        # Add equal top and bottom padding to non-header cells for symmetry
                        if cell.paragraphs:
                            para = cell.paragraphs[0]
                            para.paragraph_format.space_before = Pt(4)  # Equal padding
                            para.paragraph_format.space_after = Pt(4)   # Equal padding

                    # Set font for all cells using template font
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = self.template_fonts['body_font']
                            run.font.size = Pt(self.template_fonts['body_size'])

            # Determine column widths based on table type
            col_widths = None

            # Special case 1: RACI table (7 columns with specific headers)
            is_raci_table = False
            if num_cols == 7 and len(table_data) > 0:
                first_row_text = ' '.join(table_data[0]).lower()
                if 'task' in first_row_text and 'vendor' in first_row_text and 'client' in first_row_text:
                    is_raci_table = True

            # Special case 2: Cover/metadata table (2 columns, first column ~33%)
            is_cover_table = False
            if num_cols == 2 and len(table_data) > 3:
                first_col_text = ' '.join([row[0] for row in table_data[:3]]).lower()
                if 'client name' in first_col_text or 'document' in first_col_text:
                    is_cover_table = True

            if is_raci_table:
                # RACI table: Custom widths prioritizing Task/Role column
                base_width = available_width / num_cols
                reduced_width = base_width * 0.85
                extra_width = (base_width - reduced_width) * 3
                col_widths = [
                    base_width + extra_width,  # Task/Role
                    reduced_width,              # Vendor PM
                    reduced_width,              # Vendor Arch
                    reduced_width,              # Vendor DevOps
                    base_width,                 # Client IT
                    base_width,                 # Client Sec
                    base_width                  # SME
                ]
            elif is_cover_table:
                # Cover table: 33% label, 67% content
                col_widths = [available_width * 0.33, available_width * 0.67]
            elif num_cols == 1:
                # Single column: use full width
                col_widths = [available_width]
            else:
                # All other tables: use content-based optimization
                col_widths = self._calculate_optimal_column_widths(table_data, available_width)

            # Fallback: equal distribution if optimization fails
            if not col_widths:
                col_widths = [available_width / num_cols] * num_cols

            # Apply column widths
            for row in table.rows:
                for col_idx, cell in enumerate(row.cells):
                    if col_idx < len(col_widths):
                        cell.width = col_widths[col_idx]

            # Add caption if provided (using proper Word caption fields)
            if caption:
                self._add_table_caption(doc, caption)

            return table

        except Exception as e:
            # Log error but don't stop document generation
            print(f"  Warning: Table formatting error: {e}")
            return None

    def _add_table_caption(self, doc, caption_text):
        """Add a proper Word caption with SEQ field for automatic numbering.

        Args:
            doc: Word document object
            caption_text: Full caption text including number (e.g., "Table 1: Description")

        Creates a caption paragraph with:
        - "Table " literal text
        - SEQ field for auto-numbering
        - ": " separator
        - Caption description text
        """
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn

        # Extract table number and description from caption text
        # caption_text format: "Table 1: Description"
        import re
        match = re.match(r'Table\s+(\d+):\s*(.*)', caption_text)
        if match:
            table_number = match.group(1)
            description = match.group(2)
        else:
            table_number = '1'
            description = caption_text

        # Create caption paragraph
        caption_para = doc.add_paragraph()
        caption_para.style = 'Caption'
        caption_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        caption_para.paragraph_format.space_before = Pt(3)
        caption_para.paragraph_format.space_after = Pt(12)

        # Add "Table " text
        run1 = caption_para.add_run('Table ')
        run1.font.name = self.template_fonts['body_font']
        run1.font.size = Pt(12)

        # Add SEQ field for auto-numbering
        # Structure: fldChar(begin) -> instrText -> fldChar(separate) -> run(result) -> fldChar(end)

        # Field begin
        fld_begin = OxmlElement('w:r')
        fld_char_begin = OxmlElement('w:fldChar')
        fld_char_begin.set(qn('w:fldCharType'), 'begin')
        fld_begin.append(fld_char_begin)
        caption_para._element.append(fld_begin)

        # Field instruction text: " SEQ Table \* ARABIC "
        fld_instr = OxmlElement('w:r')
        instr_text = OxmlElement('w:instrText')
        instr_text.set(qn('xml:space'), 'preserve')
        instr_text.text = ' SEQ Table \\* ARABIC '
        fld_instr.append(instr_text)
        caption_para._element.append(fld_instr)

        # Field separator
        fld_sep = OxmlElement('w:r')
        fld_char_sep = OxmlElement('w:fldChar')
        fld_char_sep.set(qn('w:fldCharType'), 'separate')
        fld_sep.append(fld_char_sep)
        caption_para._element.append(fld_sep)

        # Field result (actual table number)
        result_run = caption_para.add_run(table_number)
        result_run.font.name = self.template_fonts['body_font']
        result_run.font.size = Pt(12)

        # Field end
        fld_end = OxmlElement('w:r')
        fld_char_end = OxmlElement('w:fldChar')
        fld_char_end.set(qn('w:fldCharType'), 'end')
        fld_end.append(fld_char_end)
        caption_para._element.append(fld_end)

        # Add ": " separator and description
        run_desc = caption_para.add_run(f': {description}')
        run_desc.font.name = self.template_fonts['body_font']
        run_desc.font.size = Pt(12)

    def _add_figure_caption(self, doc, caption_text):
        """Add a proper Word figure caption with SEQ field for automatic numbering.

        Args:
            doc: Word document object
            caption_text: Full caption text including number (e.g., "Figure 1: Description")

        Creates a caption paragraph with:
        - "Figure " literal text
        - SEQ field for auto-numbering
        - ": " separator
        - Caption description text
        """
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn

        # Extract figure number and description from caption text
        # caption_text format: "Figure 1: Description"
        import re
        match = re.match(r'Figure\s+(\d+):\s*(.*)', caption_text)
        if match:
            figure_number = match.group(1)
            description = match.group(2)
        else:
            figure_number = '1'
            description = caption_text

        # Create caption paragraph
        caption_para = doc.add_paragraph()
        caption_para.style = 'Caption'
        caption_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        caption_para.paragraph_format.space_before = Pt(3)
        caption_para.paragraph_format.space_after = Pt(12)

        # Add "Figure " text
        run1 = caption_para.add_run('Figure ')
        run1.font.name = self.template_fonts['body_font']
        run1.font.size = Pt(12)

        # Add SEQ field for auto-numbering
        # Structure: fldChar(begin) -> instrText -> fldChar(separate) -> run(result) -> fldChar(end)

        # Field begin
        fld_begin = OxmlElement('w:r')
        fld_char_begin = OxmlElement('w:fldChar')
        fld_char_begin.set(qn('w:fldCharType'), 'begin')
        fld_begin.append(fld_char_begin)
        caption_para._element.append(fld_begin)

        # Field instruction text: " SEQ Figure \* ARABIC "
        fld_instr = OxmlElement('w:r')
        instr_text = OxmlElement('w:instrText')
        instr_text.set(qn('xml:space'), 'preserve')
        instr_text.text = ' SEQ Figure \\* ARABIC '
        fld_instr.append(instr_text)
        caption_para._element.append(fld_instr)

        # Field separator
        fld_sep = OxmlElement('w:r')
        fld_char_sep = OxmlElement('w:fldChar')
        fld_char_sep.set(qn('w:fldCharType'), 'separate')
        fld_sep.append(fld_char_sep)
        caption_para._element.append(fld_sep)

        # Field result (actual figure number)
        result_run = caption_para.add_run(figure_number)
        result_run.font.name = self.template_fonts['body_font']
        result_run.font.size = Pt(12)

        # Field end
        fld_end = OxmlElement('w:r')
        fld_char_end = OxmlElement('w:fldChar')
        fld_char_end.set(qn('w:fldCharType'), 'end')
        fld_end.append(fld_char_end)
        caption_para._element.append(fld_end)

        # Add ": " separator and description
        run_desc = caption_para.add_run(f': {description}')
        run_desc.font.name = self.template_fonts['body_font']
        run_desc.font.size = Pt(12)

    def _add_signoff_table(self, doc):
        """Add sign-off table matching template format (5 rows x 2 columns)."""
        from docx.oxml import OxmlElement, parse_xml
        from docx.oxml.ns import qn
        from docx.enum.table import WD_ALIGN_VERTICAL

        # Get signature format from template
        sig_fmt = self.template_signature_format

        # Create table with 5 rows, 2 columns
        table = doc.add_table(rows=sig_fmt['rows'], cols=sig_fmt['cols'])
        table.style = 'Normal Table'

        # Set table to use full page width
        table.autofit = False
        table.allow_autofit = False
        available_width = Inches(6.5)
        table.width = available_width

        # Helper function for borders
        def set_cell_border(cell):
            """Set cell borders using template formatting"""
            tc = cell._element
            tcPr = tc.get_or_add_tcPr()
            tcBorders = OxmlElement('w:tcBorders')
            border_color = sig_fmt['border_color']
            border_size = sig_fmt['border_size']

            for edge in ('top', 'left', 'bottom', 'right'):
                edge_elem = OxmlElement(f'w:{edge}')
                edge_elem.set(qn('w:val'), 'single')
                edge_elem.set(qn('w:sz'), border_size)
                edge_elem.set(qn('w:space'), '0')
                edge_elem.set(qn('w:color'), border_color)
                tcBorders.append(edge_elem)
            tcPr.append(tcBorders)

            # Add cell margins for padding (top/bottom/left/right)
            tcMar = OxmlElement('w:tcMar')
            for margin_side in ['top', 'bottom', 'left', 'right']:
                margin_elem = OxmlElement(f'w:{margin_side}')
                margin_elem.set(qn('w:w'), '80')  # 80 twips = approx 0.06 inches (moderate padding)
                margin_elem.set(qn('w:type'), 'dxa')
                tcMar.append(margin_elem)
            tcPr.append(tcMar)

        # Row labels
        row_labels = ['Party', 'Signature', 'Name', 'Title', 'Date']

        for row_idx, label in enumerate(row_labels):
            row = table.rows[row_idx]

            # Left column - CLIENT
            left_cell = row.cells[0]
            left_para = left_cell.paragraphs[0]

            if row_idx == 0:  # Header row
                left_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = left_para.add_run('CLIENT')
                run.font.bold = True
                run.font.name = self.template_fonts['body_font']
                run.font.size = Pt(self.template_fonts['body_size'])

                # Set light blue background
                shading_elm = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(
                    'xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"',
                    sig_fmt['header_bg']
                ))
                left_cell._element.get_or_add_tcPr().append(shading_elm)

                left_para.paragraph_format.space_before = Pt(sig_fmt['header_padding'][0])
                left_para.paragraph_format.space_after = Pt(sig_fmt['header_padding'][1])
            else:
                # All signature fields - left aligned, bottom justified, same padding
                left_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
                left_para.add_run(f'{label}: ____________________')
                left_para.paragraph_format.space_before = Pt(sig_fmt['field_padding'][0])
                left_para.paragraph_format.space_after = Pt(sig_fmt['field_padding'][1])
                left_cell.vertical_alignment = WD_ALIGN_VERTICAL.BOTTOM

                left_para.runs[0].font.name = self.template_fonts['body_font']
                left_para.runs[0].font.size = Pt(self.template_fonts['body_size'])

            set_cell_border(left_cell)

            # Right column - CONSULTING COMPANY
            right_cell = row.cells[1]
            right_para = right_cell.paragraphs[0]

            if row_idx == 0:  # Header row
                right_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = right_para.add_run('CONSULTING COMPANY')
                run.font.bold = True
                run.font.name = self.template_fonts['body_font']
                run.font.size = Pt(self.template_fonts['body_size'])

                # Set light blue background
                shading_elm = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(
                    'xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"',
                    sig_fmt['header_bg']
                ))
                right_cell._element.get_or_add_tcPr().append(shading_elm)

                right_para.paragraph_format.space_before = Pt(sig_fmt['header_padding'][0])
                right_para.paragraph_format.space_after = Pt(sig_fmt['header_padding'][1])
            else:
                # All signature fields - left aligned, bottom justified, same padding
                right_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
                right_para.add_run(f'{label}: ____________________')
                right_para.paragraph_format.space_before = Pt(sig_fmt['field_padding'][0])
                right_para.paragraph_format.space_after = Pt(sig_fmt['field_padding'][1])
                right_cell.vertical_alignment = WD_ALIGN_VERTICAL.BOTTOM

                right_para.runs[0].font.name = self.template_fonts['body_font']
                right_para.runs[0].font.size = Pt(self.template_fonts['body_size'])

            set_cell_border(right_cell)

        # Set column widths to equal (half of available width each)
        col_width = available_width / 2
        for row in table.rows:
            for cell in row.cells:
                cell.width = col_width

    def _process_md_tokens(self, doc, tokens, md_content, md_file=None):
        """Process markdown tokens and add to Word document with section numbering.

        Returns:
            tuple: (table_list, figure_list) - Lists of captions for List of Tables and List of Figures
        """
        # Initialize section counters for numbering
        section_counters = [0, 0, 0, 0, 0, 0]  # Support up to H6
        first_h1_skipped = False
        current_section_title = ''
        figure_list = []  # Track figures for List of Figures
        table_list = []   # Track tables for List of Tables
        table_counter = 0  # Track table numbers

        # Metadata field patterns to skip (these are now on cover page)
        metadata_patterns = [
            'Document Version:',
            'Date:',
            'Prepared by:',
            'Client:',
            'Project:',
            'SOW Number:',
            'Project Name:',
            'Requesting Organization:',
            'Business Sponsor:',
            'Prepared By:',
            'Decision Required By:'
        ]

        # Track if we've seen tables in current section
        lines = md_content.split('\n')
        line_index = 0

        i = 0
        while i < len(tokens):
            token = tokens[i]

            if token.type == 'heading_open':
                level = int(token.tag[1])  # h1 -> 1, h2 -> 2, etc.
                i += 1
                if i < len(tokens) and tokens[i].type == 'inline':
                    text = tokens[i].content

                    # Skip the first H1 (it's on the cover page)
                    if level == 1 and not first_h1_skipped:
                        first_h1_skipped = True
                        i += 1
                        continue

                    # Update section counters
                    section_counters[level - 1] += 1
                    # Reset lower level counters
                    for j in range(level, 6):
                        section_counters[j] = 0

                    # Generate section number (e.g., 1.2.3)
                    section_parts = [str(section_counters[k]) for k in range(level) if section_counters[k] > 0]
                    section_num = '.'.join(section_parts)

                    # Add heading with section number (single space)
                    numbered_text = f"{section_num} {text}"
                    heading = doc.add_heading(numbered_text, level=level)
                    # Ensure headings are left-aligned with no indent
                    heading.paragraph_format.left_indent = Inches(0)
                    heading.paragraph_format.first_line_indent = Inches(0)
                    # Add spacing before headings (more space for H2, less for H3+)
                    if level == 2:
                        heading.paragraph_format.space_before = Pt(12)
                        heading.paragraph_format.space_after = Pt(6)
                    else:
                        heading.paragraph_format.space_before = Pt(8)
                        heading.paragraph_format.space_after = Pt(6)

                    # Store current section title for table detection
                    current_section_title = text.lower()

                    # Check for content after this heading (tables, bullets, etc.)
                    heading_marker = '#' * level + ' ' + text
                    try:
                        md_pos = md_content.find(heading_marker)
                        if md_pos != -1:
                            # Find first line after heading
                            start_line = md_content[:md_pos].count('\n')

                            # Look for special content after this heading
                            for idx in range(start_line + 1, min(len(lines), start_line + 20)):
                                line = lines[idx].strip()

                                # Skip empty lines
                                if not line:
                                    continue

                                # Check for markdown table (ANY section can have a table)
                                if line.startswith('|'):
                                    # Found table, parse it
                                    table_data, end_idx = self._parse_md_table(md_content, idx)
                                    if table_data:
                                        # Increment table counter
                                        table_counter += 1

                                        # Check for explicit caption after table
                                        explicit_caption = None
                                        for caption_idx in range(end_idx, min(len(lines), end_idx + 5)):
                                            caption_line = lines[caption_idx].strip()
                                            if caption_line.startswith('Table:'):
                                                explicit_caption = caption_line[6:].strip()
                                                explicit_caption = f"Table {table_counter}: {explicit_caption}"
                                                break
                                            elif caption_line and not caption_line.startswith('---'):
                                                # Non-empty, non-separator line found - stop looking
                                                break

                                        # Generate caption if not explicitly provided
                                        if not explicit_caption:
                                            caption = self._generate_table_caption(table_counter, current_section_title, table_data)
                                        else:
                                            caption = explicit_caption

                                        # Add table with caption
                                        self._add_word_table(doc, table_data, caption=caption)

                                        # Track table for List of Tables
                                        table_list.append(caption)
                                    break

                                # Check for Milestones bullets (convert to table)
                                elif 'milestone' in current_section_title and line.startswith('-'):
                                    # Collect all milestone bullets
                                    milestone_rows = [['Milestone', 'Target Date']]
                                    j = idx
                                    while j < len(lines) and lines[j].strip().startswith('-'):
                                        bullet = lines[j].strip()[1:].strip()
                                        # Split on — or : or - to separate name and date
                                        if '—' in bullet:
                                            parts = bullet.split('—', 1)
                                        elif ':' in bullet:
                                            parts = bullet.split(':', 1)
                                        else:
                                            parts = [bullet, '[Date]']
                                        milestone_rows.append([parts[0].strip(), parts[1].strip() if len(parts) > 1 else '[Date]'])
                                        j += 1
                                    if len(milestone_rows) > 1:
                                        # Increment table counter and generate caption
                                        table_counter += 1
                                        caption = self._generate_table_caption(table_counter, current_section_title, milestone_rows)

                                        # Add table with caption
                                        self._add_word_table(doc, milestone_rows, caption=caption)

                                        # Track table for List of Tables
                                        table_list.append(caption)
                                    break

                                # Check for Sign-off section (special formatting)
                                elif 'sign' in current_section_title:
                                    # Look ahead for signature indicators
                                    has_signature = False
                                    for check_line in lines[idx:idx+10]:
                                        if '**Client' in check_line or '**Service Provider' in check_line or 'Signature:' in check_line or 'signing below' in check_line:
                                            has_signature = True
                                            break

                                    if has_signature:
                                        # Add introductory paragraph before signature table
                                        intro_para = doc.add_paragraph('By signing below, both parties agree to the scope, approach, and terms outlined in this Statement of Work.')
                                        intro_para.runs[0].font.name = self.template_fonts['body_font']
                                        intro_para.runs[0].font.size = Pt(self.template_fonts['body_size'])
                                        intro_para.paragraph_format.space_after = Pt(12)  # Single spacing before table

                                        # Create sign-off table matching template format
                                        self._add_signoff_table(doc)

                                        # Add spacing after signature table for visual separation
                                        spacing_para = doc.add_paragraph()
                                        spacing_para.paragraph_format.space_after = Pt(12)
                                        break

                                elif line and not line.startswith('#') and not line.startswith('---'):
                                    # Non-empty line that's not a heading means no special handling
                                    break
                    except Exception as e:
                        # Silently continue if table/special content detection fails
                        pass

                i += 1  # Skip heading_close

            elif token.type == 'paragraph_open':
                i += 1
                if i < len(tokens) and tokens[i].type == 'inline':
                    text = tokens[i].content

                    # Check for markdown image syntax: ![alt text](image_path)
                    import re
                    image_match = re.search(r'!\[([^\]]*)\]\(([^\)]+)\)', text)
                    if image_match:
                        # Extract image details
                        alt_text = image_match.group(1)
                        image_path = image_match.group(2)

                        # Resolve image path relative to markdown file directory
                        from pathlib import Path
                        if md_file:
                            # Get the directory containing the markdown file
                            md_file_dir = md_file.parent
                            # Resolve the relative image path from the markdown file location
                            full_image_path = (md_file_dir / image_path).resolve()
                        else:
                            # Fallback: try as absolute path or relative to current directory
                            full_image_path = Path(image_path).resolve()

                        # Try to add the image - using table-based border approach (safer)
                        try:
                            if full_image_path.exists():
                                # Create a 1x1 table to hold the image with border
                                from docx.enum.table import WD_TABLE_ALIGNMENT
                                from docx.oxml import OxmlElement
                                from docx.oxml.ns import qn

                                # Create table with 1 cell
                                table = doc.add_table(rows=1, cols=1)
                                table.alignment = WD_TABLE_ALIGNMENT.CENTER

                                # Get the cell
                                cell = table.rows[0].cells[0]

                                # Set cell padding to create space around image
                                cell_pr = cell._element.get_or_add_tcPr()
                                tc_mar = OxmlElement('w:tcMar')
                                for margin_name in ['top', 'left', 'bottom', 'right']:
                                    node = OxmlElement(f'w:{margin_name}')
                                    node.set(qn('w:w'), '100')  # 100 twips = about 7pt padding
                                    node.set(qn('w:type'), 'dxa')
                                    tc_mar.append(node)
                                cell_pr.append(tc_mar)

                                # Add the picture to the cell
                                cell_para = cell.paragraphs[0]
                                cell_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                run = cell_para.add_run()
                                picture = run.add_picture(str(full_image_path), width=Inches(6.0))

                                # Apply border to table (1pt black border)
                                tbl_pr = table._element.tblPr
                                tbl_borders = OxmlElement('w:tblBorders')
                                for border_name in ['top', 'left', 'bottom', 'right']:
                                    border = OxmlElement(f'w:{border_name}')
                                    border.set(qn('w:val'), 'single')
                                    border.set(qn('w:sz'), '8')  # 8 eighths of a point = 1pt
                                    border.set(qn('w:space'), '0')
                                    border.set(qn('w:color'), '000000')
                                    tbl_borders.append(border)
                                tbl_pr.append(tbl_borders)

                                # Add proper caption with SEQ field for figures
                                if alt_text:
                                    # Determine figure number (count existing figure captions + 1)
                                    figure_number = len(figure_list) + 1

                                    # Create full caption text: "Figure X: Description"
                                    if alt_text.startswith('Figure'):
                                        full_caption = alt_text
                                    else:
                                        full_caption = f"Figure {figure_number}: {alt_text}"

                                    # Add caption using proper Word field
                                    self._add_figure_caption(doc, full_caption)

                                    # Add to figure list
                                    figure_list.append(full_caption)
                            else:
                                print(f"  Warning: Image file not found: {full_image_path}")
                        except Exception as e:
                            print(f"  Warning: Could not insert image {image_path}: {e}")
                            import traceback
                            traceback.print_exc()

                        i += 1
                        continue

                    # Skip table rows (they're already rendered as tables)
                    if text.startswith('|') and text.endswith('|'):
                        i += 1
                        continue

                    # Skip table separators
                    if text.startswith('|') and set(text.replace('|', '').replace('-', '').replace(' ', '').replace(':', '')) == set():
                        i += 1
                        continue

                    # Skip verbose figure descriptions (e.g., "**Figure 1: ...** - High-level overview...")
                    # These are redundant since we add simple captions programmatically
                    import re
                    if re.match(r'\*\*Figure \d+:.*\*\*\s*-\s*', text):
                        i += 1
                        continue

                    # Check if this paragraph contains metadata (skip if it does)
                    is_metadata = False
                    for pattern in metadata_patterns:
                        if pattern in text:
                            is_metadata = True
                            break

                    # Skip Sign-off section content (already rendered as table)
                    # This includes the "By signing below" paragraph since we add it programmatically
                    if 'sign' in current_section_title:
                        if 'By signing' in text or '**' in text or 'Name:' in text or 'Title:' in text or 'Signature:' in text or 'Date:' in text or '____' in text:
                            i += 1
                            continue

                    if not is_metadata:
                        # Add text with markdown formatting applied
                        if text.strip():
                            para = doc.add_paragraph()
                            self._add_formatted_text(para, text)
                            # Add spacing after paragraphs
                            para.paragraph_format.space_after = Pt(8)
                i += 1  # Skip paragraph_close

            elif token.type == 'bullet_list_open':
                i += 1
                # Skip bullets in Milestones section (they're converted to tables)
                if 'milestone' in current_section_title:
                    while i < len(tokens) and tokens[i].type != 'bullet_list_close':
                        i += 1
                    i += 1  # Skip bullet_list_close
                    continue

                while i < len(tokens) and tokens[i].type != 'bullet_list_close':
                    if tokens[i].type == 'list_item_open':
                        i += 1
                        if i < len(tokens) and tokens[i].type == 'paragraph_open':
                            i += 1
                            if i < len(tokens) and tokens[i].type == 'inline':
                                text = tokens[i].content
                                # Add bullet paragraph with manual formatting (1 tab = 0.5 inches)
                                para = doc.add_paragraph()
                                para.paragraph_format.left_indent = Inches(0.5)
                                para.paragraph_format.first_line_indent = Inches(-0.25)
                                para.paragraph_format.space_after = Pt(4)
                                # Add bullet symbol
                                run = para.add_run('• ')
                                run.font.name = self.template_fonts['body_font']
                                run.font.size = Pt(self.template_fonts['body_size'])
                                # Add formatted text
                                self._add_formatted_text(para, text)
                    i += 1
            else:
                i += 1

        return table_list, figure_list

    def _populate_list_of_tables(self, doc, table_list):
        r"""Populate the List of Tables section with a TOC field that references Table captions.

        Creates a proper Word Table of Contents field that automatically updates
        when user right-clicks and selects "Update Field".

        Field format: { TOC \h \z \c "Table" }
        - \h: Use hyperlinks
        - \z: Hide page numbers in Web Layout view
        - \c "Table": Build table from Caption label "Table"
        """
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn

        # Find the "List of Tables" heading
        list_of_tables_idx = None
        for idx, para in enumerate(doc.paragraphs):
            if para.text.strip() == 'List of Tables' and para.style.name == 'Heading 1':
                list_of_tables_idx = idx
                break

        if list_of_tables_idx is None:
            print("  Warning: 'List of Tables' heading not found in template")
            return

        # Remove template sample content between "List of Tables" heading and next heading
        # This removes entries like "Table 1: Table 1 contents    3"
        paragraphs_to_remove = []
        for idx in range(list_of_tables_idx + 1, len(doc.paragraphs)):
            para = doc.paragraphs[idx]
            # Stop at next heading
            if para.style.name.startswith('Heading'):
                break
            # Mark paragraph for removal if it's a table of figures style or contains template content
            if para.style.name == 'table of figures' or 'Table 1 contents' in para.text:
                paragraphs_to_remove.append(para)

        # Remove marked paragraphs
        for para in paragraphs_to_remove:
            p_element = para._element
            p_element.getparent().remove(p_element)

        # Insert TOC field paragraph after the List of Tables heading
        heading_element = doc.paragraphs[list_of_tables_idx]._element

        # Create a new paragraph for the TOC field
        toc_para = doc.add_paragraph()
        toc_para.style = 'Normal'

        # Move paragraph right after heading
        heading_element.addnext(toc_para._element)

        # Build TOC field: { TOC \h \z \c "Table" }
        # Field begin
        fld_begin = OxmlElement('w:r')
        fld_char_begin = OxmlElement('w:fldChar')
        fld_char_begin.set(qn('w:fldCharType'), 'begin')
        fld_begin.append(fld_char_begin)
        toc_para._element.append(fld_begin)

        # Field instruction
        fld_instr = OxmlElement('w:r')
        instr_text = OxmlElement('w:instrText')
        instr_text.set(qn('xml:space'), 'preserve')
        instr_text.text = ' TOC \\h \\z \\c "Table" '
        fld_instr.append(instr_text)
        toc_para._element.append(fld_instr)

        # Field separator
        fld_sep = OxmlElement('w:r')
        fld_char_sep = OxmlElement('w:fldChar')
        fld_char_sep.set(qn('w:fldCharType'), 'separate')
        fld_sep.append(fld_char_sep)
        toc_para._element.append(fld_sep)

        # Field result (placeholder - will be populated when field updates)
        # Add sample entries so user knows what to expect
        for i, caption in enumerate(table_list, 1):
            result_run = toc_para.add_run(f'{caption}\t{i+2}\n' if i < len(table_list) else f'{caption}\t{i+2}')
            result_run.font.name = self.template_fonts['body_font']
            result_run.font.size = Pt(12)

        # Field end
        fld_end = OxmlElement('w:r')
        fld_char_end = OxmlElement('w:fldChar')
        fld_char_end.set(qn('w:fldCharType'), 'end')
        fld_end.append(fld_char_end)
        toc_para._element.append(fld_end)

    def _populate_list_of_figures(self, doc, figure_list):
        r"""Populate the List of Figures section with a TOC field that references Figure captions.

        Creates a proper Word Table of Contents field that automatically updates
        when user right-clicks and selects "Update Field".

        Field format: { TOC \h \z \c "Figure" }
        - \h: Use hyperlinks
        - \z: Hide page numbers in Web Layout view
        - \c "Figure": Build table from Caption label "Figure"
        """
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn

        # Find the "List of Figures" heading
        list_of_figures_idx = None
        for idx, para in enumerate(doc.paragraphs):
            if para.text.strip() == 'List of Figures' and para.style.name == 'Heading 1':
                list_of_figures_idx = idx
                break

        if list_of_figures_idx is None:
            print("  Warning: 'List of Figures' heading not found in template")
            return

        # Remove template sample content between "List of Figures" heading and next heading
        paragraphs_to_remove = []
        for idx in range(list_of_figures_idx + 1, len(doc.paragraphs)):
            para = doc.paragraphs[idx]
            # Stop at next heading
            if para.style.name.startswith('Heading'):
                break
            # Mark paragraph for removal
            if para.text.strip():  # Remove any existing content
                paragraphs_to_remove.append(para)

        # Remove marked paragraphs
        for para in paragraphs_to_remove:
            p_element = para._element
            p_element.getparent().remove(p_element)

        # Insert TOC field paragraph after the List of Figures heading
        heading_element = doc.paragraphs[list_of_figures_idx]._element

        # Create a new paragraph for the TOC field
        toc_para = doc.add_paragraph()
        toc_para.style = 'Normal'

        # Move paragraph right after heading
        heading_element.addnext(toc_para._element)

        # Build TOC field: { TOC \h \z \c "Figure" }
        # Field begin
        fld_begin = OxmlElement('w:r')
        fld_char_begin = OxmlElement('w:fldChar')
        fld_char_begin.set(qn('w:fldCharType'), 'begin')
        fld_begin.append(fld_char_begin)
        toc_para._element.append(fld_begin)

        # Field instruction
        fld_instr = OxmlElement('w:r')
        instr_text = OxmlElement('w:instrText')
        instr_text.set(qn('xml:space'), 'preserve')
        instr_text.text = ' TOC \\h \\z \\c "Figure" '
        fld_instr.append(instr_text)
        toc_para._element.append(fld_instr)

        # Field separator
        fld_sep = OxmlElement('w:r')
        fld_char_sep = OxmlElement('w:fldChar')
        fld_char_sep.set(qn('w:fldCharType'), 'separate')
        fld_sep.append(fld_char_sep)
        toc_para._element.append(fld_sep)

        # Field result (placeholder - will be populated when field updates)
        # Add sample entries so user knows what to expect
        if figure_list:
            for i, caption in enumerate(figure_list, 1):
                result_run = toc_para.add_run(f'{caption}\t{i+10}\n' if i < len(figure_list) else f'{caption}\t{i+10}')
                result_run.font.name = self.template_fonts['body_font']
                result_run.font.size = Pt(12)
        else:
            # Add placeholder text if no figures
            result_run = toc_para.add_run('[No figures in document]')
            result_run.font.name = self.template_fonts['body_font']
            result_run.font.size = Pt(12)

        # Field end
        fld_end = OxmlElement('w:r')
        fld_char_end = OxmlElement('w:fldChar')
        fld_char_end.set(qn('w:fldCharType'), 'end')
        fld_end.append(fld_char_end)
        toc_para._element.append(fld_end)

    def _get_layout(self, prs, name_contains):
        """Get layout by partial name match."""
        for layout in prs.slide_layouts:
            if name_contains.lower() in layout.name.lower():
                return layout
        raise ValueError(f"Layout containing '{name_contains}' not found")

    def _insert_logos(self, slide, metadata=None):
        """Insert logos into picture placeholders on a slide.

        Args:
            slide: The slide to insert logos into
            metadata: Optional metadata dict with client_logo, footer_logo_left, and footer_logo_right paths
        """
        if not metadata:
            return

        from pptx.util import Inches

        # Get logo paths from metadata (relative to repo root)
        client_logo_path = metadata.get('client_logo', '')
        footer_logo_left_path = metadata.get('footer_logo_left', '')
        footer_logo_right_path = metadata.get('footer_logo_right', '')

        # Convert to absolute paths
        repo_root = Path(__file__).parent.parent.parent

        # Map of placeholder indices to logo paths
        # Title slide: idx=10 is client logo (top), idx=13 is footer left logo
        # Content slides: idx=13 or idx=14 is footer left logo
        logos = {}

        if client_logo_path:
            client_logo = repo_root / client_logo_path
            if client_logo.exists():
                logos[10] = str(client_logo)  # Top of slide (title slide only)

        # Footer left logo (consulting company)
        if footer_logo_left_path:
            footer_logo_left = repo_root / footer_logo_left_path
            if footer_logo_left.exists():
                logos[13] = str(footer_logo_left)  # Left footer
                logos[14] = str(footer_logo_left)  # Some layouts use idx 14

        # Insert logos into picture placeholders
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            if idx in logos and ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                try:
                    ph.insert_picture(logos[idx])
                except Exception as e:
                    # Silently continue if logo insertion fails
                    pass

        # Add footer right logo (EO Framework) as a new picture shape
        if footer_logo_right_path:
            footer_logo_right = repo_root / footer_logo_right_path
            if footer_logo_right.exists():
                try:
                    # Position and size to mirror the left logo
                    # Left logo is at: Left=0.34", Top=4.96", Width=2.34", Height=0.58"
                    # Slide width is 10.00"
                    # Right logo position: Left=7.31", Top=4.96"
                    left = Inches(7.31)
                    top = Inches(4.96)
                    width = Inches(2.34)
                    height = Inches(0.58)

                    slide.shapes.add_picture(
                        str(footer_logo_right),
                        left, top, width, height
                    )
                except Exception as e:
                    # Silently continue if logo insertion fails
                    pass

    def convert_md_to_pptx(self, md_file, output_dir):
        """Convert Markdown to PowerPoint presentation using branded template."""
        try:
            output_file = output_dir / f"{md_file.stem}.pptx"

            # Dry-run mode - just show what would be generated
            if self.dry_run:
                print(f"  🔍 Would generate: {md_file.name} → {output_file.name}")
                self.stats['md_to_pptx'] += 1
                return

            # Store current markdown file path for image path resolution
            # Solution root is 3 levels up from: solution-name/presales/raw/file.md
            self.current_solution_root = md_file.parent.parent.parent

            # Load branded template
            template_path = Path(__file__).parent.parent / 'doc-templates' / 'powerpoint' / 'EOFramework-Template-01.pptx'

            if not template_path.exists():
                raise FileNotFoundError(f"Template not found: {template_path}")

            prs = Presentation(str(template_path))

            # Clear demo slides (layouts-only approach)
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]

            # Read markdown
            with open(md_file, 'r', encoding='utf-8') as f:
                md_content = f.read()

            # Parse markdown into slides
            slides_data = self._parse_md_for_presentation(md_content)

            # Generate slides from layouts
            for slide_data in slides_data:
                self._add_slide_from_layout(prs, slide_data)

            # Set the presentation to open in Normal (Slide) view, not Slide Master view
            # This ensures it opens in presentation view instead of master view
            try:
                # Access the presentation part and set the view properties
                prs._element.attrib.pop('showMasterSp', None)  # Remove any master view settings
            except:
                pass  # If this fails, the file will still work

            prs.save(output_file)
            print(f"  ✅ {md_file.name} → {output_file.name}")
            self.stats['md_to_pptx'] += 1

        except Exception as e:
            error_msg = f"Error converting {md_file.name}: {str(e)}"
            print(f"  ❌ {error_msg}")
            self.stats['errors'].append(error_msg)

    def _parse_yaml_frontmatter(self, md_content):
        """Extract YAML frontmatter from markdown content.

        Returns: (metadata_dict, content_without_frontmatter)
        """
        metadata = {}
        content = md_content

        # Check if content starts with ---
        if md_content.strip().startswith('---'):
            parts = md_content.split('---', 2)
            if len(parts) >= 3:
                # Parse YAML frontmatter
                frontmatter = parts[1].strip()
                content = parts[2]

                # Simple YAML parser for key: value pairs
                for line in frontmatter.split('\n'):
                    line = line.strip()
                    if ':' in line:
                        key, value = line.split(':', 1)
                        metadata[key.strip()] = value.strip()

        return metadata, content

    def _parse_md_for_presentation(self, md_content):
        """Parse markdown content into presentation slides.

        Expects format:
        ---
        presentation_title: Title
        solution_name: Name
        ---
        # Main Title (H1) → Skip
        ### Slide Title (H3) → Content slides
        --- → Slide separator

        Stops processing when reaching ## Presentation Notes
        """
        # Extract metadata
        metadata, content = self._parse_yaml_frontmatter(md_content)

        slides = []
        lines = content.split('\n')
        current_slide = None
        first_h1_found = False
        slide_counter = 0

        i = 0
        while i < len(lines):
            # Keep original line with indentation for bullet parsing
            line_original = lines[i]
            line = line_original.strip()  # Stripped version for header checks

            # Stop processing when reaching Presentation Notes or similar sections
            if line.startswith('## ') and any(keyword in line.lower() for keyword in ['presentation notes', 'speaking points', 'q&a', 'appendix', 'powerpoint template']):
                break

            # Skip H2 "Slide Deck Structure" heading (don't create a slide for it)
            if line.startswith('## ') and 'slide deck structure' in line.lower():
                i += 1
                continue

            # Skip H1 (document title, not a slide)
            if line.startswith('# ') and not first_h1_found:
                first_h1_found = True
                i += 1
                continue

            # H3 creates slides (both title and content slides)
            if line.startswith('### '):
                if current_slide:
                    slides.append(current_slide)

                title_text = line[4:].strip()

                # Check if this is "Slide 1: Title Slide" - skip it (will use metadata instead)
                if 'slide 1:' in title_text.lower() and 'title slide' in title_text.lower():
                    # Skip this entire section until next --- or ###
                    i += 1
                    while i < len(lines):
                        if lines[i].strip().startswith('---') or lines[i].strip().startswith('###'):
                            i -= 1  # Back up one so we process the next section
                            break
                        i += 1
                    i += 1
                    continue

                slide_counter += 1

                # Regular content slide
                current_slide = {
                    'type': 'content',
                    'title': title_text,
                    'bullets': [],
                    'content': [],
                    'images': [],  # Store image paths for slide
                    'notes': [],  # Store speaker notes
                    'slide_number': slide_counter,
                    'layout_hint': self._determine_layout_hint(title_text, slide_counter),
                    'metadata': metadata  # Store metadata for logo insertion
                }

            # Content slide processing - check for images, bullets, numbered lists
            elif current_slide and current_slide['type'] == 'content':
                # Check for speaker notes marker
                if line.strip().startswith('**SPEAKER NOTES:**'):
                    # Set flag to capture everything after this as notes
                    current_slide['in_notes_section'] = True
                    i += 1
                    continue

                # If in notes section, capture all non-empty lines as notes (but stop at ---)
                if current_slide.get('in_notes_section', False):
                    # Stop collecting notes at slide separator or next slide header
                    if line.strip() == '---' or line.strip().startswith('###'):
                        current_slide['in_notes_section'] = False
                        # Don't continue - let it be processed normally
                    elif line.strip():  # Only add non-empty lines
                        current_slide['notes'].append(line.strip())
                        i += 1
                        continue

                # First check for markdown images: ![alt text](image_path)
                image_match = re.search(r'!\[([^\]]*)\]\(([^\)]+)\)', line)
                if image_match:
                    alt_text = image_match.group(1)
                    image_path = image_match.group(2)
                    current_slide['images'].append({
                        'alt': alt_text,
                        'path': image_path
                    })

                # Check for bullets at any indentation level - use original line with indentation
                elif line_original.lstrip().startswith('- ') or line_original.lstrip().startswith('* '):
                    stripped = line_original.lstrip()
                    # Calculate indentation level (2 spaces = 1 level)
                    indent_spaces = len(line_original) - len(stripped)
                    level = indent_spaces // 2

                    # Extract bullet text
                    bullet_text = stripped[2:].strip()
                    # Remove markdown formatting but preserve bold markers for processing
                    bullet_text = re.sub(r'🔴|✅', '', bullet_text).strip()

                    # Store bullet with its level
                    current_slide['bullets'].append((bullet_text, level))

                # Check for numbered list items (with indentation support)
                elif re.match(r'^\s*\d+\.', line_original):
                    stripped = line_original.lstrip()
                    if re.match(r'^\d+\.', stripped):
                        # Calculate indentation level
                        indent_spaces = len(line_original) - len(stripped)
                        level = indent_spaces // 2

                        item_text = re.sub(r'^\d+\.\s*', '', stripped).strip()
                        item_text = re.sub(r'\*\*(.+?)\*\*', r'\1', item_text)
                        current_slide['bullets'].append((item_text, level))

                # Check for table rows (inside content slide processing)
                elif line.startswith('|'):
                    if 'has_table' not in current_slide:
                        current_slide['has_table'] = True
                        current_slide['table_rows'] = []

                    # Parse table row (skip separator rows)
                    if not line.strip().replace('|', '').replace('-', '').replace(' ', '').replace(':', ''):
                        # Skip separator rows (|---|---|)
                        pass
                    else:
                        # Parse actual table row
                        cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                        if cells:
                            current_slide['table_rows'].append(cells)

            # Table detection and parsing (legacy - for non-content slides)
            elif line.startswith('|') and current_slide:
                if 'has_table' not in current_slide:
                    current_slide['has_table'] = True
                    current_slide['table_rows'] = []
                    # Check if it's a comparison table (Traditional vs Our Solution)
                    if 'traditional' in line.lower() or 'our solution' in line.lower():
                        current_slide['layout_hint'] = 'two_column'

                # Parse table row
                if not line.strip().replace('|', '').replace('-', '').strip():
                    # Skip separator rows (|---|---|)
                    pass
                else:
                    # Parse actual table row
                    cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                    if cells and current_slide['type'] == 'content':
                        # Store table rows for later processing
                        current_slide['table_rows'].append(cells)

            # Regular content lines
            elif line and not line.startswith('#') and not line.startswith('---') and current_slide:
                if current_slide['type'] == 'content':
                    current_slide['content'].append(line)

            i += 1

        if current_slide:
            slides.append(current_slide)

        # Post-process: Override layout_hint to 'table' if slide has table data
        # This ensures tables are rendered as PowerPoint tables, not bullets
        for slide in slides:
            if slide.get('has_table') and slide.get('table_rows') and len(slide['table_rows']) > 1:
                # Force table layout for any slide with table data
                slide['layout_hint'] = 'table'

        # Post-process: Convert table rows to formatted bullets (only for non-table layouts)
        for slide in slides:
            # Skip table-to-bullets conversion if using table layout
            if slide.get('layout_hint') == 'table':
                continue

            if slide.get('has_table') and slide.get('table_rows'):
                table_rows = slide['table_rows']
                if len(table_rows) > 1:
                    # First row is headers
                    headers = table_rows[0]
                    # Remaining rows are data
                    for row in table_rows[1:]:
                        if len(row) == len(headers):
                            # Format as "Header: Value" pairs
                            formatted_items = []
                            for i, (header, value) in enumerate(zip(headers, row)):
                                # Remove markdown bold from headers
                                clean_header = re.sub(r'\*\*(.+?)\*\*', r'\1', header)
                                clean_value = re.sub(r'\*\*(.+?)\*\*', r'\1', value)
                                if i == 0:
                                    # First column: use as main point (level 0)
                                    slide['bullets'].append((f"{clean_header}: {clean_value}", 0))
                                else:
                                    # Additional columns: indent as sub-points (level 1)
                                    slide['bullets'].append((f"{clean_header}: {clean_value}", 1))

        # Prepend title slide from metadata
        # Use current date instead of placeholder
        current_date = datetime.now().strftime('%B %d, %Y')

        title_slide = {
            'type': 'title',
            'title': metadata.get('presentation_title', 'Presentation Title'),
            'subtitle': metadata.get('solution_name', 'Solution Name'),
            'presenter': metadata.get('presenter_name', 'Presenter Name'),
            'date': current_date,
            'layout_hint': 'title',
            'metadata': metadata  # Store full metadata for logo paths
        }
        slides.insert(0, title_slide)

        return slides

    def _determine_layout_hint(self, title, slide_number):
        """Determine optimal layout based on slide title and position."""
        title_lower = title.lower()

        # Title slide keywords - only match "thank you" exactly, not "next steps"
        if any(keyword in title_lower for keyword in ['thank you', 'contact']):
            return 'thank_you'

        # Bullet points layout for next steps
        if 'next steps' in title_lower:
            return 'bullet_points'

        # Table layout keywords (timeline, milestones, schedule)
        if any(keyword in title_lower for keyword in ['timeline', 'milestone', 'schedule', 'roadmap']):
            return 'table'

        # Two-column layout keywords (comparisons, partnerships, advantages)
        if any(keyword in title_lower for keyword in ['why this', 'differentiated', 'comparison', 'vs', 'versus', 'partnership advantage', 'our partnership', 'why partner']):
            return 'two_column'

        # Visual content keywords (diagrams, architecture)
        if any(keyword in title_lower for keyword in ['overview', 'architecture', 'diagram', 'visual']):
            return 'visual'

        # Data visualization keywords (financials, metrics, charts)
        if any(keyword in title_lower for keyword in ['investment', 'roi', 'cost', 'pricing', 'financials', 'metrics']):
            return 'data_viz'

        # Key points layout keywords (highlights) - use single column
        # Note: removed 'advantage', 'partnership', 'why partner' as they typically need two columns
        if any(keyword in title_lower for keyword in ['key points', 'highlights', 'benefits summary']):
            return 'single'

        # Default to single column
        return 'single'

    def _add_slide_from_layout(self, prs, slide_data):
        """Add slide using appropriate template layout based on content."""
        try:
            if slide_data['type'] == 'title':
                # Use "EO Title Slide" layout
                layout = self._get_layout(prs, "Title")
                slide = prs.slides.add_slide(layout)

                # Fill placeholders based on template structure
                for placeholder in slide.placeholders:
                    try:
                        if 'title' in placeholder.name.lower() or placeholder.placeholder_format.idx == 12:
                            placeholder.text = slide_data['title']
                        elif 'subtitle' in placeholder.name.lower() or placeholder.placeholder_format.idx == 14:
                            placeholder.text = slide_data.get('subtitle', '')
                        elif placeholder.placeholder_format.idx == 15:
                            presenter_text = f"{slide_data.get('presenter', 'Presenter Name')} | {slide_data.get('date', datetime.now().strftime('%B %d, %Y'))}"
                            placeholder.text = presenter_text
                    except:
                        pass

                # Insert logos using metadata
                metadata = slide_data.get('metadata', {})
                self._insert_logos(slide, metadata)

            else:  # content slide - use layout hint
                layout_hint = slide_data.get('layout_hint', 'single')

                # Select appropriate layout based on hint
                if layout_hint == 'thank_you':
                    layout = self._get_layout(prs, "Thank You")
                elif layout_hint == 'bullet_points':
                    layout = self._get_layout(prs, "Bullet Points")
                elif layout_hint == 'table':
                    layout = self._get_layout(prs, "Table")
                elif layout_hint == 'two_column':
                    layout = self._get_layout(prs, "Two Column")
                elif layout_hint == 'visual':
                    layout = self._get_layout(prs, "Visual Content")
                elif layout_hint == 'data_viz':
                    layout = self._get_layout(prs, "Data Visualization")
                else:
                    layout = self._get_layout(prs, "Single Column")

                slide = prs.slides.add_slide(layout)

                # Fill title (different index for each layout)
                title_idx = 12 if layout_hint == 'thank_you' else 10
                try:
                    title_placeholder = slide.placeholders[title_idx]
                    title_placeholder.text = slide_data['title']

                    # For table layout: left-align and bold the title
                    if layout_hint == 'table':
                        from pptx.enum.text import PP_ALIGN
                        if title_placeholder.text_frame and title_placeholder.text_frame.paragraphs:
                            for paragraph in title_placeholder.text_frame.paragraphs:
                                paragraph.alignment = PP_ALIGN.LEFT
                                for run in paragraph.runs:
                                    run.font.bold = True
                except:
                    # Fallback: try to find title placeholder by name
                    for placeholder in slide.placeholders:
                        try:
                            if 'title' in placeholder.name.lower() or placeholder.placeholder_format.idx in [8, 10, 12]:
                                placeholder.text = slide_data['title']
                                # Apply same formatting for table layout
                                if layout_hint == 'table':
                                    from pptx.enum.text import PP_ALIGN
                                    if placeholder.text_frame and placeholder.text_frame.paragraphs:
                                        for paragraph in placeholder.text_frame.paragraphs:
                                            paragraph.alignment = PP_ALIGN.LEFT
                                            for run in paragraph.runs:
                                                run.font.bold = True
                                break
                        except:
                            pass

                # Fill body content based on layout type
                if layout_hint in ['thank_you']:
                    # Thank you slide - idx=15 is content area below title
                    self._fill_content_placeholder(slide, slide_data, [15])
                elif layout_hint == 'bullet_points':
                    # Bullet points slide - idx=14 is content
                    self._fill_content_placeholder(slide, slide_data, [14])
                elif layout_hint == 'table':
                    # Table slide - idx=14 is table placeholder
                    self._fill_table_placeholder(slide, slide_data, 14)
                elif layout_hint == 'two_column':
                    # Two column - idx=16 (left), idx=17 (right) - EO Two Column layout
                    self._fill_content_placeholder(slide, slide_data, [16, 17])
                elif layout_hint == 'data_viz':
                    # Data visualization - idx=11 is content
                    self._fill_content_placeholder(slide, slide_data, [11])
                elif layout_hint == 'visual':
                    # Visual content - idx=16 is content (updated template)
                    self._fill_content_placeholder(slide, slide_data, [16])
                else:
                    # Single column - idx=14 is content (EO Single Column layout)
                    self._fill_content_placeholder(slide, slide_data, [14])

                # Insert logos on content slides (footer logo)
                metadata = slide_data.get('metadata', {})
                if metadata:
                    self._insert_logos(slide, metadata)

                # Insert content images (e.g., architecture diagrams)
                images = slide_data.get('images', [])
                if images:
                    self._insert_slide_images(slide, slide_data, layout_hint)

                # Add speaker notes if present
                notes = slide_data.get('notes', [])
                if notes:
                    try:
                        notes_slide = slide.notes_slide
                        text_frame = notes_slide.notes_text_frame
                        # Join all notes with newlines and clean up markdown formatting
                        notes_text = '\n'.join(notes)
                        # Remove markdown formatting from notes
                        notes_text = re.sub(r'\*([^\*]+)\*', r'\1', notes_text)  # Remove italic *text*
                        notes_text = re.sub(r'_([^_]+)_', r'\1', notes_text)     # Remove italic _text_
                        if text_frame:
                            text_frame.text = notes_text
                    except Exception as e:
                        # Notes may not be available for all slide layouts
                        pass

        except Exception as e:
            # Log error but continue with other slides
            print(f"    ⚠️  Warning: Could not create slide '{slide_data.get('title', 'Unknown')}': {str(e)}")

    def _apply_bold_formatting(self, paragraph, text):
        """Parse text with **bold** markers and apply formatting to paragraph runs."""
        # Split text by ** markers to find bold sections
        parts = text.split('**')

        # Clear existing text
        if paragraph.runs:
            for run in paragraph.runs:
                run.text = ''

        # Add parts with alternating bold formatting
        for i, part in enumerate(parts):
            if part:  # Skip empty parts
                run = paragraph.add_run()
                run.text = part
                # Odd indices (1, 3, 5...) are between ** markers, so make them bold
                if i % 2 == 1:
                    run.font.bold = True

    def _fill_table_placeholder(self, slide, slide_data, table_placeholder_idx):
        """Fill table placeholder with data from markdown table."""
        try:
            from pptx.util import Pt
            from pptx.dml.color import RGBColor

            # Get table data
            table_rows = slide_data.get('table_rows', [])
            if not table_rows or len(table_rows) < 2:
                return

            # Get the table placeholder
            table_placeholder = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == table_placeholder_idx:
                    table_placeholder = ph
                    break

            if not table_placeholder:
                return

            # Access the table shape
            graphic_frame = table_placeholder.insert_table(rows=len(table_rows), cols=len(table_rows[0]))
            table = graphic_frame.table

            # Set column widths (proportional distribution)
            total_width = graphic_frame.width

            # Different column width distributions based on column count
            if len(table.columns) == 4:
                # Phase No | Phase Description | Timeline | Key Deliverables
                col_widths = [0.12, 0.25, 0.12, 0.51]
            elif len(table.columns) == 5:
                # Cost Item | Cost per Item | Credits | Total Cost to Client | Comments
                # Optimized for Investment Summary tables
                col_widths = [0.25, 0.15, 0.12, 0.18, 0.30]
            else:
                # Default: equal distribution for other column counts
                col_widths = [1.0 / len(table.columns)] * len(table.columns)

            # Apply column widths
            for i, width_ratio in enumerate(col_widths):
                if i < len(table.columns):
                    table.columns[i].width = int(total_width * width_ratio)

            # Fill table with data
            for row_idx, row_data in enumerate(table_rows):
                for col_idx, cell_text in enumerate(row_data):
                    cell = table.rows[row_idx].cells[col_idx]

                    # Clean markdown formatting from cell text
                    clean_text = re.sub(r'\*\*(.+?)\*\*', r'\1', cell_text)
                    cell.text = clean_text

                    # Format header row (row 0)
                    if row_idx == 0:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.bold = True
                                run.font.size = Pt(14)
                                run.font.color.rgb = RGBColor(255, 255, 255)
                        # Header cell background (brand red #A01C02)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(160, 28, 2)
                    else:
                        # Data row formatting
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(11)
                        # Data cell background (light gray #E7E6E6)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(231, 230, 230)

        except Exception as e:
            print(f"    ⚠️  Warning: Could not fill table: {str(e)}")

    def _fill_content_placeholder(self, slide, slide_data, placeholder_indices):
        """Fill content placeholders with bullets or text (supports hierarchical levels)."""
        try:
            from pptx.util import Pt
            from pptx.dml.color import RGBColor

            bullets = slide_data.get('bullets', [])
            content_lines = slide_data.get('content', [])

            # Process bullets - they are now tuples of (text, level)
            # Keep ** markers for bold formatting
            processed_bullets = []
            for item in bullets:
                if isinstance(item, tuple):
                    text, level = item
                    # Keep ** markers for bold text formatting (don't remove them)
                    processed_bullets.append((text, level))
                else:
                    # Legacy format - just text, assume level 0
                    processed_bullets.append((str(item), 0))

            # Add plain content lines as level 0
            for line in content_lines:
                if line.strip() and not line.startswith('**') and not line.startswith('|'):
                    processed_bullets.append((line, 0))

            if not processed_bullets:
                return

            # Fill placeholders
            for idx in placeholder_indices:
                try:
                    placeholder = slide.placeholders[idx]
                    tf = placeholder.text_frame
                    tf.clear()

                    # For two-column layouts, split content between columns
                    if len(placeholder_indices) == 2:
                        mid_point = len(processed_bullets) // 2
                        if idx == placeholder_indices[0]:
                            content_to_add = processed_bullets[:mid_point]
                        else:
                            content_to_add = processed_bullets[mid_point:]
                    else:
                        content_to_add = processed_bullets

                    # Add content with proper indentation levels
                    if content_to_add:
                        # First item
                        first_text, first_level = content_to_add[0]
                        # Set level for first paragraph (color inherited from template)
                        if tf.paragraphs:
                            p = tf.paragraphs[0]
                            p.level = first_level
                            # Apply bold formatting if ** markers present
                            if '**' in first_text:
                                self._apply_bold_formatting(p, first_text)
                            else:
                                p.text = first_text

                        # Remaining items
                        for text, level in content_to_add[1:]:
                            if text.strip():
                                p = tf.add_paragraph()
                                p.level = level  # Set the indentation level (color inherited from template)
                                # Apply bold formatting if ** markers present
                                if '**' in text:
                                    self._apply_bold_formatting(p, text)
                                else:
                                    p.text = text

                    # For single placeholder, we're done
                    if len(placeholder_indices) == 1:
                        break

                except KeyError as e:
                    # Placeholder index doesn't exist, try next one
                    continue
                except Exception as e:
                    # Other error, continue to next placeholder
                    continue

        except Exception as e:
            pass

    def _insert_slide_images(self, slide, slide_data, layout_hint):
        """Insert images from markdown into slide picture placeholders."""
        try:
            from pptx.util import Inches

            images = slide_data.get('images', [])
            if not images:
                return

            # Determine picture placeholder index based on layout
            picture_idx = None
            if layout_hint == 'visual':
                # Visual Content layout has main image at idx=15
                picture_idx = 15
            elif layout_hint == 'data_viz':
                # Data Visualization layout might also have image placeholder
                picture_idx = 15
            else:
                # Other layouts - try to find any picture placeholder that's not the footer
                for placeholder in slide.placeholders:
                    try:
                        if placeholder.placeholder_format.type == 18:  # PICTURE type
                            # Skip footer logos (idx 13, 14)
                            if placeholder.placeholder_format.idx not in [13, 14]:
                                picture_idx = placeholder.placeholder_format.idx
                                break
                    except:
                        continue

            if picture_idx is None:
                return

            # Insert the first image (for now, handle only one image per slide)
            image_data = images[0]
            image_path = image_data['path']

            # Resolve relative path (relative to the solution root)
            # The path in markdown is relative to the solution root
            if not image_path.startswith('/'):
                # Construct absolute path
                # We're in: solution-name/presales/raw/file.md
                # Image path is: assets/images/diagram.png (relative to solution root)
                solution_root = self.current_solution_root
                full_image_path = solution_root / image_path
            else:
                full_image_path = Path(image_path)

            # Check if image file exists
            if not full_image_path.exists():
                print(f"    ⚠️  Warning: Image not found: {full_image_path}")
                return

            # Insert image directly to slide (not using placeholder's insert_picture)
            try:
                placeholder = slide.placeholders[picture_idx]

                # Get placeholder dimensions for positioning
                ph_left = placeholder.left
                ph_top = placeholder.top
                ph_width = placeholder.width
                ph_height = placeholder.height

                # Get image natural size
                from PIL import Image as PILImage
                with PILImage.open(full_image_path) as pil_img:
                    img_width, img_height = pil_img.size

                # Calculate aspect ratios
                img_aspect = img_width / img_height
                ph_aspect = ph_width / ph_height

                # Calculate size to fit without cropping
                if img_aspect > ph_aspect:
                    # Image is wider - fit to width
                    new_width = ph_width
                    new_height = int(ph_width / img_aspect)
                else:
                    # Image is taller - fit to height
                    new_height = ph_height
                    new_width = int(ph_height * img_aspect)

                # Center the image in placeholder area
                left = ph_left + (ph_width - new_width) // 2
                top = ph_top + (ph_height - new_height) // 2

                # Remove existing picture shapes in the content area to avoid duplicates
                # Keep only footer logos (typically in bottom area)
                from pptx.util import Inches
                shapes_to_remove = []
                for shape in slide.shapes:
                    # Remove picture shapes that are in the content area (not footer)
                    if shape.shape_type == 13:  # PICTURE type
                        # Check if it's in the main content area (not footer)
                        # Footer is typically below 6.5 inches from top
                        if shape.top < Inches(6.5):
                            shapes_to_remove.append(shape)

                # Remove identified shapes
                for shape in shapes_to_remove:
                    sp = shape.element
                    sp.getparent().remove(sp)

                # Add picture directly to slide at calculated position
                pic = slide.shapes.add_picture(str(full_image_path), left, top, new_width, new_height)

            except KeyError:
                print(f"    ⚠️  Warning: Picture placeholder idx={picture_idx} not found in slide")
            except Exception as e:
                print(f"    ⚠️  Warning: Could not insert image: {str(e)}")

        except Exception as e:
            # Don't fail slide creation if image insertion fails
            pass

    def print_summary(self):
        """Print generation summary."""
        if self.quiet:
            return

        print("\n" + "="*60)
        if self.dry_run:
            print("📊 Dry Run Summary (No Files Generated)")
        else:
            print("📊 Generation Summary")
        print("="*60)

        action = "Would generate" if self.dry_run else "Generated"
        print(f"✅ CSV → Excel:      {self.stats['csv_to_xlsx']} files")
        print(f"✅ MD → Word:        {self.stats['md_to_docx']} files")
        print(f"✅ MD → PowerPoint:  {self.stats['md_to_pptx']} files")

        total = self.stats['csv_to_xlsx'] + self.stats['md_to_docx'] + self.stats['md_to_pptx']
        print(f"\n📁 Total {action.lower()}:  {total} files")

        if self.stats['skipped'] > 0:
            print(f"⏭️  Skipped (filtered): {self.stats['skipped']} files")

        if self.stats['errors']:
            print(f"\n❌ Errors: {len(self.stats['errors'])}")
            for error in self.stats['errors']:
                print(f"   • {error}")
        elif not self.dry_run:
            print("\n✨ All files generated successfully!")
        print("="*60)


def main():
    parser = argparse.ArgumentParser(
        description='Build professional output documents from raw source files.',
        epilog='Examples:\n'
               '  %(prog)s --path solutions/aws/ai/idp\n'
               '  %(prog)s --path solutions/aws/ai/idp --type excel\n'
               '  %(prog)s --path solutions/aws/ai/idp --file statement-of-work.md\n'
               '  %(prog)s --path solutions/aws/ai/idp --dir presales\n'
               '  %(prog)s --path solutions/aws/ai/idp --dry-run\n',
        formatter_class=argparse.RawDescriptionHelpFormatter
    )
    parser.add_argument(
        '--path',
        type=str,
        required=True,
        help='Path to solution directory or solution-template'
    )
    parser.add_argument(
        '--type',
        choices=['excel', 'word', 'pptx'],
        nargs='+',
        help='Generate only specific file types (can specify multiple)'
    )
    parser.add_argument(
        '--file',
        help='Generate only this specific source file (e.g., statement-of-work.md)'
    )
    parser.add_argument(
        '--files',
        help='Comma-separated list of source files to generate'
    )
    parser.add_argument(
        '--dir',
        choices=['presales', 'delivery'],
        help='Process only files in this directory'
    )
    parser.add_argument(
        '--force',
        action='store_true',
        help='Regenerate files even if output already exists'
    )
    parser.add_argument(
        '--dry-run',
        action='store_true',
        help='Show what would be generated without actually generating files'
    )
    parser.add_argument(
        '--quiet',
        action='store_true',
        help='Suppress progress output (errors only)'
    )
    parser.add_argument(
        '--verbose',
        action='store_true',
        help='Show detailed processing information including skipped files'
    )

    args = parser.parse_args()

    # Validate path exists
    path = Path(args.path)
    if not path.exists():
        print(f"❌ Error: Path does not exist: {path}")
        sys.exit(1)

    # Parse file filters
    file_filter = None
    if args.file:
        file_filter = [args.file]
    elif args.files:
        file_filter = [f.strip() for f in args.files.split(',')]

    # Generate outputs
    generator = OutputGenerator(
        base_path=path,
        type_filter=args.type,
        file_filter=file_filter,
        dir_filter=args.dir,
        force=args.force,
        dry_run=args.dry_run,
        quiet=args.quiet,
        verbose=args.verbose
    )
    success = generator.generate_all()

    sys.exit(0 if success else 1)


if __name__ == '__main__':
    main()
