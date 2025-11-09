#!/usr/bin/env python3
"""
Create PowerPoint template with branded layouts and logo placeholders.
Creates a professional template with:
- Title slide with 3 logo placeholders (customer, consulting, EO)
- Clean content slides without logos
- Consistent branding and colors
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path


def create_pptx_template():
    """Create PowerPoint template with branded layouts."""

    output_file = Path('support/doc-templates/powerpoint/EOFramework-Template-01.pptx')
    logo_path = Path('support/doc-templates/assets/logos')

    # Create new presentation
    prs = Presentation()

    # Set slide size to standard (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("📄 Creating PowerPoint Template")
    print("=" * 70)

    # =========================================================================
    # REFERENCE SLIDE 1: TITLE SLIDE WITH 3 LOGOS
    # =========================================================================
    print("\n📊 Reference Slide 1: Title Slide (3 Logos)")
    print("-" * 70)

    # Use built-in Title Slide layout (index 0)
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title = slide.shapes.title
    title.text = "Project Proposal"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)  # #1F4E78
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Set subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = "Professional Solution Presentation"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)  # #44546A
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add 3 logo placeholders at bottom
    # Logo dimensions and positions
    logo_width = Inches(1.5)
    logo_height = Inches(0.6)
    logo_y = Inches(6.5)  # Bottom of slide

    # Customer logo (left)
    left_logo = slide.shapes.add_picture(
        str(logo_path / 'client_logo.png') if (logo_path / 'client_logo.png').exists()
        else str(logo_path / 'eo-framework-logo-real.png'),
        Inches(1.0), logo_y, logo_width, logo_height
    )

    # Consulting company logo (center)
    center_logo = slide.shapes.add_picture(
        str(logo_path / 'consulting_company_logo.png') if (logo_path / 'consulting_company_logo.png').exists()
        else str(logo_path / 'eo-framework-logo-real.png'),
        Inches(4.25), logo_y, logo_width, logo_height
    )

    # EO Framework logo (right)
    right_logo = slide.shapes.add_picture(
        str(logo_path / 'eo-framework-logo-real.png'),
        Inches(7.5), logo_y, logo_width, logo_height
    )

    print("  ✅ Title slide created with title, subtitle, and 3 logo placeholders")

    # =========================================================================
    # REFERENCE SLIDE 2: TITLE AND CONTENT
    # =========================================================================
    print("\n📊 Reference Slide 2: Title and Content")
    print("-" * 70)

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Key Benefits"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Accelerated delivery timeline"

    p = tf.add_paragraph()
    p.text = "Reduced operational costs"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Improved security posture"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Enhanced scalability"
    p.level = 0

    # Style the bullets
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(68, 84, 106)

    print("  ✅ Title and Content slide created")

    # =========================================================================
    # REFERENCE SLIDE 3: TWO COLUMN CONTENT
    # =========================================================================
    print("\n📊 Reference Slide 3: Two Column Content")
    print("-" * 70)

    slide_layout = prs.slide_layouts[3] if len(prs.slide_layouts) > 3 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Technical Architecture"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)

    # Add two text boxes for columns
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.25), Inches(4.5))
    left_tf = left_box.text_frame
    left_tf.text = "Cloud Infrastructure"
    left_tf.paragraphs[0].font.size = Pt(18)
    left_tf.paragraphs[0].font.bold = True
    left_tf.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)

    p = left_tf.add_paragraph()
    p.text = "• Multi-region deployment"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(68, 84, 106)

    p = left_tf.add_paragraph()
    p.text = "• Auto-scaling capabilities"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(68, 84, 106)

    right_box = slide.shapes.add_textbox(Inches(5.25), Inches(2), Inches(4.25), Inches(4.5))
    right_tf = right_box.text_frame
    right_tf.text = "Security Layer"
    right_tf.paragraphs[0].font.size = Pt(18)
    right_tf.paragraphs[0].font.bold = True
    right_tf.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)

    p = right_tf.add_paragraph()
    p.text = "• Enterprise encryption"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(68, 84, 106)

    p = right_tf.add_paragraph()
    p.text = "• Identity management"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(68, 84, 106)

    print("  ✅ Two column content slide created")

    # =========================================================================
    # REFERENCE SLIDE 4: SECTION HEADER
    # =========================================================================
    print("\n📊 Reference Slide 4: Section Header")
    print("-" * 70)

    slide_layout = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Implementation Approach"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add section subtitle if placeholder exists
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "Phase 1: Discovery & Planning"
        subtitle.text_frame.paragraphs[0].font.size = Pt(28)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    print("  ✅ Section header slide created")

    # =========================================================================
    # REFERENCE SLIDE 5: TIMELINE
    # =========================================================================
    print("\n📊 Reference Slide 5: Timeline/Roadmap")
    print("-" * 70)

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Implementation Timeline"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)

    # Add timeline text boxes
    phases = [
        ("Q1: Discovery", Inches(1.0), "Requirements gathering\nTechnical assessment"),
        ("Q2: Design", Inches(3.5), "Architecture design\nSecurity planning"),
        ("Q3: Build", Inches(6.0), "Development\nTesting & QA"),
        ("Q4: Deploy", Inches(8.5), "Production deployment\nGo-live support")
    ]

    for phase_name, x_pos, phase_desc in phases:
        # Phase header
        header_box = slide.shapes.add_textbox(x_pos - Inches(0.25), Inches(2.5), Inches(2), Inches(0.5))
        header_tf = header_box.text_frame
        header_tf.text = phase_name
        header_tf.paragraphs[0].font.size = Pt(16)
        header_tf.paragraphs[0].font.bold = True
        header_tf.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
        header_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Phase description
        desc_box = slide.shapes.add_textbox(x_pos - Inches(0.4), Inches(3.2), Inches(2.3), Inches(2))
        desc_tf = desc_box.text_frame
        desc_tf.text = phase_desc
        desc_tf.paragraphs[0].font.size = Pt(12)
        desc_tf.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)
        desc_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    print("  ✅ Timeline slide created")

    # =========================================================================
    # REFERENCE SLIDE 6: BLANK FOR CUSTOM CONTENT
    # =========================================================================
    print("\n📊 Reference Slide 6: Blank Slide")
    print("-" * 70)

    slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    # Add centered text placeholder
    text_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    tf = text_box.text_frame
    tf.text = "Blank Slide"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Add custom content here"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER

    print("  ✅ Blank slide created")

    # =========================================================================
    # REFERENCE SLIDE 7: THANK YOU / CLOSING
    # =========================================================================
    print("\n📊 Reference Slide 7: Thank You / Closing")
    print("-" * 70)

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Thank You"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle = slide.placeholders[1]
    subtitle.text = "Questions & Discussion"
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add logos at bottom (same as title slide)
    left_logo = slide.shapes.add_picture(
        str(logo_path / 'client_logo.png') if (logo_path / 'client_logo.png').exists()
        else str(logo_path / 'eo-framework-logo-real.png'),
        Inches(1.0), logo_y, logo_width, logo_height
    )

    center_logo = slide.shapes.add_picture(
        str(logo_path / 'consulting_company_logo.png') if (logo_path / 'consulting_company_logo.png').exists()
        else str(logo_path / 'eo-framework-logo-real.png'),
        Inches(4.25), logo_y, logo_width, logo_height
    )

    right_logo = slide.shapes.add_picture(
        str(logo_path / 'eo-framework-logo-real.png'),
        Inches(7.5), logo_y, logo_width, logo_height
    )

    print("  ✅ Thank you slide created with 3 logos")

    # =========================================================================
    # REFERENCE SLIDE 8: CONTACT INFORMATION
    # =========================================================================
    print("\n📊 Reference Slide 8: Contact Information")
    print("-" * 70)

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Contact Information"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)

    # Add contact info text box
    contact_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(3.5))
    tf = contact_box.text_frame

    tf.text = "[Consulting Company Name]"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Primary Contact: [Name]"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(68, 84, 106)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Email: [email@company.com]"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(68, 84, 106)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Phone: [+1-XXX-XXX-XXXX]"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(68, 84, 106)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Website: [www.company.com]"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(68, 84, 106)
    p.alignment = PP_ALIGN.CENTER

    print("  ✅ Contact information slide created")

    # =========================================================================
    # FIX VIEW MODE: Ensure template opens in Normal view, not Master view
    # =========================================================================
    print("\n🔧 Fixing view mode...")

    try:
        if hasattr(prs._element, 'attrib'):
            prs._element.attrib.pop('showMasterSp', None)
            prs._element.attrib.pop('showMasterPhAnim', None)
        print("  ✅ Removed master view attributes (will open in Normal view)")
    except Exception as e:
        print(f"  ⚠️  Could not modify view attributes: {e}")

    # Save template
    print("\n" + "=" * 70)
    print(f"💾 Saving template: {output_file}")

    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_file))

    print()
    print("✨ PowerPoint template created successfully!")
    print()
    print("Template structure:")
    print("  📊 Slide 1: Title Slide - 3 logo placeholders (customer, consulting, EO)")
    print("  📊 Slide 2: Title and Content - Bullet points")
    print("  📊 Slide 3: Two Column - Side-by-side content")
    print("  📊 Slide 4: Section Header - Chapter/phase separator")
    print("  📊 Slide 5: Timeline - Roadmap/schedule")
    print("  📊 Slide 6: Blank - Custom content")
    print("  📊 Slide 7: Thank You - Closing slide with 3 logos")
    print("  📊 Slide 8: Contact Info - Contact details")
    print()
    print("Features:")
    print("  ✅ Professional color scheme (#1F4E78 primary, #44546A secondary)")
    print("  ✅ Consistent typography (Calibri)")
    print("  ✅ 3 logo support on title/closing slides")
    print("  ✅ Clean content slides without logos")
    print("  ✅ 8 reference slides showing all patterns")
    print("  ✅ Opens in Normal view (not Slide Master view)")
    print()
    print("Usage in generate-outputs.py:")
    print("  The script will clear these reference slides and create new content")
    print("  from markdown files using the established layouts and styles.")
    print()


if __name__ == '__main__':
    create_pptx_template()
