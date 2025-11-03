#!/usr/bin/env python3
"""
Test v2: Clone the demo slide and replace text instead of using placeholders.
"""

from pptx import Presentation
from pathlib import Path
from copy import deepcopy
import re

def extract_title_from_markdown(md_file):
    """Extract the first H1 title from markdown file."""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Find first H1 header
    match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
    if match:
        return match.group(1).strip()

    return "Untitled Presentation"

def create_title_slide_from_demo():
    """Create title slide by cloning demo slide."""

    # Paths
    template_path = Path('support/doc-templates/powerpoint/EOFramework-Template-01.pptx')
    md_file = Path('solution-template/sample-provider/sample-category/sample-solution/presales/raw/executive-presentation.md')
    output_file = Path('test-title-slide-v2.pptx')

    # Load template
    print(f"📂 Loading template: {template_path}")
    template_prs = Presentation(str(template_path))

    # Get the first demo slide (has logos already)
    print(f"✅ Template has {len(template_prs.slides)} demo slides")
    demo_slide = template_prs.slides[0]

    print(f"\n📋 Demo slide structure:")
    for i, shape in enumerate(demo_slide.shapes):
        shape_type = str(shape.shape_type).split('.')[-1]
        name = shape.name
        text = shape.text[:30] if hasattr(shape, 'text') and shape.text else ''
        print(f"  {i}. {shape_type:20} | {name:25} | {text}")

    # Extract title from markdown
    title_text = extract_title_from_markdown(md_file)
    print(f"\n📝 Title from markdown: {title_text}")

    # Create new presentation from template
    prs = Presentation(str(template_path))

    # Remove all slides except the first one
    print(f"\n🗑️  Removing slides 2-{len(prs.slides)}...")
    while len(prs.slides) > 1:
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]

    # Now we have just the first slide with logos
    slide = prs.slides[0]

    # Replace the title text
    print(f"\n✏️  Replacing title text...")
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            # Look for the shape with "Presentation Title"
            if 'Presentation Title' in shape.text:
                shape.text = title_text
                print(f"  ✅ Updated: '{shape.text}' (in {shape.name})")

    # Save
    print(f"\n💾 Saving: {output_file}")
    prs.save(str(output_file))

    print(f"\n✨ Done! Open {output_file} to verify.")
    print(f"   Windows path: {output_file.absolute()}")
    print(f"\nThis version should have:")
    print(f"  ✓ 1 slide only")
    print(f"  ✓ All 3 logos from the demo slide")
    print(f"  ✓ Your title: '{title_text}'")

if __name__ == '__main__':
    create_title_slide_from_demo()
