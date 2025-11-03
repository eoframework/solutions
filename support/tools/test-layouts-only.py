#!/usr/bin/env python3
"""
Test: Generate slides purely from layouts (no demo slides).
Following the clean approach with insert_picture() on picture placeholders.
"""

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pathlib import Path
import re

def get_layout(prs, name_contains):
    """Get layout by partial name match."""
    for layout in prs.slide_layouts:
        if name_contains.lower() in layout.name.lower():
            return layout
    raise ValueError(f"Layout containing '{name_contains}' not found")

def extract_metadata_from_markdown(md_file):
    """Extract title and subtitle from markdown."""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    metadata = {'title': 'Untitled', 'subtitle': ''}

    # Find H1 (title)
    match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
    if match:
        metadata['title'] = match.group(1).strip()

    # Find H2 (subtitle)
    match = re.search(r'^##\s+(.+)$', content, re.MULTILINE)
    if match:
        metadata['subtitle'] = match.group(1).strip()

    return metadata

def test_layouts_only():
    """Test generating slides from layouts without demo slides."""

    # Paths
    template_path = Path('support/doc-templates/powerpoint/EOFramework-Template-01.pptx')
    logo_path = Path('support/doc-templates/powerpoint')
    md_file = Path('solution-template/sample-provider/sample-category/sample-solution/presales/raw/executive-presentation.md')
    output_file = Path('test-layouts-only.pptx')

    print(f"📂 Loading template: {template_path}")
    prs = Presentation(str(template_path))

    # Clear all demo slides (we're starting fresh from layouts)
    print(f"🗑️  Removing {len(prs.slides)} demo slides...")
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    print(f"✅ Starting fresh with {len(prs.slide_layouts)} layouts")
    print()

    # Extract metadata
    metadata = extract_metadata_from_markdown(md_file)
    print(f"📝 Metadata: {metadata['title']}")
    print()

    # TEST 1: Create Title Slide from layout
    print("=" * 60)
    print("TEST 1: Title Slide from Layout")
    print("=" * 60)

    title_layout = get_layout(prs, "Title")
    slide = prs.slides.add_slide(title_layout)
    print(f"✅ Created slide from layout: {title_layout.name}")

    # Set title (index 12 based on template analysis)
    try:
        title_ph = slide.placeholders[12]
        title_ph.text = metadata['title']
        print(f"✅ Set title [12]: {metadata['title']}")
    except:
        print(f"⚠️  Could not set title")

    # Set subtitle (index 14 based on template analysis)
    try:
        subtitle_ph = slide.placeholders[14]
        subtitle_ph.text = metadata['subtitle']
        print(f"✅ Set subtitle [14]: {metadata['subtitle']}")
    except:
        print(f"⚠️  Could not set subtitle")

    # Set presenter info (index 15)
    try:
        presenter_ph = slide.placeholders[15]
        presenter_ph.text = "Presenter Name | November 02, 2025"
        print(f"✅ Set presenter [15]")
    except:
        print(f"⚠️  Could not set presenter")

    # Find picture placeholders and insert logos
    print(f"\\n🖼️  Looking for picture placeholders...")

    pic_placeholders = [
        ph for ph in slide.placeholders
        if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE
    ]

    print(f"Found {len(pic_placeholders)} picture placeholders")

    logos = [
        ('client_logo.png', 'Client Logo'),
        ('consulting_company_logo.png', 'Consulting Company Logo'),
    ]

    for i, (logo_file, logo_name) in enumerate(logos):
        if i >= len(pic_placeholders):
            print(f"⚠️  Not enough picture placeholders for {logo_name}")
            break

        logo_path_full = logo_path / logo_file

        if not logo_path_full.exists():
            print(f"❌ Logo not found: {logo_path_full}")
            continue

        try:
            pic_ph = pic_placeholders[i]
            print(f"  [{pic_ph.placeholder_format.idx}] {logo_name}:")
            print(f"      Placeholder: {pic_ph.name}")

            # THIS IS THE KEY: insert_picture() on picture placeholder
            pic = pic_ph.insert_picture(str(logo_path_full))

            print(f"      ✅ Inserted {logo_file}")

        except Exception as e:
            print(f"      ❌ Error: {e}")
            import traceback
            traceback.print_exc()

    # TEST 2: Create Content Slide (if Single Column layout exists)
    print()
    print("=" * 60)
    print("TEST 2: Content Slide from Layout")
    print("=" * 60)

    try:
        content_layout = get_layout(prs, "Single")
        slide2 = prs.slides.add_slide(content_layout)
        print(f"✅ Created slide from layout: {content_layout.name}")

        slide2.shapes.title.text = "Test Content Slide"

        # Find body placeholder
        body_ph = next(
            ph for ph in slide2.placeholders
            if ph.placeholder_format.type == PP_PLACEHOLDER.BODY
        )

        tf = body_ph.text_frame
        tf.clear()
        tf.text = "First bullet point"

        for text in ["Second bullet point", "Third bullet point"]:
            p = tf.add_paragraph()
            p.text = text
            p.level = 0

        print(f"✅ Added content with bullets")

    except ValueError as e:
        print(f"⚠️  Single Column layout not found, skipping")
    except Exception as e:
        print(f"❌ Error: {e}")

    # Save
    print()
    print(f"💾 Saving: {output_file}")
    prs.save(str(output_file))

    print()
    print(f"✨ Done! Open {output_file} to verify.")
    print(f"   Windows path: {output_file.absolute()}")
    print()
    print("Expected results:")
    print("  ✓ Title slide with logos inserted via insert_picture()")
    print("  ✓ Content slide with bullets (if layout exists)")

if __name__ == '__main__':
    test_layouts_only()
