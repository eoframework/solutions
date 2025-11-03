#!/usr/bin/env python3
"""
Generate test presentation with all 7 layouts populated with sample content.
Includes random text, images, and charts.
"""

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pathlib import Path
from datetime import datetime


def get_layout(prs, name_contains):
    """Get layout by partial name match."""
    for layout in prs.slide_layouts:
        if name_contains.lower() in layout.name.lower():
            return layout
    raise ValueError(f"Layout containing '{name_contains}' not found")


def insert_logos(slide, logo_path):
    """Insert logos into picture placeholders."""
    logos = {
        10: 'client_logo.png',
        13: 'consulting_company_logo.png',
        14: 'consulting_company_logo.png',  # Some layouts use idx 14 for logo
    }

    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx in logos and ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            logo_file = logo_path / logos[idx]
            if logo_file.exists():
                try:
                    ph.insert_picture(str(logo_file))
                    print(f"      ✅ Logo inserted at [{idx}]")
                except Exception as e:
                    print(f"      ❌ Logo failed at [{idx}]: {e}")


def create_sample_chart(placeholder):
    """Create a sample bar chart in a chart placeholder."""
    try:
        # Chart data
        chart_data = CategoryChartData()
        chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
        chart_data.add_series('Revenue ($M)', (45.2, 52.7, 61.3, 68.9))
        chart_data.add_series('Expenses ($M)', (32.1, 35.4, 38.2, 39.8))

        # Add chart to placeholder
        chart = placeholder.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart

        # Customize chart
        chart.has_legend = True
        chart.chart_title.text_frame.text = "Financial Performance"

        print(f"      ✅ Chart inserted")
        return True
    except Exception as e:
        print(f"      ⚠️  Chart insertion failed: {e}")
        return False


def test_all_7_layouts():
    """Generate test presentation with all 7 layouts."""

    template_path = Path('support/doc-templates/powerpoint/EOFramework-Template-01.pptx')
    logo_path = Path('support/doc-templates/powerpoint')
    output_file = Path('test-all-7-layouts.pptx')

    print(f"📂 Loading template: {template_path}")
    prs = Presentation(str(template_path))

    # Clear demo slides
    print(f"🗑️  Removing {len(prs.slides)} demo slides...")
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    print(f"✅ Starting fresh with {len(prs.slide_layouts)} layouts")
    print()

    # =================================================================
    # LAYOUT 0: EO Title Slide
    # =================================================================
    print("=" * 70)
    print("LAYOUT 0: EO Title Slide")
    print("=" * 70)

    layout0 = get_layout(prs, "Title")
    slide = prs.slides.add_slide(layout0)

    # [10] = Picture (Client Logo)
    # [12] = Text (Title)
    # [13] = Picture (Consulting Logo)
    # [14] = Text (Subtitle)
    # [15] = Text (Presenter | Date)

    slide.placeholders[12].text = "Digital Transformation Initiative"
    slide.placeholders[14].text = "Cloud Migration & Modernization Strategy"
    slide.placeholders[15].text = f"Solutions Architecture Team | {datetime.now().strftime('%B %d, %Y')}"
    print("  ✅ Text placeholders filled")

    insert_logos(slide, logo_path)
    print()

    # =================================================================
    # LAYOUT 1: EO Single Column
    # =================================================================
    print("=" * 70)
    print("LAYOUT 1: EO Single Column")
    print("=" * 70)

    layout1 = get_layout(prs, "Single")
    slide = prs.slides.add_slide(layout1)

    # [10] = Text (Title)
    # [11] = Text (Body/Content)
    # [13] = Picture (Logo)

    slide.placeholders[10].text = "Project Objectives"

    body = slide.placeholders[11]
    tf = body.text_frame
    tf.clear()
    tf.text = "Migrate 50+ legacy applications to cloud infrastructure"
    for text in [
        "Reduce operational costs by 40% within 18 months",
        "Improve system performance and scalability",
        "Enhance security posture and compliance",
        "Enable modern DevOps practices and CI/CD pipelines"
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0

    print("  ✅ Title and bullets filled")
    insert_logos(slide, logo_path)
    print()

    # =================================================================
    # LAYOUT 2: EO Two Column
    # =================================================================
    print("=" * 70)
    print("LAYOUT 2: EO Two Column")
    print("=" * 70)

    layout2 = get_layout(prs, "Two Column")
    slide = prs.slides.add_slide(layout2)

    # [10] = Text (Title)
    # [11] = Text (Left column header)
    # [12] = Text (Right column header)
    # [13] = Text (Left column content)
    # [14] = Picture (Logo)
    # [15] = Text (Right column content)

    slide.placeholders[10].text = "Current State vs. Future State"
    slide.placeholders[11].text = "Current Challenges"
    slide.placeholders[12].text = "Proposed Solutions"

    left_content = slide.placeholders[13]
    left_content.text = "Legacy monolithic architecture\nManual deployment processes\nLimited scalability\nHigh maintenance costs\nSlow time-to-market"

    right_content = slide.placeholders[15]
    right_content.text = "Microservices architecture\nAutomated CI/CD pipelines\nAuto-scaling capabilities\n40% cost reduction\n3x faster deployments"

    print("  ✅ Title and two columns filled")
    insert_logos(slide, logo_path)
    print()

    # =================================================================
    # LAYOUT 3: EO Key Points
    # =================================================================
    print("=" * 70)
    print("LAYOUT 3: EO Key Points")
    print("=" * 70)

    layout3 = get_layout(prs, "Key Points")
    slide = prs.slides.add_slide(layout3)

    # [10] = Text (Title)
    # [13] = Picture (Logo)
    # [14] = Text (Key points content)

    slide.placeholders[10].text = "Critical Success Factors"

    content = slide.placeholders[14]
    tf = content.text_frame
    tf.clear()
    tf.text = "Executive sponsorship and stakeholder buy-in across all business units"
    for text in [
        "Dedicated cross-functional team with cloud expertise and domain knowledge",
        "Comprehensive change management and training program for 500+ users",
        "Phased migration approach with minimal business disruption"
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0

    print("  ✅ Title and key points filled")
    insert_logos(slide, logo_path)
    print()

    # =================================================================
    # LAYOUT 4: EO Visual Content
    # =================================================================
    print("=" * 70)
    print("LAYOUT 4: EO Visual Content")
    print("=" * 70)

    layout4 = get_layout(prs, "Visual")
    slide = prs.slides.add_slide(layout4)

    # [10] = Text (Title)
    # [11] = Text (Subtitle/Section)
    # [13] = Text (Description)
    # [14] = Picture (Logo)
    # [15] = Picture (Main visual content)

    slide.placeholders[10].text = "Solution Architecture"
    slide.placeholders[11].text = "High-Level System Design"
    slide.placeholders[13].text = "The architecture leverages cloud-native services including containerization, serverless computing, and managed databases to deliver a scalable, resilient, and cost-effective solution."

    print("  ✅ Title and description filled")

    # Insert visual content (use logo as placeholder image)
    visual_logo = logo_path / 'eo-framework-logo-real.png'
    if visual_logo.exists():
        try:
            slide.placeholders[15].insert_picture(str(visual_logo))
            print("  ✅ Visual content image inserted")
        except Exception as e:
            print(f"  ⚠️  Visual image failed: {e}")

    insert_logos(slide, logo_path)
    print()

    # =================================================================
    # LAYOUT 5: EO Data Visualization
    # =================================================================
    print("=" * 70)
    print("LAYOUT 5: EO Data Visualization")
    print("=" * 70)

    layout5 = get_layout(prs, "Data")
    slide = prs.slides.add_slide(layout5)

    # [10] = Text (Title)
    # [11] = Text (Subtitle)
    # [13] = Text (Description/Insights)
    # [14] = Picture (Logo)
    # [15] = Chart (placeholder)

    slide.placeholders[10].text = "Financial Impact Analysis"
    slide.placeholders[11].text = "Quarterly Performance Metrics"
    slide.placeholders[13].text = "Key Insights:\n• Revenue increased 53% year-over-year\n• Expense management improved by 24%\n• Net margin expansion of 15 percentage points\n• Strong growth trajectory maintained"

    print("  ✅ Title, subtitle and insights filled")

    # Try to insert chart
    try:
        # Find chart placeholder
        chart_placeholder = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 15:
                chart_placeholder = ph
                break

        if chart_placeholder:
            create_sample_chart(chart_placeholder)
        else:
            print("  ⚠️  Chart placeholder [15] not found")
    except Exception as e:
        print(f"  ⚠️  Chart creation failed: {e}")

    insert_logos(slide, logo_path)
    print()

    # =================================================================
    # LAYOUT 6: EO Thank You
    # =================================================================
    print("=" * 70)
    print("LAYOUT 6: EO Thank You")
    print("=" * 70)

    layout6 = get_layout(prs, "Thank")
    slide = prs.slides.add_slide(layout6)

    # [10] = Picture (Top logo)
    # [12] = Text (Thank you message)
    # [13] = Picture (Bottom logo)
    # [14] = Text (Contact info)

    slide.placeholders[12].text = "Thank You"
    slide.placeholders[14].text = "Questions?\n\nContact: solutions@eoframework.com\nWeb: www.eoframework.com\nPhone: +1 (555) 123-4567"

    print("  ✅ Thank you message and contact filled")
    insert_logos(slide, logo_path)
    print()

    # Save
    print("=" * 70)
    print(f"💾 Saving: {output_file}")
    prs.save(str(output_file))

    print()
    print("✨ Done! Open test-all-7-layouts.pptx to verify.")
    print(f"   Windows path: {output_file.absolute()}")
    print()
    print("Generated slides:")
    print("  1. Title Slide - Digital Transformation Initiative")
    print("  2. Single Column - Project Objectives with 5 bullets")
    print("  3. Two Column - Current vs. Future State comparison")
    print("  4. Key Points - Critical Success Factors")
    print("  5. Visual Content - Solution Architecture with image")
    print("  6. Data Visualization - Financial Impact with chart")
    print("  7. Thank You - Closing slide with contact info")
    print()
    print(f"Total slides: 7")
    print(f"All layouts demonstrated with logos and sample content!")


if __name__ == '__main__':
    test_all_7_layouts()
