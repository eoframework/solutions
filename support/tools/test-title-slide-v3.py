#!/usr/bin/env python3
"""
Test v3: Use picture fill method to replace placeholder images.
"""

from pptx import Presentation
from pathlib import Path
import re

def extract_title_from_markdown(md_file):
    """Extract the first H1 title from markdown file."""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
    if match:
        return match.group(1).strip()
    return "Untitled Presentation"

def create_title_slide_v3():
    """Create title slide using picture fill method."""

    # Paths
    template_path = Path('support/doc-templates/powerpoint/EOFramework-Template-01.pptx')
    logo_path = Path('support/doc-templates/powerpoint')
    md_file = Path('solution-template/sample-provider/sample-category/sample-solution/presales/raw/executive-presentation.md')
    output_file = Path('test-title-slide-v3.pptx')

    # Load template
    print(f"📂 Loading template: {template_path}")
    prs = Presentation(str(template_path))

    # Create slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    print(f"✅ Created slide from layout: {title_layout.name}")

    # Extract title
    title_text = extract_title_from_markdown(md_file)
    print(f"📝 Title from markdown: {title_text}")
    print()

    print("🎨 Filling placeholders using picture fill:")

    # Define logo mappings
    logo_mapping = {
        1: ('acme_logo_rectangle.png', 'Acme Inc'),
        3: ('zplus_consulting_logo.png', 'ZplusConsulting'),
        4: ('eo-framework-logo-converted.png', 'EO Framework')
    }

    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue

        idx = shape.placeholder_format.idx

        # Handle logo placeholders
        if idx in logo_mapping:
            logo_file, logo_name = logo_mapping[idx]
            logo_full_path = logo_path / logo_file

            if logo_full_path.exists():
                try:
                    # Method 1: Try fill.solid() then fore_color.rgb with picture
                    # This doesn't work for pictures, so let's try the direct approach

                    # Method 2: Delete placeholder and add picture in same position
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height

                    # Remove the placeholder
                    sp = shape.element
                    sp.getparent().remove(sp)

                    # Add picture in exact same position
                    pic = slide.shapes.add_picture(
                        str(logo_full_path),
                        left, top, width, height
                    )

                    print(f"  [{idx}] ✅ {logo_name}: Replaced placeholder with picture")

                except Exception as e:
                    print(f"  [{idx}] ❌ {logo_name}: Failed - {e}")
            else:
                print(f"  [{idx}] ❌ {logo_name}: File not found - {logo_full_path}")

        # Handle title placeholder
        elif idx == 12:
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = title_text
                print(f"  [{idx}] ✅ Title: Set to '{title_text}'")

    # Save
    print(f"\n💾 Saving: {output_file}")
    prs.save(str(output_file))

    print(f"\n✨ Done! Open {output_file} to verify.")
    print(f"   Windows path: {output_file.absolute()}")

if __name__ == '__main__':
    create_title_slide_v3()
