#!/usr/bin/env python3
"""
Test v4: Collect placeholder positions first, then replace them.
"""

from pptx import Presentation
from pathlib import Path
import re

def extract_title_from_markdown(md_file):
    """Extract the first H1 title from markdown file."""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
    if match:
        return match.group(1).strip()
    return "Untitled Presentation"

def create_title_slide_v4():
    """Create title slide - collect positions first, then replace."""

    # Paths
    template_path = Path('support/doc-templates/powerpoint/EOFramework-Template-01.pptx')
    logo_path = Path('support/doc-templates/powerpoint')
    md_file = Path('solution-template/sample-provider/sample-category/sample-solution/presales/raw/executive-presentation.md')
    output_file = Path('test-title-slide-v4.pptx')

    # Load template
    print(f"📂 Loading template: {template_path}")
    prs = Presentation(str(template_path))

    # Create slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    print(f"✅ Created slide from layout: {title_layout.name}")

    # Extract title
    title_text = extract_title_from_markdown(md_file)
    print(f"📝 Title from markdown: {title_text}")
    print()

    # Define logo mappings
    logo_mapping = {
        1: ('acme_logo_rectangle.png', 'Acme Inc'),
        3: ('zplus_consulting_logo.png', 'ZplusConsulting'),
        4: ('eo-framework-logo-converted.png', 'EO Framework')
    }

    # STEP 1: Collect placeholder information
    print("📋 Step 1: Collecting placeholder positions...")
    placeholder_info = []
    title_placeholder = None

    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue

        idx = shape.placeholder_format.idx

        # Collect logo placeholder info
        if idx in logo_mapping:
            info = {
                'idx': idx,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'element': shape.element,
                'logo_file': logo_mapping[idx][0],
                'logo_name': logo_mapping[idx][1]
            }
            placeholder_info.append(info)
            print(f"  [{idx}] {logo_mapping[idx][1]}: Position ({shape.left/914400:.2f}\", {shape.top/914400:.2f}\") Size ({shape.width/914400:.2f}\" x {shape.height/914400:.2f}\")")

        # Handle title placeholder
        elif idx == 12:
            title_placeholder = shape

    # STEP 2: Remove logo placeholders
    print(f"\n🗑️  Step 2: Removing {len(placeholder_info)} logo placeholders...")
    for info in placeholder_info:
        info['element'].getparent().remove(info['element'])

    # STEP 3: Add pictures at the correct positions
    print(f"\n🎨 Step 3: Adding logo pictures at placeholder positions...")
    for info in placeholder_info:
        logo_full_path = logo_path / info['logo_file']

        if logo_full_path.exists():
            try:
                pic = slide.shapes.add_picture(
                    str(logo_full_path),
                    info['left'],
                    info['top'],
                    info['width'],
                    info['height']
                )
                print(f"  [{info['idx']}] ✅ {info['logo_name']}: Added at position ({info['left']/914400:.2f}\", {info['top']/914400:.2f}\")")
            except Exception as e:
                print(f"  [{info['idx']}] ❌ {info['logo_name']}: Failed - {e}")
        else:
            print(f"  [{info['idx']}] ❌ {info['logo_name']}: File not found - {logo_full_path}")

    # STEP 4: Set title text
    print(f"\n✏️  Step 4: Setting title text...")
    if title_placeholder and hasattr(title_placeholder, 'text_frame'):
        title_placeholder.text_frame.text = title_text
        print(f"  ✅ Title set to: '{title_text}'")

    # Save
    print(f"\n💾 Saving: {output_file}")
    prs.save(str(output_file))

    print(f"\n✨ Done! Open {output_file} to verify.")
    print(f"   Windows path: {output_file.absolute()}")
    print(f"\nLogos should appear in their correct positions (not at the bottom)")

if __name__ == '__main__':
    create_title_slide_v4()
