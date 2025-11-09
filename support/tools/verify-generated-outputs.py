#!/usr/bin/env python3
"""
Verify generated output files quality.
"""

from pathlib import Path
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import sys

def verify_word_doc(docx_path):
    """Verify Word document has proper formatting."""
    doc = Document(docx_path)

    # Check for basic elements
    has_title = False
    has_headings = False
    has_footer = False
    has_logo = False

    for para in doc.paragraphs:
        if para.style.name == 'Title':
            has_title = True
        if para.style.name.startswith('Heading'):
            has_headings = True

        # Check for logo
        for run in para.runs:
            if hasattr(run, '_element'):
                drawings = run._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing')
                if drawings:
                    has_logo = True

    # Check footer
    for section in doc.sections:
        if section.footer.paragraphs:
            has_footer = True

    return {
        'title': has_title,
        'headings': has_headings,
        'footer': has_footer,
        'logo': has_logo,
        'paragraphs': len(doc.paragraphs),
        'tables': len(doc.tables)
    }

def verify_excel(xlsx_path):
    """Verify Excel file has proper formatting."""
    wb = load_workbook(xlsx_path)

    # Get first sheet
    ws = wb.active

    # Check for headers
    has_headers = ws['A1'].value is not None
    has_formatting = ws['A1'].fill.start_color.rgb is not None if ws['A1'].fill else False

    return {
        'sheets': len(wb.sheetnames),
        'rows': ws.max_row,
        'cols': ws.max_column,
        'has_headers': has_headers,
        'has_formatting': has_formatting
    }

def verify_pptx(pptx_path):
    """Verify PowerPoint presentation."""
    prs = Presentation(pptx_path)

    return {
        'slides': len(prs.slides),
        'has_content': len(prs.slides) > 0
    }

if __name__ == '__main__':
    solution_path = Path(sys.argv[1]) if len(sys.argv) > 1 else Path('solution-template/sample-provider/sample-category/sample-solution')

    print("=" * 80)
    print("GENERATED OUTPUTS VERIFICATION")
    print("=" * 80)

    # Check presales
    presales_dir = solution_path / 'presales'
    print(f"\n📁 PRESALES: {presales_dir}")
    print("-" * 80)

    if presales_dir.exists():
        for file in sorted(presales_dir.glob('*.docx')):
            result = verify_word_doc(file)
            print(f"\n📄 {file.name}")
            print(f"   Title: {'✅' if result['title'] else '❌'}")
            print(f"   Headings: {'✅' if result['headings'] else '❌'}")
            print(f"   Footer: {'✅' if result['footer'] else '❌'}")
            print(f"   Logo: {'✅' if result['logo'] else '❌'}")
            print(f"   Paragraphs: {result['paragraphs']}")
            print(f"   Tables: {result['tables']}")

        for file in sorted(presales_dir.glob('*.xlsx')):
            result = verify_excel(file)
            print(f"\n📊 {file.name}")
            print(f"   Sheets: {result['sheets']}")
            print(f"   Rows: {result['rows']}, Cols: {result['cols']}")
            print(f"   Headers: {'✅' if result['has_headers'] else '❌'}")

        for file in sorted(presales_dir.glob('*.pptx')):
            result = verify_pptx(file)
            print(f"\n📽️  {file.name}")
            print(f"   Slides: {result['slides']}")
            print(f"   Content: {'✅' if result['has_content'] else '❌'}")

    # Check delivery
    delivery_dir = solution_path / 'delivery'
    print(f"\n\n📁 DELIVERY: {delivery_dir}")
    print("-" * 80)

    if delivery_dir.exists():
        for file in sorted(delivery_dir.glob('*.docx')):
            result = verify_word_doc(file)
            print(f"\n📄 {file.name}")
            print(f"   Title: {'✅' if result['title'] else '❌'}")
            print(f"   Headings: {'✅' if result['headings'] else '❌'}")
            print(f"   Footer: {'✅' if result['footer'] else '❌'}")
            print(f"   Logo: {'✅' if result['logo'] else '❌'}")
            print(f"   Paragraphs: {result['paragraphs']}")
            print(f"   Tables: {result['tables']}")

        for file in sorted(delivery_dir.glob('*.xlsx')):
            result = verify_excel(file)
            print(f"\n📊 {file.name}")
            print(f"   Sheets: {result['sheets']}")
            print(f"   Rows: {result['rows']}, Cols: {result['cols']}")
            print(f"   Headers: {'✅' if result['has_headers'] else '❌'}")

        for file in sorted(delivery_dir.glob('*.pptx')):
            result = verify_pptx(file)
            print(f"\n📽️  {file.name}")
            print(f"   Slides: {result['slides']}")
            print(f"   Content: {'✅' if result['has_content'] else '❌'}")

    print("\n" + "=" * 80)
    print("VERIFICATION COMPLETE")
    print("=" * 80)
