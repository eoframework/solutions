#!/usr/bin/env python3
"""
Test v6: Use the demo slide and replace image relationships.
"""

from pptx import Presentation
from pptx.parts.image import ImagePart
from pathlib import Path
import re

def extract_title_from_markdown(md_file):
    """Extract the first H1 title from markdown file."""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
    if match:
        return match.group(1).strip()
    return "Untitled Presentation"

def replace_image_in_placeholder(slide_part, relationship_id, new_image_path, shape):
    """Replace image in a relationship (similar to Word example)."""
    try:
        # Get the current relationship
        image_rel = slide_part.rels[relationship_id]

        # Get the current image part
        old_image_part = image_rel.target_part

        # Read new image
        with open(new_image_path, 'rb') as f:
            image_blob = f.read()

        # Replace the blob in the existing image part
        old_image_part._blob = image_blob

        # Remove cropping (srcRect) to prevent image from being cropped
        pic = shape._element
        blipFill = pic.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}blipFill')
        if blipFill is not None:
            srcRect = blipFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srcRect')
            if srcRect is not None:
                blipFill.remove(srcRect)
                print(f"      Removed cropping")

        return True
    except Exception as e:
        print(f"      ❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return False

def create_title_slide_v6():
    """Use demo slide and replace image relationships."""

    # Paths
    template_path = Path('support/doc-templates/powerpoint/EOFramework-Template-01.pptx')
    logo_path = Path('support/doc-templates/powerpoint')
    md_file = Path('solution-template/sample-provider/sample-category/sample-solution/presales/raw/executive-presentation.md')
    output_file = Path('test-title-slide-v6.pptx')

    # Load template
    print(f"📂 Loading template: {template_path}")
    prs = Presentation(str(template_path))

    # Use the demo slide (don't create new one)
    if len(prs.slides) == 0:
        print("❌ No demo slide found in template!")
        return

    slide = prs.slides[0]
    print(f"✅ Using demo slide: {slide.slide_layout.name}")

    # Extract title
    title_text = extract_title_from_markdown(md_file)
    print(f"📝 Title from markdown: {title_text}")
    print()

    # Map placeholder indices to logos
    # Based on current template analysis:
    logo_mapping = {
        10: ('acme_logo_rectangle.png', 'Acme Inc'),
        13: ('zplus_consulting_logo.png', 'ZplusConsulting'),
    }

    print("🎨 Replacing placeholder images:")

    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx

        # Handle logo placeholders
        if idx in logo_mapping:
            logo_file, logo_name = logo_mapping[idx]
            logo_full_path = logo_path / logo_file

            print(f"  [{idx}] {logo_name}:")

            if not logo_full_path.exists():
                print(f"      ❌ File not found: {logo_full_path}")
                continue

            # Find the relationship ID
            pic = shape._element
            blipFill = pic.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}blipFill')

            if blipFill is None:
                print(f"      ❌ No blipFill found")
                continue

            blip = blipFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')

            if blip is None:
                print(f"      ❌ No blip found")
                continue

            embed_rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')

            if not embed_rId:
                print(f"      ❌ No embed relationship found")
                continue

            print(f"      Relationship ID: {embed_rId}")

            # Replace the image
            success = replace_image_in_placeholder(slide.part, embed_rId, logo_full_path, shape)

            if success:
                print(f"      ✅ Replaced with {logo_file}")

        # Handle title placeholder
        elif idx == 12:
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = title_text
                print(f"  [{idx}] ✅ Title: Set to '{title_text}'")

    # Save
    print(f"\n💾 Saving: {output_file}")
    prs.save(str(output_file))

    print(f"\n✨ Done! Open {output_file} to verify.")
    print(f"   Windows path: {output_file.absolute()}")
    print(f"\nThe 2 placeholder images should now show:")
    print(f"  - Acme Inc logo")
    print(f"  - ZplusConsulting logo")

if __name__ == '__main__':
    create_title_slide_v6()
